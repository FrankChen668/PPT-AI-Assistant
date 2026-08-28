#!/usr/bin/env python3
"""Quality gate for AI-PPT projects.

The QA layer validates the contracts between architecture stages:

- blueprint.json is structurally renderable.
- State 2.5 artifacts (art direction / reference pack / slide plan) are consumed before State 3.
- style_route/style_drafts gating is satisfied for low-confidence style routing.
- SVG slides exist, parse, and stay within canvas constraints.
- Optional PPTX exports contain the expected number of slides and native objects.
- Optional snapshots provide a visual review surface without changing artifacts.
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import sys
import xml.etree.ElementTree as ET
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable, Sequence, TypedDict

SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPT_DIR))

from aesthetic_critic import run_aesthetic_critic  # noqa: E402
from blueprint_schema import validate_blueprint_schema  # noqa: E402
from canvas_context import CanvasContext, load_canvas_context  # noqa: E402
from check_template_binding_consistency import evaluate_template_binding_consistency  # noqa: E402
from check_visual_baseline import baseline_grade_for_project  # noqa: E402
from check_visual_diversity import check_visual_diversity  # noqa: E402
from copyfit_contract import is_dense_slide  # noqa: E402
from design_spec_tokens import parse_data_palette, split_design_spec_lines  # noqa: E402
from layout_contracts import CORE_SLIDE_KEYS, layout_tags, required_content_keys, validate_content_shape  # noqa: E402
from layout_renderers import LayoutRenderer  # noqa: E402
from pipeline.design_token_guard import run_design_token_guard  # noqa: E402
from pipeline.qa_non_critical import run_snapshot_and_visual_checks  # noqa: E402
from pipeline.slide_budget import build_slide_budget  # noqa: E402
from profile_policy import resolve_profile_policy  # noqa: E402
from quality.json_contracts import load_json_object  # noqa: E402
from quality.reporting_outputs import write_reports as write_reports_output  # noqa: E402
from quality.warning_classification import (  # noqa: E402
    ASSET_NON_BLOCKING_WARNING_CODES,
    BUDGET_NON_BLOCKING_WARNING_CODES,
    DESIGN_TOKEN_NON_BLOCKING_WARNING_CODES,
    QUALITY_MODES,
    RELEASE_SAFE_NON_BLOCKING_WARNING_CODES,
)
from quality.warning_classification import (  # noqa: E402
    is_visual_delivery_code as _is_visual_delivery_code,
)
from quality.warning_classification import (  # noqa: E402
    is_warning_non_blocking as _is_warning_non_blocking,
)
from quality.warning_classification import (  # noqa: E402
    normalize_quality_mode as _normalize_quality_mode,
)
from render_theme import Theme, visual_width  # noqa: E402
from svg_quality_checker import check_svg_file  # noqa: E402

BAD_TAGS = {"script", "foreignObject"}
TEXT_TAGS = {"text", "tspan"}
# These mojibake literals are intentional: they catch historical Chinese text
# that was decoded with the wrong codec. Do not normalize them without tests.
CONCLUSION_ANCHOR_RE = re.compile(
    r"(鏈€缁堜环鍊紎鏈€缁堢粨璁簗鍏抽敭缁撹|缁撹锛殀Conclusion|Final Value)",
    re.IGNORECASE,
)
ENGINE_TITLE_RE = re.compile(r"(璇佹槑閾惧紩鎿巪杩芥函寮曟搸|寮曟搸|Engine|Proof Chain)", re.IGNORECASE)
CHART_HINT_RE = re.compile(
    r"(chart|kpi|metric|trend|data\s*visual|visualization|鏁版嵁|鎸囨爣|鍥捐〃|瓒嬪娍)",
    re.IGNORECASE,
)
ICON_REF_RE = re.compile(r"icon_ref\s*[:=]\s*([A-Za-z0-9:_./-]+)")
ICON_REF_ATTR_RE = re.compile(r"data-icon-ref\s*=\s*['\"]([A-Za-z0-9:_./-]+)['\"]", re.IGNORECASE)
TRANSLATE_RE = re.compile(
    r"translate\(\s*([-+]?(?:\d+(?:\.\d+)?|\.\d+))(?:[\s,]+([-+]?(?:\d+(?:\.\d+)?|\.\d+)))?\s*\)"
)
SAFE_AREA_PROFILES: dict[str, dict[str, float]] = {
    "legacy": {"pad_x": 32.0, "pad_top": 32.0, "pad_bottom": 48.0},
    "presentation": {"pad_x": 24.0, "pad_top": 24.0, "pad_bottom": 20.0},
}
DEV_FAST_WARNING_TO_ADVISORY_CODES = {
    "style-over-cardization",
    "style-rhythm-monotony",
    "style-conclusion-first-weak",
    "style-takeaway-bar-missing",
    "visual-diversity",
    "prompt-pattern-missing",
    "prompt-pattern-incomplete",
}
PREMIUM_ADVISORY_TO_WARNING_CODES = {
    "visual-headline-weak",
    "visual-density-high",
    "visual-text-fragmented",
    "visual-hierarchy-flat",
    "visual-alignment-chaos",
    "style-over-cardization",
    "style-rhythm-monotony",
}
PREMIUM_WARNING_TO_ERROR_CODES = {
    "style-over-cardization",
    "style-rhythm-monotony",
    "style-takeaway-bar-missing",
    "prompt-pattern-missing",
    "prompt-pattern-incomplete",
    "visual-baseline-below-b",
    "visual-diversity",
    "style-hard-token-color-limit",
    "style-hard-token-consecutive-homogeneous",
    "style-hard-token-conclusion-hierarchy-weak",
}
DEV_FAST_WARNING_TO_ADVISORY_CODES.update(
    {
        "style-hard-token-color-limit",
        "style-hard-token-consecutive-homogeneous",
        "style-hard-token-conclusion-hierarchy-weak",
    }
)
CONSULTING_PROFILE_TOKENS = ("executive_exhibit", "strategy", "consult", "proposal", "tender", "bid")
TOKEN_BUDGET_POLICY = "default-v1"
TOKEN_BUDGET_CAPS = {
    "strategist": 3000,
    "designer": 4000,
    "executor_per_slide": 2500,
    "checks": 8000,
}
CONTENT_PLACEHOLDER_PATTERNS: tuple[tuple[str, re.Pattern[str]], ...] = (
    ("generic-point", re.compile(r"要点\s*[一二三四五六七八九十\d]+")),
    ("placeholder-word", re.compile(r"(占位|待补充|待完善|待定|示例文本)")),
    ("english-placeholder", re.compile(r"\b(TBD|TODO|placeholder|lorem\s+ipsum)\b", re.IGNORECASE)),
)
# B3 (AB-05/AB-06) advisory-only semantic codes. Deliberately excluded from
# PROPOSAL_CRITIC_BLOCKING_CODES and every quality-mode escalation set so they
# stay advisory in all three quality modes (AB-07).
CONCLUSION_HEADLINE_MIN_FONT = 28.0
CONCLUSION_ZONE_MIN_TEXT_CHARS = 8
CONCLUSION_ZONE_CLUSTER_GAP = 140.0
CLAIM_BOUNDARY_QUALIFIER_TERMS = (
    "预计",
    "计划",
    "假设",
    "拟",
    "待验证",
    "预期",
    "有望",
    "尚待",
    "estimated",
    "planned",
    "expected",
    "assumption",
    "hypothesis",
    "pending",
)
PROPOSAL_CRITIC_BLOCKING_CODES = {
    "style-over-cardization",
    "style-rhythm-monotony",
    "style-conclusion-first-weak",
    "style-takeaway-bar-missing",
    "visual-density-high",
    "visual-text-fragmented",
    "visual-hierarchy-flat",
    "visual-alignment-chaos",
    "visual-headline-weak",
    "visual-whitespace-low",
    "visual-dominance-weak",
    "text-cjk-longline-in-narrow-column",
}
TEMPLATE_CONSISTENCY_HARD_CONFLICT_CODES = {
    "template-binding-reference-pack-mismatch",
    "template-binding-slide-visual-plan-mismatch",
}
TEMPLATE_CONSISTENCY_MISSING_CONTEXT_CODES = {
    "missing-reference-pack-for-binding",
    "missing-slide-visual-plan-for-binding",
}
TEMPLATE_CONSISTENCY_PARSE_OR_SHAPE_CODES = {
    "invalid-template-binding-json",
    "invalid-reference-pack-json",
    "invalid-slide-visual-plan-json",
    "template-binding-missing-layout-id",
}
TEMPLATE_CONSISTENCY_SOFT_SIGNAL_CODES = {
    "template-binding-override-reason",
    "template-binding-slide-visual-plan-unspecified",
}
REFERENCE_FIRST_GATE_CODES = {
    "reference-pack-empty",
    "reference-pack-free-design-reason-missing",
}
VISUAL_CONTRACT_REQUIRED_KEYS = {
    "scene_type",
    "generation_strategy",
    "focal_point",
    "primary_read_path",
    "composition_grammar",
    "hierarchy_ladder",
    "density_budget",
    "whitespace_target",
    "template_inheritance",
    "anti_patterns",
    "critic_checks",
}
VISUAL_CONTRACT_V2_REQUIRED_KEYS = {
    "layout_intent",
    "bbox_budget",
    "text_budget",
    "deterministic_scaffold",
    "must_avoid",
    "pre_authoring_checks",
}
EXECUTION_POLICY_REQUIRED_KEYS = {
    "scene_type",
    "generation_strategy",
    "risk_level",
    "required_loop",
    "qa_strictness",
    "expected_first_pass_rules",
}
PLANNING_REQUIRED_FIELDS = {
    "layout_objective",
    "density_budget",
    "dominance_map",
    "must_keep_claims",
}
ARCHITECTURE_SLOT_GROUPS = (
    {
        "content_key": "left_systems",
        "slot_prefix": "left_systems",
        "missing_slot_code": "contract-svg-architecture-missing-left-slot",
        "missing_text_code": "contract-svg-architecture-missing-left-text",
        "label": "left systems",
        "strict_slot_ids": False,
        "require_slot_box": False,
    },
    {
        "content_key": "core_modules",
        "slot_prefix": "core_modules",
        "missing_slot_code": "contract-svg-architecture-missing-core-slot",
        "missing_text_code": "contract-svg-architecture-missing-core-title",
        "label": "core modules",
        "strict_slot_ids": True,
        "require_slot_box": True,
    },
    {
        "content_key": "right_modules",
        "slot_prefix": "right_modules",
        "missing_slot_code": "contract-svg-architecture-missing-right-slot",
        "missing_text_code": "contract-svg-architecture-missing-right-text",
        "label": "right modules",
        "strict_slot_ids": True,
        "require_slot_box": True,
    },
    {
        "content_key": "storage",
        "slot_prefix": "storage",
        "missing_slot_code": "contract-svg-architecture-missing-storage-slot",
        "missing_text_code": "contract-svg-architecture-missing-storage-text",
        "label": "storage",
        "strict_slot_ids": True,
        "require_slot_box": True,
    },
)
ROADMAP_PHASE_SLOT_PREFIX = "phases"
ROADMAP_SUMMARY_SLOT_ID = "summary"
COMPARISON_ROW_SLOT_PREFIX = "rows"
COMPARISON_SUMMARY_SLOT_ID = "summary"
COMPARISON_DIMENSION_HEADER_ANCHOR = "comparison-dimension-header"
COMPARISON_OPTION_LEFT_ANCHOR = "comparison-option-left"
COMPARISON_OPTION_RIGHT_ANCHOR = "comparison-option-right"
COMPARISON_RECOMMENDATION_ANCHOR = "comparison-recommendation"
CORE_ORBIT_CORE_SLOT_ID = "core-node"
CORE_ORBIT_SATELLITE_SLOT_PREFIX = "satellite"
CORE_ORBIT_CORE_SHAPE_ANCHOR = "core-node-shape"
CORE_ORBIT_RELATIONSHIP_EDGE_ANCHOR = "relationship-edge"


@dataclass
class Finding:
    severity: str
    code: str
    path: str
    message: str


@dataclass
class QaReport:
    project: str
    ok: bool
    errors: int
    warnings: int
    findings: list[Finding]
    metrics: dict[str, Any]
    visual_score: float | None = None
    visual_findings: list[dict[str, Any]] | None = None
    repair_recommendation: list[dict[str, Any]] | None = None
    density_flag: bool = False
    hierarchy_flag: bool = False
    advisories: int = 0


class VisualMetrics(TypedDict):
    visual_score: float | None
    visual_findings: list[dict[str, Any]]
    repair_recommendation: list[dict[str, Any]]
    density_flag: bool
    hierarchy_flag: bool
    visual_whitespace_ratio: float
    visual_hierarchy_depth_score: float
    visual_dominant_point_count: int
    visual_repetition_penalty: float
    visual_alignment_quality_score: float
    visual_fragmentation_by_slide: dict[str, dict[str, Any]]


def local_name(tag: str) -> str:
    if "}" in tag:
        return tag.rsplit("}", 1)[1]
    return tag


def number(value: str | None, default: float = 0.0) -> float:
    if value is None:
        return default
    text = value.strip().replace("px", "")
    try:
        return float(text)
    except ValueError:
        return default


def _project_relative_path(project_dir: Path, path: Path) -> str:
    try:
        return str(path.relative_to(project_dir))
    except ValueError:
        return str(path)


def _read_design_spec_values(project_dir: Path) -> dict[str, str]:
    design_spec = project_dir / "design_spec.md"
    if not design_spec.exists():
        return {}
    return split_design_spec_lines(design_spec.read_text(encoding="utf-8-sig"))


def _iter_text_values(node: Any):
    if isinstance(node, str):
        value = node.strip()
        if value:
            yield value
        return
    if isinstance(node, dict):
        for value in node.values():
            yield from _iter_text_values(value)
        return
    if isinstance(node, list):
        for value in node:
            yield from _iter_text_values(value)


def _parse_banned_terms(raw: str | None) -> list[str]:
    text = (raw or "").strip()
    if not text:
        return []
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError:
        parsed = None
    if isinstance(parsed, list):
        candidates = [str(item).strip() for item in parsed]
    else:
        normalized = text.strip("[]")
        candidates = [item.strip().strip("'\"") for item in re.split(r"[,锛?锛泑\n]+", normalized)]
    seen: set[str] = set()
    terms: list[str] = []
    for item in candidates:
        if item and item not in seen:
            seen.add(item)
            terms.append(item)
    return terms


def _contains_term(text: str, term: str) -> bool:
    if term.isascii():
        return term.lower() in text.lower()
    return term in text


def validate_banned_terms(
    project_dir: Path,
    slides: list[dict[str, Any]],
    findings: list[Finding],
) -> dict[str, Any]:
    values = _read_design_spec_values(project_dir)
    terms = _parse_banned_terms(values.get("banned_terms") or values.get("forbidden_terms"))
    metrics: dict[str, Any] = {
        "banned_term_count": len(terms),
        "banned_term_hit_count": 0,
        "banned_term_hit_slides": [],
    }
    if not terms:
        return metrics

    hit_slides: list[int] = []
    for index, slide in enumerate(slides, start=1):
        sid = slide.get("id")
        slide_id = sid if isinstance(sid, int) else index
        text_blob = "\n".join(_iter_text_values(slide))
        hits = [term for term in terms if _contains_term(text_blob, term)]
        if not hits:
            continue
        hit_slides.append(slide_id)
        metrics["banned_term_hit_count"] += len(hits)
        emit(
            findings,
            "error",
            "banned-term-hit",
            f"blueprint.json/slides/{slide_id}",
            f"Slide {slide_id} contains banned term(s): {', '.join(hits)}.",
        )
    metrics["banned_term_hit_slides"] = hit_slides
    return metrics


def _slide_has_chart_hint(slide: dict[str, Any]) -> bool:
    layout_tag = str(slide.get("layout_tag") or "").lower()
    if "chart" in layout_tag or "kpi" in layout_tag or layout_tag.startswith("data-"):
        return True
    title = str(slide.get("title") or "")
    if CHART_HINT_RE.search(title):
        return True
    content = slide.get("content")
    if isinstance(content, dict):
        for key in ("insight", "headline", "subtitle", "title", "label"):
            value = content.get(key)
            if isinstance(value, str) and CHART_HINT_RE.search(value):
                return True
        if isinstance(content.get("bars"), list) or isinstance(content.get("points"), list):
            return True
    return False


def validate_data_palette_presence(
    project_dir: Path,
    slides: list[dict[str, Any]],
    findings: list[Finding],
) -> dict[str, Any]:
    values = _read_design_spec_values(project_dir)
    data_palette = parse_data_palette(values.get("data_palette"))
    chart_hint_slides: list[int] = []
    for slide in slides:
        slide_id = slide.get("id")
        if isinstance(slide_id, int) and _slide_has_chart_hint(slide):
            chart_hint_slides.append(slide_id)
    metrics: dict[str, Any] = {
        "data_palette_count": len(data_palette),
        "data_palette_defined": len(data_palette) > 0,
        "chart_hint_slide_count": len(chart_hint_slides),
        "chart_hint_slide_ids": chart_hint_slides,
    }
    if chart_hint_slides and not data_palette:
        emit(
            findings,
            "warning",
            "theme-data-palette-missing",
            str(project_dir / "design_spec.md"),
            "Chart/KPI content detected but design_spec.md has no valid data_palette.",
        )
    return metrics


def _estimate_tokens_from_bytes(size_bytes: int) -> int:
    return max(0, int(round(size_bytes / 4.0)))


def _qa_context_profile(project_dir: Path, slide_id: int | None = None) -> dict[str, Any]:
    files: list[Path] = [
        project_dir / "design_spec.md",
        project_dir / "blueprint.json",
    ]
    optional = [
        project_dir / "style_route.json",
        project_dir / "slide_visual_plan.json",
        project_dir / "art_direction.md",
    ]
    files.extend(path for path in optional if path.exists())
    svg_dir = project_dir / "svg_output"
    if slide_id is not None:
        files.append(svg_dir / f"slide_{slide_id:02d}.svg")
    else:
        files.extend(sorted(svg_dir.glob("slide_*.svg"))[:5])
    existing = [path for path in files if path.exists() and path.is_file()]
    total_bytes = sum(path.stat().st_size for path in existing)
    return {
        "context_file_count": len(existing),
        "context_bytes_estimate": total_bytes,
        "context_files": [_project_relative_path(project_dir, path) for path in existing],
    }


def emit(findings: list[Finding], severity: str, code: str, path: str, message: str) -> None:
    findings.append(Finding(severity, code, path, message))


def _template_consistency_severity(code: str, quality_mode: str, profile: str) -> str:
    mode = _normalize_quality_mode(quality_mode)
    profile_normalized = str(profile or "presentation").strip().lower()
    proposal_profile = profile_normalized == "proposal_consulting"

    if mode == "premium":
        if code in TEMPLATE_CONSISTENCY_HARD_CONFLICT_CODES:
            return "error"
        if code in TEMPLATE_CONSISTENCY_MISSING_CONTEXT_CODES or code in TEMPLATE_CONSISTENCY_PARSE_OR_SHAPE_CODES:
            return "warning"
        return "advisory"

    if mode == "release-safe":
        if proposal_profile and code in TEMPLATE_CONSISTENCY_HARD_CONFLICT_CODES:
            return "warning"
        if proposal_profile and (
            code in TEMPLATE_CONSISTENCY_MISSING_CONTEXT_CODES
            or code in TEMPLATE_CONSISTENCY_PARSE_OR_SHAPE_CODES
        ):
            return "advisory"
        if code in TEMPLATE_CONSISTENCY_HARD_CONFLICT_CODES:
            return "advisory"
        if code in TEMPLATE_CONSISTENCY_MISSING_CONTEXT_CODES or code in TEMPLATE_CONSISTENCY_PARSE_OR_SHAPE_CODES:
            return "advisory"
        return "advisory"

    if code in TEMPLATE_CONSISTENCY_HARD_CONFLICT_CODES:
        return "warning"
    if code in TEMPLATE_CONSISTENCY_MISSING_CONTEXT_CODES or code in TEMPLATE_CONSISTENCY_PARSE_OR_SHAPE_CODES:
        return "warning"
    return "advisory"


def validate_template_consistency(
    project_dir: Path,
    findings: list[Finding],
    *,
    quality_mode: str,
    profile: str,
) -> dict[str, Any]:
    result = evaluate_template_binding_consistency(project_dir)
    raw_findings = result.get("findings")
    normalized_raw = raw_findings if isinstance(raw_findings, list) else []
    template_findings: list[dict[str, Any]] = []

    for item in normalized_raw:
        if not isinstance(item, dict):
            continue
        code = str(item.get("code") or "").strip()
        path = str(item.get("path") or project_dir)
        message = str(item.get("message") or "").strip()
        if not code or not message:
            continue
        severity = _template_consistency_severity(code, quality_mode, profile)
        emit(findings, severity, code, path, message)
        template_findings.append({"severity": severity, "code": code, "path": path, "message": message})

    warning_codes = [item["code"] for item in template_findings if item["severity"] == "warning"]
    error_codes = [item["code"] for item in template_findings if item["severity"] == "error"]
    advisory_codes = [item["code"] for item in template_findings if item["severity"] == "advisory"]
    metrics_payload = {
        "ok": bool(result.get("ok", True)),
        "binding_present": bool((result.get("metrics") or {}).get("binding_present")),
        "binding_template_id": (result.get("metrics") or {}).get("binding_template_id"),
        "override_reason": (result.get("metrics") or {}).get("override_reason"),
        "warning_count": len(warning_codes),
        "error_count": len(error_codes),
        "advisory_count": len(advisory_codes),
        "warnings": warning_codes,
        "errors": error_codes,
        "advisories": advisory_codes,
        "findings": template_findings,
    }
    return {"template_consistency": metrics_payload}


def _apply_quality_mode(findings: list[Finding], quality_mode: str) -> None:
    mode = _normalize_quality_mode(quality_mode)
    for item in findings:
        if item.code.startswith("overflow-risk-"):
            risk_level = item.code.removeprefix("overflow-risk-").strip().lower()
            if risk_level == "low":
                item.severity = "advisory"
            elif risk_level in {"medium", "high"}:
                item.severity = "warning"
            continue
        if mode == "dev-fast":
            if item.severity == "warning" and item.code in DEV_FAST_WARNING_TO_ADVISORY_CODES:
                item.severity = "advisory"
        elif mode == "premium":
            if item.severity == "advisory" and item.code in PREMIUM_ADVISORY_TO_WARNING_CODES:
                item.severity = "warning"
            elif item.severity == "warning" and item.code in PREMIUM_WARNING_TO_ERROR_CODES:
                item.severity = "error"


def _apply_visual_critic_gate(findings: list[Finding], profile: str) -> list[str]:
    normalized_profile = str(profile or "").strip().lower()
    if normalized_profile != "proposal_consulting":
        return []
    escalated_codes: set[str] = set()
    for item in findings:
        if item.code in PROPOSAL_CRITIC_BLOCKING_CODES and item.severity in {"advisory", "warning"}:
            item.severity = "error"
            escalated_codes.add(item.code)
    return sorted(escalated_codes)


def _classify_overflow_risk(ratio: float, profile: str) -> str | None:
    policy = resolve_profile_policy(profile)
    if ratio >= policy.overflow_risk_high_ratio:
        return "high"
    if ratio >= policy.overflow_risk_medium_ratio:
        return "medium"
    if ratio >= policy.overflow_risk_low_ratio:
        return "low"
    return None


def _hex_to_linear_rgb(hex_color: str) -> tuple[float, float, float] | None:
    s = hex_color.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(ch * 2 for ch in s)
    if len(s) != 6 or not re.fullmatch(r"[0-9a-fA-F]{6}", s):
        return None
    r = int(s[0:2], 16) / 255.0
    g = int(s[2:4], 16) / 255.0
    b = int(s[4:6], 16) / 255.0

    def lin(c: float) -> float:
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    return lin(r), lin(g), lin(b)


def _relative_luminance(rgb: tuple[float, float, float]) -> float:
    r, g, b = rgb
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_ratio(a: tuple[float, float, float], b: tuple[float, float, float]) -> float:
    la = _relative_luminance(a)
    lb = _relative_luminance(b)
    light, dark = max(la, lb), min(la, lb)
    return (light + 0.05) / (dark + 0.05)


def validate_theme_readability(
    project_dir: Path,
    findings: list[Finding],
    *,
    profile: str = "presentation",
) -> dict[str, Any]:
    policy = resolve_profile_policy(profile)
    theme = Theme.from_design_spec(project_dir / "design_spec.md")
    text_rgb = _hex_to_linear_rgb(theme.text)
    bg_rgb = _hex_to_linear_rgb(theme.background)
    card_rgb = _hex_to_linear_rgb(theme.card)
    metrics: dict[str, Any] = {}

    if text_rgb and bg_rgb:
        contrast_bg = _contrast_ratio(text_rgb, bg_rgb)
        metrics["theme_text_background_contrast"] = round(contrast_bg, 2)
        if contrast_bg < policy.min_theme_contrast:
            emit(
                findings,
                "warning",
                "theme-low-contrast-background",
                str(project_dir / "design_spec.md"),
                (
                    f"Theme text/background contrast is low ({contrast_bg:.2f}:1, "
                    f"required>={policy.min_theme_contrast:.1f})."
                ),
            )

    if text_rgb and card_rgb:
        contrast_card = _contrast_ratio(text_rgb, card_rgb)
        metrics["theme_text_card_contrast"] = round(contrast_card, 2)
        if contrast_card < policy.min_theme_contrast:
            emit(
                findings,
                "warning",
                "theme-low-contrast-card",
                str(project_dir / "design_spec.md"),
                f"Theme text/card contrast is low ({contrast_card:.2f}:1, required>={policy.min_theme_contrast:.1f}).",
            )
    return metrics


def validate_design_token_guard(
    project_dir: Path,
    findings: list[Finding],
    *,
    profile: str = "presentation",
) -> dict[str, Any]:
    design_spec_path = project_dir / "design_spec.md"
    if not design_spec_path.exists():
        return {
            "design_token_profile": profile,
            "design_token_checked_count": 0,
            "design_token_invalid_color_count": 0,
            "design_token_missing_core_count": 0,
            "design_token_low_contrast_count": 0,
            "design_token_font_scale_warning_count": 0,
        }

    payload = run_design_token_guard(design_spec_path.read_text(encoding="utf-8"), profile=profile)
    for item in payload.get("findings", []):
        if not isinstance(item, dict):
            continue
        code = str(item.get("code") or "design-token-check")
        severity = str(item.get("severity") or "warning")
        message = str(item.get("message") or "Design token check finding.")
        recommendation = str(item.get("recommendation") or "").strip()
        if recommendation:
            message = f"{message} Recommendation: {recommendation}"
        emit(findings, severity, code, str(design_spec_path), message)

    return {
        "design_token_profile": str(payload.get("profile") or profile),
        "design_token_checked_count": int(payload.get("checked_token_count", 0)),
        "design_token_invalid_color_count": int(payload.get("invalid_color_count", 0)),
        "design_token_missing_core_count": int(payload.get("missing_token_count", 0)),
        "design_token_low_contrast_count": int(payload.get("contrast_warning_count", 0)),
        "design_token_font_scale_warning_count": int(payload.get("font_scale_warning_count", 0)),
    }


def validate_slide_budget(
    slides: list[dict[str, Any]],
    findings: list[Finding],
    *,
    profile: str = "presentation",
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "slide_budget_profile": profile,
        "slide_budget_high_risk_count": 0,
        "slide_budget_medium_risk_count": 0,
        "slide_budget_low_risk_count": 0,
        "slide_budget_by_slide": {},
    }

    by_slide: dict[str, dict[str, Any]] = {}
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        slide_identifier = slide.get("id")
        sid = slide_identifier if isinstance(slide_identifier, int) else index
        budget = build_slide_budget(slide, profile=profile)
        by_slide[str(sid)] = budget
        metrics["slide_budget_profile"] = budget["profile"]
        risk = str(budget.get("risk") or "none")
        if risk == "high":
            metrics["slide_budget_high_risk_count"] += 1
            emit(
                findings,
                "warning",
                "slide-budget-high",
                f"blueprint.json/slides/{sid}",
                (
                    f"Slide {sid} budget high: chars={budget['actual_chars']}/{budget['max_chars']}, "
                    f"text_nodes={budget['actual_text_nodes']}/{budget['max_text_nodes']}, "
                    f"overflow_action={budget['overflow_action']}."
                ),
            )
        elif risk == "medium":
            metrics["slide_budget_medium_risk_count"] += 1
            emit(
                findings,
                "warning",
                "slide-budget-medium",
                f"blueprint.json/slides/{sid}",
                (
                    f"Slide {sid} budget medium: chars={budget['actual_chars']}/{budget['max_chars']}, "
                    f"text_nodes={budget['actual_text_nodes']}/{budget['max_text_nodes']}, "
                    f"overflow_action={budget['overflow_action']}."
                ),
            )
        elif risk == "low":
            metrics["slide_budget_low_risk_count"] += 1
            emit(
                findings,
                "advisory",
                "slide-budget-low",
                f"blueprint.json/slides/{sid}",
                (
                    f"Slide {sid} budget low: chars={budget['actual_chars']}/{budget['max_chars']}, "
                    f"text_nodes={budget['actual_text_nodes']}/{budget['max_text_nodes']}, "
                    f"overflow_action={budget['overflow_action']}."
                ),
            )

    metrics["slide_budget_by_slide"] = by_slide
    return metrics


def _iter_content_text(value: Any, path: str = "content") -> Sequence[tuple[str, str]]:
    results: list[tuple[str, str]] = []
    if isinstance(value, str):
        text = value.strip()
        if text:
            results.append((path, text))
    elif isinstance(value, dict):
        for key, child in value.items():
            if isinstance(key, str):
                results.extend(_iter_content_text(child, f"{path}.{key}"))
    elif isinstance(value, list):
        for index, child in enumerate(value):
            results.extend(_iter_content_text(child, f"{path}[{index}]"))
    return results


def validate_content_coverage(slides: list[dict[str, Any]], findings: list[Finding]) -> dict[str, Any]:
    """Catch generic fallback copy before a deck passes QA as a polished result."""
    placeholder_slides: list[int] = []
    occurrence_count = 0
    examples: dict[str, list[dict[str, str]]] = {}

    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        raw_sid = slide.get("id")
        sid = raw_sid if isinstance(raw_sid, int) else index
        slide_matches: list[dict[str, str]] = []
        for field_path, text in _iter_content_text(slide.get("content", {})):
            for label, pattern in CONTENT_PLACEHOLDER_PATTERNS:
                if pattern.search(text):
                    slide_matches.append({"path": field_path, "pattern": label, "text": text[:80]})
                    break
        if not slide_matches:
            continue
        placeholder_slides.append(sid)
        occurrence_count += len(slide_matches)
        examples[str(sid)] = slide_matches[:5]
        rendered = "; ".join(f"{item['path']}={item['text']!r}" for item in slide_matches[:3])
        emit(
            findings,
            "warning",
            "content-placeholder-fallback",
            f"blueprint.json/slides/{sid}",
            (
                f"Slide {sid} contains placeholder or fallback copy that should be resolved before final delivery: "
                f"{rendered}."
            ),
        )

    return {
        "content_placeholder_slide_count": len(placeholder_slides),
        "content_placeholder_occurrence_count": occurrence_count,
        "content_placeholder_slides": placeholder_slides,
        "content_placeholder_examples": examples,
    }


def load_blueprint(project_dir: Path, findings: list[Finding]) -> dict[str, Any]:
    path = project_dir / "blueprint.json"
    if not path.exists():
        emit(findings, "error", "missing-blueprint", str(path), "blueprint.json is required.")
        return {"slides": []}
    loaded = load_json_object(path, encoding="utf-8-sig")
    if loaded.status in {"invalid", "read-failed"}:
        emit(findings, "error", "invalid-blueprint-json", str(path), loaded.message)
        return {"slides": []}
    if loaded.status == "schema_mismatch":
        emit(findings, "error", "invalid-blueprint-root", str(path), loaded.message)
        return {"slides": []}
    data = loaded.payload if isinstance(loaded.payload, dict) else {}
    for issue in validate_blueprint_schema(data):
        emit(findings, "error", issue.code, issue.path, issue.message)
    return data


def supported_layout_tags(project_dir: Path) -> set[str]:
    renderer = LayoutRenderer(Theme.from_design_spec(project_dir / "design_spec.md"))
    return set(renderer.registry)


def validate_layout_contracts(project_dir: Path, findings: list[Finding]) -> set[str]:
    renderer_tags = supported_layout_tags(project_dir)
    contract_tags = set(layout_tags())
    for tag in sorted(contract_tags - renderer_tags):
        emit(
            findings,
            "error",
            "layout-contract-without-renderer",
            "scripts/layout_contracts.py",
            f"{tag} has a contract but no renderer.",
        )
    for tag in sorted(renderer_tags - contract_tags):
        emit(
            findings,
            "error",
            "renderer-without-layout-contract",
            "scripts/render_svg.py",
            f"{tag} has a renderer but no layout contract.",
        )
    return renderer_tags


def validate_blueprint(blueprint: dict[str, Any], findings: list[Finding], tags: set[str]) -> list[dict[str, Any]]:
    slides = blueprint.get("slides")
    if not isinstance(slides, list):
        emit(findings, "error", "missing-slides", "blueprint.json", "blueprint.json must contain a slides array.")
        return []

    seen_ids: set[int] = set()
    valid_slides: list[dict[str, Any]] = []
    for index, slide in enumerate(slides, start=1):
        path = f"blueprint.json/slides/{index}"
        if not isinstance(slide, dict):
            emit(findings, "error", "invalid-slide", path, "Each slide must be an object.")
            continue

        missing = CORE_SLIDE_KEYS - set(slide)
        if missing:
            emit(findings, "error", "missing-slide-key", path, f"Missing keys: {', '.join(sorted(missing))}.")

        slide_id = slide.get("id")
        if not isinstance(slide_id, int):
            emit(findings, "error", "invalid-slide-id", path, "Slide id must be an integer.")
        elif slide_id in seen_ids:
            emit(findings, "error", "duplicate-slide-id", path, f"Duplicate slide id: {slide_id}.")
        else:
            seen_ids.add(slide_id)
            if slide_id != index:
                emit(findings, "warning", "nonsequential-slide-id", path, f"Expected id {index}, got {slide_id}.")

        tag = slide.get("layout_tag")
        if tag not in tags:
            emit(findings, "error", "unsupported-layout", path, f"Unsupported layout_tag: {tag!r}.")

        content = slide.get("content")
        if not isinstance(content, dict):
            emit(findings, "error", "invalid-content", path, "content must be an object.")
            content = {}

        for key in required_content_keys(str(tag)):
            if key not in content or content[key] in (None, "", []):
                emit(findings, "warning", "missing-content-key", path, f"Expected content key for {tag}: {key}.")

        for message in validate_content_shape(str(tag), content):
            emit(findings, "warning", "content-shape", path, f"{tag}: {message}")

        valid_slides.append(slide)

    return valid_slides


def _load_source_manifest_ids(project_dir: Path, findings: list[Finding]) -> tuple[set[str], bool]:
    manifest_path = project_dir / "sources" / "manifest.json"
    loaded = load_json_object(manifest_path, encoding="utf-8")
    if loaded.status == "missing":
        return set(), False
    if loaded.status in {"invalid", "read-failed"}:
        emit(findings, "warning", "invalid-source-manifest-json", str(manifest_path), loaded.message)
        return set(), True
    if loaded.status == "schema_mismatch":
        emit(findings, "warning", "invalid-source-manifest-shape", str(manifest_path), loaded.message)
        return set(), True
    payload = loaded.payload if isinstance(loaded.payload, dict) else {}
    records = payload.get("records") if isinstance(payload, dict) else None
    if not isinstance(records, list):
        emit(
            findings,
            "warning",
            "invalid-source-manifest-shape",
            str(manifest_path),
            "sources/manifest.json must contain records[].",
        )
        return set(), True
    ids = {
        str(item.get("id"))
        for item in records
        if isinstance(item, dict) and isinstance(item.get("id"), str) and item.get("id")
    }
    return ids, True


def validate_blueprint_source_refs(
    project_dir: Path,
    slides: list[dict[str, Any]],
    findings: list[Finding],
) -> dict[str, Any]:
    source_refs_present = False
    resolved_refs = 0
    missing_refs = 0
    manifest_ids, has_manifest = _load_source_manifest_ids(project_dir, findings)

    for slide in slides:
        slide_id = slide.get("id")
        source_refs = slide.get("source_refs")
        if source_refs is None:
            continue
        if not isinstance(source_refs, list):
            continue
        non_empty_refs = [
            ref for ref in source_refs if isinstance(ref, str) and ref.strip()
        ]
        if not non_empty_refs:
            continue
        source_refs_present = True
        for ref in non_empty_refs:
            if ref in manifest_ids:
                resolved_refs += 1
            else:
                missing_refs += 1
                emit(
                    findings,
                    "warning",
                    "missing-source-ref",
                    f"blueprint.json/slides/{slide_id}",
                    f"source_refs includes unknown id {ref!r}; expected an id from sources/manifest.json.",
                )

    if source_refs_present and not has_manifest:
        emit(
            findings,
            "warning",
            "missing-source-manifest",
            str(project_dir / "sources" / "manifest.json"),
            "Blueprint contains source_refs but sources/manifest.json is missing.",
        )

    return {
        "source_refs_present": source_refs_present,
        "source_ref_resolved": resolved_refs,
        "source_ref_missing": missing_refs,
        "source_manifest_ids": len(manifest_ids),
    }


def parse_slide_plan(project_dir: Path, findings: list[Finding]) -> dict[int, list[dict[str, Any]]]:
    path = project_dir / "slide_plan.json"
    loaded = load_json_object(path, encoding="utf-8")
    if loaded.status == "missing":
        return {}
    if loaded.status in {"invalid", "read-failed"}:
        emit(findings, "warning", "invalid-slide-plan-json", str(path), loaded.message)
        return {}
    if loaded.status == "schema_mismatch":
        emit(findings, "warning", "invalid-slide-plan-shape", str(path), loaded.message)
        return {}
    payload = loaded.payload if isinstance(loaded.payload, dict) else {}
    slides = payload.get("slides")
    if not isinstance(slides, list):
        emit(findings, "warning", "invalid-slide-plan-shape", str(path), "slide_plan.json must contain slides list.")
        return {}
    parsed: dict[int, list[dict[str, Any]]] = {}
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_id = slide.get("slide_id", slide.get("id"))
        blocks = slide.get("blocks")
        if not isinstance(slide_id, int) or not isinstance(blocks, list):
            continue
        valid_blocks: list[dict[str, Any]] = []
        for block in blocks:
            if not isinstance(block, dict):
                continue
            box = block.get("box")
            if isinstance(box, list) and len(box) == 4 and all(isinstance(v, (int, float)) for v in box):
                valid_blocks.append(block)
        parsed[slide_id] = valid_blocks
    return parsed


def _load_slide_plan_contract(project_dir: Path) -> dict[int, dict[str, Any]]:
    path = project_dir / "slide_plan.json"
    loaded = load_json_object(path, encoding="utf-8")
    if loaded.status != "ok":
        return {}
    payload = loaded.payload if isinstance(loaded.payload, dict) else {}
    slides = payload.get("slides")
    if not isinstance(slides, list):
        return {}
    by_id: dict[int, dict[str, Any]] = {}
    for row in slides:
        if not isinstance(row, dict):
            continue
        slide_id = row.get("slide_id", row.get("id"))
        if not isinstance(slide_id, int):
            continue
        by_id[slide_id] = row
    return by_id


def _load_slide_visual_plan_contract(project_dir: Path) -> dict[int, dict[str, Any]]:
    path = project_dir / "slide_visual_plan.json"
    loaded = load_json_object(path, encoding="utf-8-sig")
    if loaded.status != "ok":
        return {}
    payload = loaded.payload if isinstance(loaded.payload, dict) else {}
    slides = payload.get("slides")
    if not isinstance(slides, list):
        return {}
    by_id: dict[int, dict[str, Any]] = {}
    for row in slides:
        if not isinstance(row, dict):
            continue
        slide_id = row.get("slide_id", row.get("id"))
        if not isinstance(slide_id, int):
            continue
        by_id[slide_id] = row
    return by_id


def _collect_slide_text_fragments(value: Any) -> list[str]:
    if isinstance(value, str):
        text = " ".join(value.split()).strip()
        return [text] if text else []
    if isinstance(value, dict):
        dict_output: list[str] = []
        for item in value.values():
            dict_output.extend(_collect_slide_text_fragments(item))
        return dict_output
    if isinstance(value, list):
        list_output: list[str] = []
        for item in value:
            list_output.extend(_collect_slide_text_fragments(item))
        return list_output
    return []


def _is_high_density_planning_slide(slide: dict[str, Any]) -> bool:
    if is_dense_slide(slide):
        return True
    content_density = str(slide.get("content_density") or "").strip().lower()
    if content_density in {"high", "dense"}:
        return True
    layout_tag = str(slide.get("layout_tag") or "")
    if layout_tag.startswith(("Data-", "Chart-", "Roadmap-", "Flow-", "Grid-")):
        return True
    content = slide.get("content")
    text_total = len(" ".join(_collect_slide_text_fragments(content)))
    return text_total >= 380


def _find_missing_planning_fields(slide_contract: dict[str, Any]) -> list[str]:
    missing: list[str] = []
    for field in sorted(PLANNING_REQUIRED_FIELDS):
        value = slide_contract.get(field)
        if value in (None, "", [], {}):
            missing.append(field)
    return missing


def _load_svg_text_blob(svg_dir: Path, slide_id: int) -> str:
    path = svg_dir / f"slide_{slide_id:02d}.svg"
    if not path.exists():
        return ""
    try:
        raw = path.read_text(encoding="utf-8-sig")
    except Exception:
        raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        root = ET.fromstring(raw)
    except ET.ParseError:
        return raw
    visible_text = " ".join(
        " ".join(" ".join(node.itertext()).split()) for node in root.iter() if local_name(node.tag) == "text"
    )
    return f"{raw}\n{visible_text}"


def _normalize_text_fragment(value: str | None) -> str:
    return " ".join(str(value or "").split()).strip()


def _parse_slot_index(slot_id: str | None, slot_prefix: str) -> int | None:
    text = str(slot_id or "").strip()
    if not text.startswith(f"{slot_prefix}-"):
        return None
    suffix = text.removeprefix(f"{slot_prefix}-")
    if not suffix.isdigit():
        return None
    return int(suffix)


def _format_box_values(values: Sequence[float]) -> str:
    return ",".join(f"{float(value):g}" for value in values)


def _parse_data_slot_box(raw_value: str | None) -> tuple[float, float, float, float] | None:
    if not raw_value:
        return None
    parts = [part.strip() for part in str(raw_value).split(",")]
    if len(parts) != 4:
        return None
    try:
        parsed_values = [float(part) for part in parts]
    except ValueError:
        return None
    return (parsed_values[0], parsed_values[1], parsed_values[2], parsed_values[3])


def _boxes_match(expected: Sequence[float], actual: Sequence[float], *, tolerance: float = 0.01) -> bool:
    if len(expected) != 4 or len(actual) != 4:
        return False
    return all(math.isclose(float(expected[idx]), float(actual[idx]), abs_tol=tolerance) for idx in range(4))


def _load_svg_text_nodes(svg_path: Path) -> list[dict[str, str]] | None:
    if not svg_path.exists():
        return []
    try:
        root = ET.fromstring(svg_path.read_text(encoding="utf-8-sig"))
    except (ET.ParseError, OSError):
        return None
    text_nodes: list[dict[str, str]] = []
    for elem in root.iter():
        if local_name(elem.tag) != "text":
            continue
        text_nodes.append(
            {
                "slot_id": str(elem.get("data-slot-id") or "").strip(),
                "slot_box": str(elem.get("data-slot-box") or "").strip(),
                "structural_anchor": str(elem.get("data-structural-anchor") or "").strip(),
                "text": _normalize_text_fragment(" ".join(elem.itertext())),
            }
        )
    return text_nodes


def _load_svg_elements(svg_path: Path) -> list[dict[str, str]] | None:
    if not svg_path.exists():
        return []
    try:
        root = ET.fromstring(svg_path.read_text(encoding="utf-8-sig"))
    except (ET.ParseError, OSError):
        return None
    elements: list[dict[str, str]] = []
    for elem in root.iter():
        elements.append(
            {
                "tag": local_name(elem.tag),
                "slot_id": str(elem.get("data-slot-id") or "").strip(),
                "slot_box": str(elem.get("data-slot-box") or "").strip(),
                "structural_anchor": str(elem.get("data-structural-anchor") or "").strip(),
                "text": _normalize_text_fragment(" ".join(elem.itertext())),
            }
        )
    return elements


def _slot_ids_for_prefix(slot_ids: Iterable[str | None], slot_prefix: str) -> list[str]:
    return sorted(
        {
            str(slot_id or "").strip()
            for slot_id in slot_ids
            if _parse_slot_index(slot_id, slot_prefix) is not None
        },
        key=lambda value: _parse_slot_index(value, slot_prefix) or 0,
    )


def _missing_contiguous_slot_ids(
    slot_ids: Sequence[str],
    slot_prefix: str,
    *,
    expected_count: int | None = None,
) -> list[str]:
    present_slot_ids = _slot_ids_for_prefix(slot_ids, slot_prefix)
    if not present_slot_ids and not expected_count:
        return []
    present_indices = [_parse_slot_index(slot_id, slot_prefix) or 0 for slot_id in present_slot_ids]
    max_index = max([*present_indices, int(expected_count or 0)])
    expected_slot_ids = [f"{slot_prefix}-{index}" for index in range(1, max_index + 1)]
    present_slot_id_set = set(present_slot_ids)
    return [slot_id for slot_id in expected_slot_ids if slot_id not in present_slot_id_set]


def _missing_expected_texts(expected_texts: Sequence[str], actual_text_blob: str) -> list[str]:
    normalized_text_blob = _normalize_text_fragment(actual_text_blob)
    compact_text_blob = _compact_claim_text(normalized_text_blob)
    missing_texts: list[str] = []
    for expected_text in expected_texts:
        compact_text = _compact_claim_text(expected_text)
        if expected_text not in normalized_text_blob and compact_text not in compact_text_blob:
            missing_texts.append(expected_text)
    return missing_texts


def _text_node_has_anchor(node: dict[str, str], anchor: str) -> bool:
    raw = str(node.get("structural_anchor") or "").strip()
    return anchor in {token.strip() for token in raw.split() if token.strip()}


def _architecture_expected_slot_texts(content: dict[str, Any], content_key: str) -> list[str]:
    raw_items = content.get(content_key)
    if not isinstance(raw_items, list):
        return []
    expected: list[str] = []
    for item in raw_items:
        if content_key == "core_modules":
            if isinstance(item, dict):
                title = _normalize_text_fragment(str(item.get("title") or ""))
                if title:
                    expected.append(title)
            continue
        if isinstance(item, str):
            text = _normalize_text_fragment(item)
            if text:
                expected.append(text)
    return expected


def _roadmap_expected_phase_texts(content: dict[str, Any]) -> list[str]:
    raw_items = content.get("phases")
    if not isinstance(raw_items, list):
        return []
    expected: list[str] = []
    for item in raw_items:
        if not isinstance(item, dict):
            continue
        for key in ("phase", "title", "body"):
            text = _normalize_text_fragment(str(item.get(key) or ""))
            if text:
                expected.append(text)
    return expected


def _comparison_expected_row_texts(content: dict[str, Any]) -> list[str]:
    raw_items = content.get("rows")
    if not isinstance(raw_items, list):
        return []
    expected: list[str] = []
    for item in raw_items:
        if isinstance(item, dict):
            for key in ("dimension", "left", "right"):
                text = _normalize_text_fragment(str(item.get(key) or ""))
                if text:
                    expected.append(text)
            continue
        text = _normalize_text_fragment(str(item))
        if text:
            expected.append(text)
    return expected


def check_comparison_matrix_summarybar_svg_contract(
    project_dir: Path,
    slides: list[dict[str, Any]],
    svg_dir: Path,
    slide_plan: dict[int, list[dict[str, Any]]],
    findings: list[Finding],
    *,
    slide_id: int | None = None,
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "comparison_svg_contract_checked_slides": 0,
        "comparison_svg_contract_applicable_slides": 0,
        "comparison_svg_contract_not_applicable_slides": 0,
        "comparison_svg_contract_missing_slot_count": 0,
        "comparison_svg_contract_missing_anchor_count": 0,
        "comparison_svg_contract_missing_text_count": 0,
        "comparison_svg_contract_skipped_parse_count": 0,
        "comparison_svg_contract_checked_slide_ids": [],
        "comparison_svg_contract_applicable_slide_ids": [],
        "comparison_svg_contract_not_applicable_slide_ids": [],
    }

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        sid = slide.get("id")
        if not isinstance(sid, int):
            continue
        if slide_id is not None and sid != slide_id:
            continue
        if str(slide.get("layout_tag") or "").strip() != "Comparison-Matrix-SummaryBar":
            continue

        metrics["comparison_svg_contract_checked_slides"] += 1
        metrics["comparison_svg_contract_checked_slide_ids"].append(sid)
        content = slide.get("content")
        if not isinstance(content, dict):
            content = {}

        plan_blocks = slide_plan.get(sid, [])
        plan_row_slot_ids = _slot_ids_for_prefix(
            [
                str(block.get("id") or "").strip()
                for block in plan_blocks
                if isinstance(block, dict)
            ],
            COMPARISON_ROW_SLOT_PREFIX,
        )
        has_summary_slot = any(
            isinstance(block, dict) and str(block.get("id") or "").strip() == COMPARISON_SUMMARY_SLOT_ID
            for block in plan_blocks
        )
        if not plan_row_slot_ids and not has_summary_slot:
            metrics["comparison_svg_contract_not_applicable_slides"] += 1
            metrics["comparison_svg_contract_not_applicable_slide_ids"].append(sid)
            continue

        metrics["comparison_svg_contract_applicable_slides"] += 1
        metrics["comparison_svg_contract_applicable_slide_ids"].append(sid)

        raw_rows = content.get("rows")
        row_count = len(raw_rows) if isinstance(raw_rows, list) else 0
        missing_plan_row_slots = _missing_contiguous_slot_ids(
            plan_row_slot_ids,
            COMPARISON_ROW_SLOT_PREFIX,
            expected_count=row_count,
        )
        expected_row_slot_ids = _slot_ids_for_prefix(
            [*plan_row_slot_ids, *missing_plan_row_slots],
            COMPARISON_ROW_SLOT_PREFIX,
        )

        if missing_plan_row_slots:
            metrics["comparison_svg_contract_missing_slot_count"] += len(missing_plan_row_slots)
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-row-slot",
                str(project_dir / "slide_plan.json"),
                f"Slide {sid} has non-contiguous comparison row slots in slide_plan: missing {missing_plan_row_slots}.",
            )

        if not has_summary_slot:
            metrics["comparison_svg_contract_missing_slot_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-summary-slot",
                str(project_dir / "slide_plan.json"),
                f"Slide {sid} slide_plan is missing comparison summary slot: {COMPARISON_SUMMARY_SLOT_ID}.",
            )

        svg_path = svg_dir / f"slide_{sid:02d}.svg"
        text_nodes = _load_svg_text_nodes(svg_path)
        if text_nodes is None:
            metrics["comparison_svg_contract_skipped_parse_count"] += 1
            continue

        row_nodes = [
            node
            for node in text_nodes
            if _parse_slot_index(node.get("slot_id"), COMPARISON_ROW_SLOT_PREFIX) is not None
        ]
        summary_nodes = [
            node for node in text_nodes if str(node.get("slot_id") or "").strip() == COMPARISON_SUMMARY_SLOT_ID
        ]
        found_row_slot_ids = set(
            _slot_ids_for_prefix((node.get("slot_id") for node in row_nodes), COMPARISON_ROW_SLOT_PREFIX)
        )
        missing_svg_row_slots = [slot_id for slot_id in expected_row_slot_ids if slot_id not in found_row_slot_ids]
        if missing_svg_row_slots:
            metrics["comparison_svg_contract_missing_slot_count"] += len(missing_svg_row_slots)
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-row-slot",
                str(svg_path),
                f"Slide {sid} SVG is missing comparison row slot(s): {missing_svg_row_slots}.",
            )

        if not summary_nodes:
            metrics["comparison_svg_contract_missing_slot_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-summary-slot",
                str(svg_path),
                f"Slide {sid} SVG is missing comparison summary slot: {COMPARISON_SUMMARY_SLOT_ID}.",
            )

        if not any(_text_node_has_anchor(node, COMPARISON_DIMENSION_HEADER_ANCHOR) for node in text_nodes):
            metrics["comparison_svg_contract_missing_anchor_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-dimension-anchor",
                str(svg_path),
                f"Slide {sid} SVG is missing comparison dimension header anchor.",
            )

        option_anchor_specs = (
            (COMPARISON_OPTION_LEFT_ANCHOR, "left option"),
            (COMPARISON_OPTION_RIGHT_ANCHOR, "right option"),
        )
        for anchor_name, label in option_anchor_specs:
            if any(_text_node_has_anchor(node, anchor_name) for node in text_nodes):
                continue
            metrics["comparison_svg_contract_missing_anchor_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-option-anchor",
                str(svg_path),
                f"Slide {sid} SVG is missing comparison {label} anchor.",
            )

        if not any(_text_node_has_anchor(node, COMPARISON_RECOMMENDATION_ANCHOR) for node in text_nodes):
            metrics["comparison_svg_contract_missing_anchor_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-recommendation-anchor",
                str(svg_path),
                f"Slide {sid} SVG is missing comparison recommendation anchor.",
            )

        comparison_text_blob = " ".join(node.get("text") or "" for node in text_nodes)
        for expected_text in _missing_expected_texts(_comparison_expected_row_texts(content), comparison_text_blob):
            metrics["comparison_svg_contract_missing_text_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-row-text",
                str(svg_path),
                f"Slide {sid} SVG is missing comparison row text: {expected_text}.",
            )

        for option_key in ("left_title", "right_title"):
            option_text = _normalize_text_fragment(str(content.get(option_key) or ""))
            if not option_text:
                continue
            if _missing_expected_texts([option_text], comparison_text_blob):
                metrics["comparison_svg_contract_missing_text_count"] += 1
                emit(
                    findings,
                    "error",
                    "contract-svg-comparison-missing-option-text",
                    str(svg_path),
                    f"Slide {sid} SVG is missing comparison option text: {option_text}.",
                )

        summary_text = _normalize_text_fragment(str(content.get("summary") or ""))
        if summary_text and _missing_expected_texts([summary_text], comparison_text_blob):
            metrics["comparison_svg_contract_missing_text_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-comparison-missing-summary-text",
                str(svg_path),
                f"Slide {sid} SVG is missing comparison summary text: {summary_text}.",
            )

    return metrics


def _core_orbit_expected_satellite_texts(content: dict[str, Any]) -> list[str]:
    raw_rows = content.get("rows")
    if not isinstance(raw_rows, list):
        return []
    expected: list[str] = []
    for row in raw_rows:
        if isinstance(row, dict):
            for key in ("dimension", "left", "right"):
                text = _normalize_text_fragment(str(row.get(key) or ""))
                if text:
                    expected.append(text)
            continue
        if isinstance(row, str):
            text = _normalize_text_fragment(row)
            if text:
                expected.append(text)
    return expected


def _slot_center(block: dict[str, Any]) -> tuple[float, float] | None:
    box = block.get("box")
    if not (isinstance(box, list) and len(box) == 4 and all(isinstance(value, (int, float)) for value in box)):
        return None
    return float(box[0]) + float(box[2]) / 2.0, float(box[1]) + float(box[3]) / 2.0


def _slot_area(block: dict[str, Any]) -> float:
    box = block.get("box")
    if not (isinstance(box, list) and len(box) == 4 and all(isinstance(value, (int, float)) for value in box)):
        return 0.0
    return float(box[2]) * float(box[3])


def _check_core_orbit_layout_plan(
    project_dir: Path,
    sid: int,
    content: dict[str, Any],
    plan_blocks: list[dict[str, Any]],
    findings: list[Finding],
    metrics: dict[str, Any],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[str], Any, bool]:
    core_blocks = [block for block in plan_blocks if str(block.get("id") or "").strip() == CORE_ORBIT_CORE_SLOT_ID]
    satellite_blocks = [
        block
        for block in plan_blocks
        if _parse_slot_index(str(block.get("id") or "").strip(), CORE_ORBIT_SATELLITE_SLOT_PREFIX) is not None
    ]
    if not plan_blocks:
        metrics["core_orbit_svg_contract_missing_slot_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-missing-layout-plan",
            str(project_dir / "slide_plan.json"),
            f"Slide {sid} is explicit core_orbit_relationship but has no layout plan blocks.",
        )
        return core_blocks, satellite_blocks, [], None, False

    if len(core_blocks) != 1:
        metrics["core_orbit_svg_contract_missing_slot_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-core-slot-count",
            str(project_dir / "slide_plan.json"),
            f"Slide {sid} slide_plan must contain exactly one core-node slot, got {len(core_blocks)}.",
        )
    core_block = core_blocks[0] if core_blocks else {}
    core_center = _slot_center(core_block)
    content_rows = content.get("rows")
    if not isinstance(content_rows, list) or not 4 <= len(content_rows) <= 6:
        metrics["core_orbit_svg_contract_shape_error_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-invalid-content-rows",
            str(project_dir / "blueprint.json"),
            f"Slide {sid} core_orbit_relationship requires 4 to 6 content rows.",
        )
    expected_satellite_count = len(content_rows) if isinstance(content_rows, list) else 0
    expected_satellite_ids = [f"{CORE_ORBIT_SATELLITE_SLOT_PREFIX}-{idx}" for idx in range(1, expected_satellite_count + 1)]
    raw_plan_satellite_ids = [str(block.get("id") or "").strip() for block in satellite_blocks]
    duplicate_satellite_ids = sorted(
        {
            slot_id
            for slot_id in raw_plan_satellite_ids
            if raw_plan_satellite_ids.count(slot_id) > 1
        },
        key=lambda value: _parse_slot_index(value, CORE_ORBIT_SATELLITE_SLOT_PREFIX) or 0,
    )
    if duplicate_satellite_ids:
        metrics["core_orbit_svg_contract_shape_error_count"] += len(duplicate_satellite_ids)
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-duplicate-satellite-slot",
            str(project_dir / "slide_plan.json"),
            f"Slide {sid} slide_plan has duplicate satellite slot(s): {duplicate_satellite_ids}.",
        )
    plan_satellite_ids = _slot_ids_for_prefix(
        raw_plan_satellite_ids,
        CORE_ORBIT_SATELLITE_SLOT_PREFIX,
    )
    extra_satellite_slots = [slot_id for slot_id in plan_satellite_ids if slot_id not in expected_satellite_ids]
    if extra_satellite_slots:
        metrics["core_orbit_svg_contract_shape_error_count"] += len(extra_satellite_slots)
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-extra-satellite-slot",
            str(project_dir / "slide_plan.json"),
            f"Slide {sid} slide_plan has extra satellite slot(s): {extra_satellite_slots}.",
        )
    missing_satellite_slots = [slot_id for slot_id in expected_satellite_ids if slot_id not in plan_satellite_ids]
    if missing_satellite_slots:
        metrics["core_orbit_svg_contract_missing_slot_count"] += len(missing_satellite_slots)
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-missing-satellite-slot",
            str(project_dir / "slide_plan.json"),
            f"Slide {sid} slide_plan is missing core orbit satellite slot(s): {missing_satellite_slots}.",
        )
    if not 4 <= len(satellite_blocks) <= 6:
        metrics["core_orbit_svg_contract_shape_error_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-satellite-count",
            str(project_dir / "slide_plan.json"),
            f"Slide {sid} core orbit requires 4 to 6 satellite slots, got {len(satellite_blocks)}.",
        )

    if core_block and any(_slot_area(core_block) <= _slot_area(block) for block in satellite_blocks):
        metrics["core_orbit_svg_contract_shape_error_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-core-not-dominant",
            str(project_dir / "slide_plan.json"),
            f"Slide {sid} core-node slot area must be larger than every satellite slot.",
        )

    if core_center and satellite_blocks:
        centers: list[tuple[float, float]] = []
        for block in satellite_blocks:
            center = _slot_center(block)
            if center is not None:
                centers.append(center)
        cx, cy = core_center
        has_left = any(center[0] < cx for center in centers)
        has_right = any(center[0] > cx for center in centers)
        has_top = any(center[1] < cy for center in centers)
        has_bottom = any(center[1] > cy for center in centers)
        if not (has_left and has_right and has_top and has_bottom):
            metrics["core_orbit_svg_contract_shape_error_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-core-orbit-one-sided-slots",
                str(project_dir / "slide_plan.json"),
                f"Slide {sid} satellite slots must cover left/right and top/bottom around core-node.",
            )
    return core_blocks, satellite_blocks, expected_satellite_ids, content_rows, True


def _check_core_orbit_svg_slots_and_text(
    sid: int,
    svg_path: Path,
    content: dict[str, Any],
    content_rows: Any,
    core_blocks: list[dict[str, Any]],
    satellite_blocks: list[dict[str, Any]],
    expected_satellite_ids: list[str],
    text_nodes: list[dict[str, str]],
    findings: list[Finding],
    metrics: dict[str, Any],
) -> None:
    expected_box_by_slot = {
        str(block.get("id") or "").strip(): [float(value) for value in block.get("box", [])]
        for block in [*core_blocks, *satellite_blocks]
        if isinstance(block.get("box"), list)
        and len(block.get("box", [])) == 4
        and all(isinstance(value, (int, float)) for value in block.get("box", []))
    }
    for slot_id, expected_box in expected_box_by_slot.items():
        slot_nodes = [node for node in text_nodes if str(node.get("slot_id") or "").strip() == slot_id]
        if not slot_nodes:
            metrics["core_orbit_svg_contract_missing_slot_count"] += 1
            code = (
                "contract-svg-core-orbit-missing-core-slot"
                if slot_id == CORE_ORBIT_CORE_SLOT_ID
                else "contract-svg-core-orbit-missing-satellite-slot"
            )
            emit(
                findings,
                "error",
                code,
                str(svg_path),
                f"Slide {sid} SVG is missing core orbit slot: {slot_id}.",
            )
            continue
        if not any(
            (actual_box := _parse_data_slot_box(node.get("slot_box"))) is not None
            and _boxes_match(expected_box, actual_box)
            for node in slot_nodes
        ):
            metrics["core_orbit_svg_contract_slot_box_mismatch_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-core-orbit-slot-box-mismatch",
                str(svg_path),
                f"Slide {sid} slot {slot_id} data-slot-box mismatch: expected {_format_box_values(expected_box)}.",
            )

    text_nodes_by_slot: dict[str, list[dict[str, str]]] = {}
    for node in text_nodes:
        slot_id = str(node.get("slot_id") or "").strip()
        if not slot_id:
            continue
        text_nodes_by_slot.setdefault(slot_id, []).append(node)

    summary_text = _normalize_text_fragment(str(content.get("summary") or ""))
    core_text_blob = " ".join(node.get("text") or "" for node in text_nodes_by_slot.get(CORE_ORBIT_CORE_SLOT_ID, []))
    if summary_text and _missing_expected_texts([summary_text], core_text_blob):
        metrics["core_orbit_svg_contract_missing_text_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-missing-core-text",
            str(svg_path),
            f"Slide {sid} SVG is missing core text: {summary_text}.",
        )
    if isinstance(content_rows, list):
        for idx, row in enumerate(content_rows, start=1):
            slot_id = f"{CORE_ORBIT_SATELLITE_SLOT_PREFIX}-{idx}"
            expected_texts = _core_orbit_expected_satellite_texts({"rows": [row]})
            slot_text_blob = " ".join(node.get("text") or "" for node in text_nodes_by_slot.get(slot_id, []))
            for expected_text in _missing_expected_texts(expected_texts, slot_text_blob):
                metrics["core_orbit_svg_contract_missing_text_count"] += 1
                emit(
                    findings,
                    "error",
                    "contract-svg-core-orbit-missing-satellite-text",
                    str(svg_path),
                    f"Slide {sid} SVG slot {slot_id} is missing satellite text: {expected_text}.",
                )

    found_svg_satellite_ids = _slot_ids_for_prefix(text_nodes_by_slot.keys(), CORE_ORBIT_SATELLITE_SLOT_PREFIX)
    extra_svg_satellite_slots = [slot_id for slot_id in found_svg_satellite_ids if slot_id not in expected_satellite_ids]
    if extra_svg_satellite_slots:
        metrics["core_orbit_svg_contract_shape_error_count"] += len(extra_svg_satellite_slots)
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-extra-satellite-slot",
            str(svg_path),
            f"Slide {sid} SVG has extra satellite slot(s): {extra_svg_satellite_slots}.",
        )


def _check_core_orbit_svg_anchors(
    sid: int,
    svg_path: Path,
    elements: list[dict[str, str]],
    satellite_blocks: list[dict[str, Any]],
    findings: list[Finding],
    metrics: dict[str, Any],
) -> None:
    all_core_anchor_elements = [
        elem
        for elem in elements
        if elem.get("tag") != "text" and _text_node_has_anchor(elem, CORE_ORBIT_CORE_SHAPE_ANCHOR)
    ]
    invalid_core_shape_tags = sorted(
        {str(elem.get("tag") or "") for elem in all_core_anchor_elements if elem.get("tag") not in {"circle", "ellipse"}}
    )
    if invalid_core_shape_tags:
        metrics["core_orbit_svg_contract_missing_anchor_count"] += len(invalid_core_shape_tags)
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-invalid-core-shape-tag",
            str(svg_path),
            (
                f"Slide {sid} core-node-shape must be circle or ellipse, "
                f"got {invalid_core_shape_tags}."
            ),
        )
    shape_elements = [elem for elem in all_core_anchor_elements if elem.get("tag") in {"circle", "ellipse"}]
    if len(shape_elements) != 1:
        metrics["core_orbit_svg_contract_missing_anchor_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-core-shape-count",
            str(svg_path),
            f"Slide {sid} SVG must contain exactly one {CORE_ORBIT_CORE_SHAPE_ANCHOR} shape, got {len(shape_elements)}.",
        )

    all_edge_anchor_elements = [
        elem
        for elem in elements
        if elem.get("tag") != "text" and _text_node_has_anchor(elem, CORE_ORBIT_RELATIONSHIP_EDGE_ANCHOR)
    ]
    invalid_edge_tags = sorted(
        {str(elem.get("tag") or "") for elem in all_edge_anchor_elements if elem.get("tag") not in {"line", "path", "polyline"}}
    )
    if invalid_edge_tags:
        metrics["core_orbit_svg_contract_missing_anchor_count"] += len(invalid_edge_tags)
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-invalid-relationship-edge-tag",
            str(svg_path),
            (
                f"Slide {sid} relationship-edge must be line, path, or polyline, "
                f"got {invalid_edge_tags}."
            ),
        )
    edge_elements = [elem for elem in all_edge_anchor_elements if elem.get("tag") in {"line", "path", "polyline"}]
    if len(edge_elements) < len(satellite_blocks):
        metrics["core_orbit_svg_contract_missing_anchor_count"] += 1
        emit(
            findings,
            "error",
            "contract-svg-core-orbit-missing-relationship-edge",
            str(svg_path),
            (
                f"Slide {sid} SVG must contain at least {len(satellite_blocks)} "
                f"{CORE_ORBIT_RELATIONSHIP_EDGE_ANCHOR} anchors, got {len(edge_elements)}."
            ),
        )


def check_core_orbit_relationship_svg_contract(
    project_dir: Path,
    slides: list[dict[str, Any]],
    svg_dir: Path,
    slide_plan: dict[int, list[dict[str, Any]]],
    findings: list[Finding],
    *,
    slide_id: int | None = None,
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "core_orbit_svg_contract_checked_slides": 0,
        "core_orbit_svg_contract_applicable_slides": 0,
        "core_orbit_svg_contract_not_applicable_slides": 0,
        "core_orbit_svg_contract_missing_slot_count": 0,
        "core_orbit_svg_contract_missing_text_count": 0,
        "core_orbit_svg_contract_missing_anchor_count": 0,
        "core_orbit_svg_contract_slot_box_mismatch_count": 0,
        "core_orbit_svg_contract_shape_error_count": 0,
        "core_orbit_svg_contract_skipped_parse_count": 0,
        "core_orbit_svg_contract_checked_slide_ids": [],
        "core_orbit_svg_contract_applicable_slide_ids": [],
        "core_orbit_svg_contract_not_applicable_slide_ids": [],
    }

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        sid = slide.get("id")
        if not isinstance(sid, int):
            continue
        if slide_id is not None and sid != slide_id:
            continue
        if str(slide.get("scene_type") or "").strip() != "core_orbit_relationship":
            continue

        metrics["core_orbit_svg_contract_checked_slides"] += 1
        metrics["core_orbit_svg_contract_checked_slide_ids"].append(sid)
        content = slide.get("content")
        if not isinstance(content, dict):
            content = {}

        plan_blocks = [block for block in slide_plan.get(sid, []) if isinstance(block, dict)]
        metrics["core_orbit_svg_contract_applicable_slides"] += 1
        metrics["core_orbit_svg_contract_applicable_slide_ids"].append(sid)
        core_blocks, satellite_blocks, expected_satellite_ids, content_rows, has_plan = _check_core_orbit_layout_plan(
            project_dir,
            sid,
            content,
            plan_blocks,
            findings,
            metrics,
        )
        if not has_plan:
            continue

        svg_path = svg_dir / f"slide_{sid:02d}.svg"
        text_nodes = _load_svg_text_nodes(svg_path)
        elements = _load_svg_elements(svg_path)
        if text_nodes is None or elements is None:
            metrics["core_orbit_svg_contract_skipped_parse_count"] += 1
            continue

        _check_core_orbit_svg_slots_and_text(
            sid,
            svg_path,
            content,
            content_rows,
            core_blocks,
            satellite_blocks,
            expected_satellite_ids,
            text_nodes,
            findings,
            metrics,
        )
        _check_core_orbit_svg_anchors(
            sid,
            svg_path,
            elements,
            satellite_blocks,
            findings,
            metrics,
        )

    return metrics


def check_roadmap_lane_milestones_svg_contract(
    project_dir: Path,
    slides: list[dict[str, Any]],
    svg_dir: Path,
    slide_plan: dict[int, list[dict[str, Any]]],
    findings: list[Finding],
    *,
    slide_id: int | None = None,
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "roadmap_svg_contract_checked_slides": 0,
        "roadmap_svg_contract_applicable_slides": 0,
        "roadmap_svg_contract_not_applicable_slides": 0,
        "roadmap_svg_contract_missing_slot_count": 0,
        "roadmap_svg_contract_missing_text_count": 0,
        "roadmap_svg_contract_skipped_parse_count": 0,
        "roadmap_svg_contract_checked_slide_ids": [],
        "roadmap_svg_contract_applicable_slide_ids": [],
        "roadmap_svg_contract_not_applicable_slide_ids": [],
    }

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        sid = slide.get("id")
        if not isinstance(sid, int):
            continue
        if slide_id is not None and sid != slide_id:
            continue
        if str(slide.get("layout_tag") or "").strip() != "Roadmap-Lane-Milestones":
            continue

        metrics["roadmap_svg_contract_checked_slides"] += 1
        metrics["roadmap_svg_contract_checked_slide_ids"].append(sid)
        content = slide.get("content")
        if not isinstance(content, dict):
            content = {}

        plan_blocks = slide_plan.get(sid, [])
        plan_phase_slot_ids = _slot_ids_for_prefix(
            [
                str(block.get("id") or "").strip()
                for block in plan_blocks
                if isinstance(block, dict)
            ],
            ROADMAP_PHASE_SLOT_PREFIX,
        )
        has_summary_slot = any(
            isinstance(block, dict) and str(block.get("id") or "").strip() == ROADMAP_SUMMARY_SLOT_ID
            for block in plan_blocks
        )
        if not plan_phase_slot_ids or not has_summary_slot:
            metrics["roadmap_svg_contract_not_applicable_slides"] += 1
            metrics["roadmap_svg_contract_not_applicable_slide_ids"].append(sid)
            continue

        raw_phase_items = content.get("phases")
        phase_count = len(raw_phase_items) if isinstance(raw_phase_items, list) else 0
        missing_plan_phase_slots = _missing_contiguous_slot_ids(
            plan_phase_slot_ids,
            ROADMAP_PHASE_SLOT_PREFIX,
            expected_count=phase_count,
        )
        expected_phase_slot_ids = _slot_ids_for_prefix(
            [*plan_phase_slot_ids, *missing_plan_phase_slots],
            ROADMAP_PHASE_SLOT_PREFIX,
        )

        metrics["roadmap_svg_contract_applicable_slides"] += 1
        metrics["roadmap_svg_contract_applicable_slide_ids"].append(sid)

        if missing_plan_phase_slots:
            metrics["roadmap_svg_contract_missing_slot_count"] += len(missing_plan_phase_slots)
            emit(
                findings,
                "error",
                "contract-svg-roadmap-missing-phase-slot",
                str(project_dir / "slide_plan.json"),
                f"Slide {sid} has non-contiguous roadmap phase slots in slide_plan: missing {missing_plan_phase_slots}.",
            )

        svg_path = svg_dir / f"slide_{sid:02d}.svg"
        text_nodes = _load_svg_text_nodes(svg_path)
        if text_nodes is None:
            metrics["roadmap_svg_contract_skipped_parse_count"] += 1
            continue

        phase_nodes = [
            node
            for node in text_nodes
            if _parse_slot_index(node.get("slot_id"), ROADMAP_PHASE_SLOT_PREFIX) is not None
        ]
        summary_nodes = [
            node for node in text_nodes if str(node.get("slot_id") or "").strip() == ROADMAP_SUMMARY_SLOT_ID
        ]
        found_phase_slot_ids = set(
            _slot_ids_for_prefix((node.get("slot_id") for node in phase_nodes), ROADMAP_PHASE_SLOT_PREFIX)
        )
        missing_svg_phase_slots = [slot_id for slot_id in expected_phase_slot_ids if slot_id not in found_phase_slot_ids]
        if missing_svg_phase_slots:
            metrics["roadmap_svg_contract_missing_slot_count"] += len(missing_svg_phase_slots)
            emit(
                findings,
                "error",
                "contract-svg-roadmap-missing-phase-slot",
                str(svg_path),
                f"Slide {sid} SVG is missing roadmap phase slot(s): {missing_svg_phase_slots}.",
            )

        if not summary_nodes:
            metrics["roadmap_svg_contract_missing_slot_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-roadmap-missing-summary-slot",
                str(svg_path),
                f"Slide {sid} SVG is missing roadmap summary slot: {ROADMAP_SUMMARY_SLOT_ID}.",
            )

        roadmap_text_blob = " ".join(node.get("text") or "" for node in [*phase_nodes, *summary_nodes])
        for expected_text in _missing_expected_texts(_roadmap_expected_phase_texts(content), roadmap_text_blob):
            metrics["roadmap_svg_contract_missing_text_count"] += 1
            emit(
                findings,
                "error",
                "contract-svg-roadmap-missing-phase-text",
                str(svg_path),
                f"Slide {sid} SVG is missing roadmap phase text: {expected_text}.",
            )

        summary_text = _normalize_text_fragment(str(content.get("summary") or ""))
        if summary_text:
            if _missing_expected_texts([summary_text], roadmap_text_blob):
                metrics["roadmap_svg_contract_missing_text_count"] += 1
                emit(
                    findings,
                    "error",
                    "contract-svg-roadmap-missing-summary-text",
                    str(svg_path),
                    f"Slide {sid} SVG is missing roadmap summary text: {summary_text}.",
                )

    return metrics


def check_architecture_three_zones_svg_contract(
    project_dir: Path,
    slides: list[dict[str, Any]],
    svg_dir: Path,
    slide_plan: dict[int, list[dict[str, Any]]],
    findings: list[Finding],
    *,
    slide_id: int | None = None,
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "architecture_svg_contract_checked_slides": 0,
        "architecture_svg_contract_applicable_slides": 0,
        "architecture_svg_contract_not_applicable_slides": 0,
        "architecture_svg_contract_applicable_group_count": 0,
        "architecture_svg_contract_missing_slot_count": 0,
        "architecture_svg_contract_missing_text_count": 0,
        "architecture_svg_contract_missing_title_count": 0,
        "architecture_svg_contract_slot_box_mismatch_count": 0,
        "architecture_svg_contract_skipped_parse_count": 0,
        "architecture_svg_contract_checked_slide_ids": [],
        "architecture_svg_contract_applicable_slide_ids": [],
        "architecture_svg_contract_not_applicable_slide_ids": [],
    }

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        sid = slide.get("id")
        if not isinstance(sid, int):
            continue
        if slide_id is not None and sid != slide_id:
            continue
        if str(slide.get("layout_tag") or "").strip() != "Architecture-Three-Zones":
            continue

        metrics["architecture_svg_contract_checked_slides"] += 1
        metrics["architecture_svg_contract_checked_slide_ids"].append(sid)
        content = slide.get("content")
        if not isinstance(content, dict):
            content = {}
        slot_groups: list[dict[str, Any]] = []
        for group in ARCHITECTURE_SLOT_GROUPS:
            expected_box_by_slot: dict[str, list[float]] = {}
            for block in slide_plan.get(sid, []):
                if not isinstance(block, dict):
                    continue
                slot_id = str(block.get("id") or "").strip()
                if _parse_slot_index(slot_id, str(group["slot_prefix"])) is None:
                    continue
                box = block.get("box")
                if isinstance(box, list) and len(box) == 4 and all(isinstance(value, (int, float)) for value in box):
                    expected_box_by_slot[slot_id] = [float(value) for value in box]
            if expected_box_by_slot:
                expected_texts = _architecture_expected_slot_texts(content, str(group["content_key"]))
                slot_groups.append(
                    {
                        "slot_prefix": str(group["slot_prefix"]),
                        "missing_slot_code": str(group["missing_slot_code"]),
                        "missing_text_code": str(group["missing_text_code"]),
                        "label": str(group["label"]),
                        "strict_slot_ids": bool(group["strict_slot_ids"]),
                        "require_slot_box": bool(group["require_slot_box"]),
                        "expected_box_by_slot": expected_box_by_slot,
                        "expected_texts": expected_texts,
                    }
                )

        if not slot_groups:
            metrics["architecture_svg_contract_not_applicable_slides"] += 1
            metrics["architecture_svg_contract_not_applicable_slide_ids"].append(sid)
            continue

        metrics["architecture_svg_contract_applicable_slides"] += 1
        metrics["architecture_svg_contract_applicable_slide_ids"].append(sid)
        metrics["architecture_svg_contract_applicable_group_count"] += len(slot_groups)

        svg_path = svg_dir / f"slide_{sid:02d}.svg"
        text_nodes = _load_svg_text_nodes(svg_path)
        if text_nodes is None:
            metrics["architecture_svg_contract_skipped_parse_count"] += 1
            continue

        for group in slot_groups:
            slot_prefix = str(group["slot_prefix"])
            raw_expected_box_by_slot = group.get("expected_box_by_slot")
            if not isinstance(raw_expected_box_by_slot, dict):
                continue
            expected_box_by_slot = {
                str(key): value
                for key, value in raw_expected_box_by_slot.items()
                if isinstance(key, str) and isinstance(value, list)
            }
            expected_slot_ids = _slot_ids_for_prefix(expected_box_by_slot.keys(), slot_prefix)
            if bool(group["strict_slot_ids"]):
                missing_plan_slots = _missing_contiguous_slot_ids(expected_slot_ids, slot_prefix)
                if missing_plan_slots:
                    metrics["architecture_svg_contract_missing_slot_count"] += len(missing_plan_slots)
                    emit(
                        findings,
                        "error",
                        str(group["missing_slot_code"]),
                        str(project_dir / "slide_plan.json"),
                        (
                            f"Slide {sid} has non-contiguous {group['label']} slots in slide_plan: "
                            f"missing {missing_plan_slots}."
                        ),
                    )

            group_text_nodes = [
                node
                for node in text_nodes
                if _parse_slot_index(node.get("slot_id"), slot_prefix) is not None
            ]
            found_slot_ids = set(_slot_ids_for_prefix((node.get("slot_id") for node in group_text_nodes), slot_prefix))
            missing_svg_slots: list[str] = []
            if bool(group["strict_slot_ids"]):
                missing_svg_slots = [slot_id for slot_id in expected_slot_ids if slot_id not in found_slot_ids]
            elif not found_slot_ids:
                missing_svg_slots = [f"{slot_prefix}-*"]
            if missing_svg_slots:
                metrics["architecture_svg_contract_missing_slot_count"] += len(missing_svg_slots)
                emit(
                    findings,
                    "error",
                    str(group["missing_slot_code"]),
                    str(svg_path),
                    f"Slide {sid} SVG is missing {group['label']} slot(s): {missing_svg_slots}.",
                )

            if bool(group["require_slot_box"]):
                for slot_id in expected_slot_ids:
                    slot_nodes = [node for node in group_text_nodes if str(node.get("slot_id") or "").strip() == slot_id]
                    if not slot_nodes:
                        continue
                    expected_box = expected_box_by_slot[slot_id]
                    has_matching_box = False
                    last_raw_slot_box = ""
                    for node in slot_nodes:
                        raw_slot_box = node.get("slot_box")
                        last_raw_slot_box = str(raw_slot_box or "")
                        actual_box = _parse_data_slot_box(raw_slot_box)
                        if actual_box is not None and _boxes_match(expected_box, actual_box):
                            has_matching_box = True
                            break
                    if has_matching_box:
                        continue
                    metrics["architecture_svg_contract_slot_box_mismatch_count"] += 1
                    emit(
                        findings,
                        "error",
                        "contract-svg-architecture-slot-box-mismatch",
                        str(svg_path),
                        (
                            f"Slide {sid} slot {slot_id} data-slot-box mismatch: "
                            f"expected {_format_box_values(expected_box)}, got {last_raw_slot_box or '<missing>'}."
                        ),
                    )

            group_text_blob = " ".join(node.get("text") or "" for node in group_text_nodes)
            raw_expected_texts = group.get("expected_texts")
            if not isinstance(raw_expected_texts, list):
                continue
            missing_expected_texts = _missing_expected_texts(
                [text for text in raw_expected_texts if isinstance(text, str)],
                group_text_blob,
            )
            for expected_text in missing_expected_texts:
                metrics["architecture_svg_contract_missing_text_count"] += 1
                if slot_prefix == "core_modules":
                    metrics["architecture_svg_contract_missing_title_count"] += 1
                emit(
                    findings,
                    "error",
                    str(group["missing_text_code"]),
                    str(svg_path),
                    f"Slide {sid} SVG is missing {group['label']} text: {expected_text}.",
                )

    return metrics


def _compact_claim_text(value: str) -> str:
    return re.sub(r"[^0-9a-zA-Z\u4e00-\u9fff]+", "", str(value or "").lower())


def _char_ngrams(value: str, size: int = 2) -> set[str]:
    compact = _compact_claim_text(value)
    if len(compact) < size:
        return set()
    return {compact[index : index + size] for index in range(0, len(compact) - size + 1)}


def _claim_present_in_svg_text(claim: str, svg_text: str) -> bool:
    normalized = " ".join(claim.split()).strip().lower()
    if not normalized:
        return False
    if normalized in svg_text:
        return True
    compact_claim = _compact_claim_text(normalized)
    compact_svg = _compact_claim_text(svg_text)
    if len(compact_claim) < 16:
        return compact_claim in compact_svg
    claim_grams = _char_ngrams(compact_claim, 2)
    if not claim_grams:
        return False
    svg_grams = _char_ngrams(compact_svg, 2)
    overlap = claim_grams & svg_grams
    return len(overlap) >= 8 and (len(overlap) / max(1, len(claim_grams))) >= 0.55


def validate_planning_consistency(
    project_dir: Path,
    slides: list[dict[str, Any]],
    svg_dir: Path,
    findings: list[Finding],
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "planning_consistency_checked_slides": 0,
        "planning_consistency_dense_slides": 0,
        "planning_consistency_missing_contract_count": 0,
        "planning_consistency_missing_field_count": 0,
        "planning_consistency_missing_claim_count": 0,
    }
    slide_plan_contract = _load_slide_plan_contract(project_dir)
    visual_plan_contract = _load_slide_visual_plan_contract(project_dir)
    if not slide_plan_contract and not visual_plan_contract:
        return metrics

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_id = slide.get("id")
        if not isinstance(slide_id, int):
            continue
        metrics["planning_consistency_checked_slides"] += 1
        if not _is_high_density_planning_slide(slide):
            continue

        metrics["planning_consistency_dense_slides"] += 1
        plan_row = slide_plan_contract.get(slide_id) or visual_plan_contract.get(slide_id)
        if not isinstance(plan_row, dict):
            metrics["planning_consistency_missing_contract_count"] += 1
            emit(
                findings,
                "warning",
                "planning-consistency-missing-contract",
                str(project_dir / "slide_plan.json"),
                (
                    f"Slide {slide_id} is dense/high-risk but no planning contract row was found in "
                    "slide_plan.json or slide_visual_plan.json."
                ),
            )
            continue

        missing_fields = _find_missing_planning_fields(plan_row)
        if missing_fields:
            metrics["planning_consistency_missing_field_count"] += 1
            emit(
                findings,
                "warning",
                "planning-consistency-missing-fields",
                str(project_dir / "slide_plan.json"),
                f"Slide {slide_id} planning contract is missing fields: {missing_fields}.",
            )

        must_keep_claims = plan_row.get("must_keep_claims")
        if not isinstance(must_keep_claims, list) or not must_keep_claims:
            continue
        svg_text = _load_svg_text_blob(svg_dir, slide_id).lower()
        for claim in must_keep_claims:
            if not isinstance(claim, str):
                continue
            normalized = " ".join(claim.split()).strip().lower()
            if len(normalized) < 8:
                continue
            if not _claim_present_in_svg_text(normalized, svg_text):
                metrics["planning_consistency_missing_claim_count"] += 1
                emit(
                    findings,
                    "warning",
                    "planning-consistency-claim-missing",
                    str(svg_dir / f"slide_{slide_id:02d}.svg"),
                    (
                        f"Slide {slide_id} must_keep_claim was not found in target SVG text: "
                        f"{claim[:80]!r}. Avoid dropping core conclusions during repair."
                    ),
                )
                break

    return metrics


def _load_json_file(path: Path, findings: list[Finding], *, code: str) -> dict[str, Any] | None:
    loaded = load_json_object(path, encoding="utf-8-sig")
    if loaded.status == "missing":
        return None
    if loaded.status in {"invalid", "read-failed"}:
        emit(findings, "warning", code, str(path), loaded.message)
        return None
    if loaded.status == "schema_mismatch":
        emit(findings, "warning", code, str(path), loaded.message)
        return None
    payload = loaded.payload
    return payload if isinstance(payload, dict) else None


def _truthy_string(value: Any) -> str:
    if not isinstance(value, str):
        return ""
    text = value.strip()
    if text.lower() in {"", "none", "null", "free-design", "free_design"}:
        return ""
    return text


def _has_nonempty_list(value: Any) -> bool:
    return isinstance(value, list) and any(item for item in value)


def _reference_pack_has_selected_reference(payload: dict[str, Any]) -> bool:
    return any(
        [
            bool(_truthy_string(payload.get("primary_template"))),
            _has_nonempty_list(payload.get("secondary_templates")),
            _has_nonempty_list(payload.get("reference_files")),
            _has_nonempty_list(payload.get("selected_templates")),
            _has_nonempty_list(payload.get("references")),
        ]
    )


def _reference_pack_override_reason(payload: dict[str, Any]) -> str:
    return _truthy_string(payload.get("free_design_override_reason")) or _truthy_string(payload.get("override_reason"))


def _validate_reference_first_gate(
    project_dir: Path,
    reference_pack: dict[str, Any] | None,
    findings: list[Finding],
) -> dict[str, Any]:
    metrics = {
        "reference_first_gate_present": False,
        "reference_pack_has_selected_reference": False,
        "reference_pack_free_design_override_present": False,
        "reference_pack_mode": "",
    }
    if reference_pack is None:
        return metrics

    mode = str(reference_pack.get("mode") or "").strip().lower()
    has_selected_reference = _reference_pack_has_selected_reference(reference_pack)
    override_reason = _reference_pack_override_reason(reference_pack)

    metrics["reference_pack_mode"] = mode
    metrics["reference_pack_has_selected_reference"] = has_selected_reference
    metrics["reference_pack_free_design_override_present"] = bool(override_reason)
    metrics["reference_first_gate_present"] = has_selected_reference or bool(override_reason)

    if mode == "free-design" and not override_reason:
        emit(
            findings,
            "warning",
            "reference-pack-free-design-reason-missing",
            str(project_dir / "reference_pack.json"),
            (
                "reference_pack.json uses free-design mode but does not explain why "
                "template/history references were skipped."
            ),
        )
    elif not has_selected_reference and not override_reason:
        emit(
            findings,
            "warning",
            "reference-pack-empty",
            str(project_dir / "reference_pack.json"),
            "reference_pack.json has no selected templates, reference files, or free_design_override_reason.",
        )
    return metrics


def validate_art_direction_consumption(
    project_dir: Path,
    slides: list[dict[str, Any]],
    findings: list[Finding],
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "art_direction_gate_present": False,
        "art_direction_gate_files": {},
        "style_route_present": False,
        "requires_style_drafts": False,
        "style_drafts_present": False,
        "style_draft_selected": False,
        "style_draft_selected_by": None,
        "slide_visual_plan_coverage": 0.0,
        "prompt_pattern_coverage": 0.0,
        "reference_first_gate_present": False,
        "reference_pack_has_selected_reference": False,
        "reference_pack_free_design_override_present": False,
        "reference_pack_mode": "",
    }
    required = {
        "art_direction.md": project_dir / "art_direction.md",
        "reference_pack.json": project_dir / "reference_pack.json",
        "slide_visual_plan.json": project_dir / "slide_visual_plan.json",
    }
    missing: list[str] = []
    for name, path in required.items():
        exists = path.exists()
        metrics["art_direction_gate_files"][name] = exists
        if not exists:
            missing.append(name)
            emit(
                findings,
                "warning",
                "missing-art-direction-artifact",
                str(path),
                f"Missing required State 2.5 artifact: {name}.",
            )
    metrics["art_direction_gate_present"] = len(missing) == 0

    reference_pack = _load_json_file(required["reference_pack.json"], findings, code="invalid-reference-pack-json")
    metrics.update(_validate_reference_first_gate(project_dir, reference_pack, findings))

    blueprint_slide_ids = set()
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_identifier = slide.get("id")
        if isinstance(slide_identifier, int):
            blueprint_slide_ids.add(slide_identifier)

    plan_payload = _load_json_file(required["slide_visual_plan.json"], findings, code="invalid-slide-visual-plan-json")
    if plan_payload is not None:
        plan_slides = plan_payload.get("slides")
        if isinstance(plan_slides, list):
            prompt_pattern_ok = 0
            prompt_pattern_total = 0
            plan_ids = set()
            for item in plan_slides:
                if not isinstance(item, dict):
                    continue
                plan_slide_id = item.get("slide_id")
                if isinstance(plan_slide_id, int):
                    plan_ids.add(plan_slide_id)
            missing_ids = sorted(blueprint_slide_ids - plan_ids)
            extra_ids = sorted(plan_ids - blueprint_slide_ids)
            if blueprint_slide_ids:
                metrics["slide_visual_plan_coverage"] = round(
                    len(plan_ids & blueprint_slide_ids) / len(blueprint_slide_ids),
                    3,
                )
            if missing_ids:
                emit(
                    findings,
                    "warning",
                    "slide-visual-plan-missing-slide",
                    str(required["slide_visual_plan.json"]),
                    f"slide_visual_plan.json does not cover blueprint slide ids: {missing_ids}.",
                )
            if extra_ids:
                emit(
                    findings,
                    "warning",
                    "slide-visual-plan-extra-slide",
                    str(required["slide_visual_plan.json"]),
                    f"slide_visual_plan.json references non-blueprint slide ids: {extra_ids}.",
                )

            required_pattern_keys = {
                "pattern_id",
                "conclusion_formula",
                "block_structure",
                "composition_cues",
                "anti_patterns",
            }
            for item in plan_slides:
                if not isinstance(item, dict):
                    continue
                if not isinstance(item.get("slide_id"), int):
                    continue
                prompt_pattern_total += 1
                pattern = item.get("page_prompt_pattern")
                if not isinstance(pattern, dict):
                    emit(
                        findings,
                        "warning",
                        "prompt-pattern-missing",
                        str(required["slide_visual_plan.json"]),
                        f"slide_id={item.get('slide_id')} is missing page_prompt_pattern.",
                    )
                    continue
                missing_keys = sorted(required_pattern_keys - set(pattern.keys()))
                if missing_keys:
                    emit(
                        findings,
                        "warning",
                        "prompt-pattern-incomplete",
                        str(required["slide_visual_plan.json"]),
                        f"slide_id={item.get('slide_id')} page_prompt_pattern missing keys: {missing_keys}.",
                    )
                    continue
                prompt_pattern_ok += 1
            if prompt_pattern_total > 0:
                metrics["prompt_pattern_coverage"] = round(prompt_pattern_ok / prompt_pattern_total, 3)
        else:
            emit(
                findings,
                "warning",
                "invalid-slide-visual-plan-shape",
                str(required["slide_visual_plan.json"]),
                "slide_visual_plan.json must contain slides[].",
            )

    style_route_path = project_dir / "style_route.json"
    style_route = _load_json_file(style_route_path, findings, code="invalid-style-route-json")
    if style_route is not None:
        metrics["style_route_present"] = True
        metrics["style_profile"] = style_route.get("style_profile")
        metrics["style_route_confidence"] = style_route.get("confidence")
        requires_style_drafts = bool(style_route.get("requires_style_drafts", False))
        metrics["requires_style_drafts"] = requires_style_drafts
        if requires_style_drafts:
            drafts_path = project_dir / "style_drafts.json"
            drafts_payload = _load_json_file(drafts_path, findings, code="invalid-style-drafts-json")
            if drafts_payload is None:
                emit(
                    findings,
                    "warning",
                    "style-drafts-required",
                    str(drafts_path),
                    "style_route.json requires style drafts, but style_drafts.json is missing or invalid.",
                )
            else:
                metrics["style_drafts_present"] = True
                drafts = drafts_payload.get("drafts")
                if not isinstance(drafts, list) or len(drafts) < 2:
                    emit(
                        findings,
                        "warning",
                        "style-drafts-insufficient",
                        str(drafts_path),
                        "Low-confidence routing requires at least two style drafts.",
                    )
                draft_index: dict[str, str] = {}
                if isinstance(drafts, list):
                    for item in drafts:
                        if not isinstance(item, dict):
                            continue
                        draft_id = str(item.get("draft_id") or "").strip()
                        template_id = str(item.get("template_id") or "").strip()
                        if draft_id and template_id:
                            draft_index[draft_id] = template_id

                selected_template_raw = drafts_payload.get("selected_template")
                selected_draft_id_raw = drafts_payload.get("selected_draft_id")
                selected_template = selected_template_raw.strip() if isinstance(selected_template_raw, str) else ""
                selected_draft_id = selected_draft_id_raw.strip() if isinstance(selected_draft_id_raw, str) else ""
                selected_from_draft = draft_index.get(selected_draft_id, "")

                if not selected_template and not selected_draft_id:
                    emit(
                        findings,
                        "warning",
                        "style-draft-not-selected",
                        str(drafts_path),
                        (
                            "Low-confidence routing requires selected_template or "
                            "selected_draft_id before full-deck execution."
                        ),
                    )
                else:
                    if selected_draft_id and selected_draft_id not in draft_index:
                        emit(
                            findings,
                            "warning",
                            "style-draft-selection-invalid",
                            str(drafts_path),
                            f"selected_draft_id={selected_draft_id!r} does not exist in drafts[].",
                        )
                    if selected_template and selected_template not in set(draft_index.values()):
                        emit(
                            findings,
                            "warning",
                            "style-draft-selection-invalid",
                            str(drafts_path),
                            f"selected_template={selected_template!r} is not present in drafts[].template_id.",
                        )
                    if (
                        selected_draft_id
                        and selected_template
                        and selected_from_draft
                        and selected_from_draft != selected_template
                    ):
                        emit(
                            findings,
                            "warning",
                            "style-draft-selection-mismatch",
                            str(drafts_path),
                            (
                                f"selected_draft_id={selected_draft_id!r} maps to template {selected_from_draft!r}, "
                                f"but selected_template is {selected_template!r}."
                            ),
                        )
                    if selected_draft_id and selected_draft_id in draft_index:
                        metrics["style_draft_selected"] = True
                        metrics["style_draft_selected_by"] = "selected_draft_id"
                    elif selected_template and selected_template in set(draft_index.values()):
                        metrics["style_draft_selected"] = True
                        metrics["style_draft_selected_by"] = "selected_template"
    return metrics


def check_visual_contracts(
    project_dir: Path,
    findings: list[Finding],
    *,
    quality_mode: str,
    profile: str,
) -> dict[str, Any]:
    plan_path = project_dir / "slide_visual_plan.json"
    payload = _load_json_file(plan_path, findings, code="invalid-slide-visual-plan-json")
    metrics: dict[str, Any] = {
        "visual_contract_total": 0,
        "visual_contract_valid": 0,
        "visual_contract_missing_count": 0,
        "visual_contract_invalid_count": 0,
        "visual_contract_coverage": 0.0,
        "execution_policy_total": 0,
        "execution_policy_valid": 0,
        "execution_policy_missing_count": 0,
        "execution_policy_invalid_count": 0,
        "execution_policy_coverage": 0.0,
    }
    if not isinstance(payload, dict):
        return metrics
    raw_slides = payload.get("slides")
    if not isinstance(raw_slides, list):
        return metrics

    for item in raw_slides:
        if not isinstance(item, dict):
            continue
        slide_id = item.get("slide_id")
        metrics["visual_contract_total"] += 1
        metrics["execution_policy_total"] += 1
        execution_policy = item.get("execution_policy")
        policy_is_high_risk_blocking = False
        if not isinstance(execution_policy, dict):
            metrics["execution_policy_missing_count"] += 1
            emit(
                findings,
                "warning",
                "execution-policy-missing",
                str(plan_path),
                f"slide_id={slide_id!r} is missing execution_policy.",
            )
        else:
            policy_missing = sorted(EXECUTION_POLICY_REQUIRED_KEYS - set(execution_policy.keys()))
            expected_rules = execution_policy.get("expected_first_pass_rules")
            if policy_missing or not isinstance(expected_rules, list) or not expected_rules:
                metrics["execution_policy_invalid_count"] += 1
                emit(
                    findings,
                    "warning",
                    "execution-policy-incomplete",
                    str(plan_path),
                    f"slide_id={slide_id!r} execution_policy missing or invalid fields: {policy_missing}.",
                )
            else:
                metrics["execution_policy_valid"] += 1
            policy_is_high_risk_blocking = (
                str(execution_policy.get("risk_level") or "").lower() == "high"
                and str(execution_policy.get("qa_strictness") or "").lower() == "blocking"
            )
        contract = item.get("visual_contract")
        if not isinstance(contract, dict):
            metrics["visual_contract_missing_count"] += 1
            emit(
                findings,
                "warning",
                "visual-contract-missing",
                str(plan_path),
                f"slide_id={slide_id!r} is missing visual_contract.",
            )
            continue

        missing = sorted(VISUAL_CONTRACT_REQUIRED_KEYS - set(contract.keys()))
        if missing:
            metrics["visual_contract_invalid_count"] += 1
            emit(
                findings,
                "warning",
                "visual-contract-incomplete",
                str(plan_path),
                f"slide_id={slide_id!r} visual_contract missing keys: {missing}.",
            )
            continue

        primary_read_path = contract.get("primary_read_path")
        if not isinstance(primary_read_path, list) or not primary_read_path:
            metrics["visual_contract_invalid_count"] += 1
            emit(
                findings,
                "warning",
                "visual-contract-invalid-read-path",
                str(plan_path),
                f"slide_id={slide_id!r} visual_contract primary_read_path must be a non-empty list.",
            )
            continue

        density_budget = contract.get("density_budget")
        if not isinstance(density_budget, dict):
            metrics["visual_contract_invalid_count"] += 1
            emit(
                findings,
                "warning",
                "visual-contract-invalid-density-budget",
                str(plan_path),
                f"slide_id={slide_id!r} visual_contract density_budget must be an object.",
            )
            continue

        anti_patterns = contract.get("anti_patterns")
        if not isinstance(anti_patterns, list):
            metrics["visual_contract_invalid_count"] += 1
            emit(
                findings,
                "warning",
                "visual-contract-invalid-anti-patterns",
                str(plan_path),
                f"slide_id={slide_id!r} visual_contract anti_patterns must be a list.",
            )
            continue

        v2_missing = sorted(VISUAL_CONTRACT_V2_REQUIRED_KEYS - set(contract.keys()))
        v2_invalid = []
        if not isinstance(contract.get("bbox_budget"), dict):
            v2_invalid.append("bbox_budget")
        if not isinstance(contract.get("text_budget"), dict):
            v2_invalid.append("text_budget")
        if not isinstance(contract.get("deterministic_scaffold"), dict):
            v2_invalid.append("deterministic_scaffold")
        if not isinstance(contract.get("must_avoid"), list):
            v2_invalid.append("must_avoid")
        if not isinstance(contract.get("pre_authoring_checks"), list) or not contract.get("pre_authoring_checks"):
            v2_invalid.append("pre_authoring_checks")
        if v2_missing or v2_invalid:
            metrics["visual_contract_invalid_count"] += 1
            is_blocking_v2 = (
                quality_mode in {"release-safe", "premium"}
                and any(token in str(profile).lower() for token in CONSULTING_PROFILE_TOKENS)
                and policy_is_high_risk_blocking
            )
            emit(
                findings,
                "error" if is_blocking_v2 else "warning",
                "visual-contract-v2-incomplete",
                str(plan_path),
                (
                    f"slide_id={slide_id!r} visual_contract v2 fields missing={v2_missing}, "
                    f"invalid={v2_invalid}."
                ),
            )
            if is_blocking_v2:
                continue

        scene_route = item.get("scene_route")
        if isinstance(scene_route, dict):
            route_scene = scene_route.get("scene_type")
            route_strategy = scene_route.get("generation_strategy")
            contract_scene = contract.get("scene_type")
            contract_strategy = contract.get("generation_strategy")
            if (route_scene and contract_scene and route_scene != contract_scene) or (
                route_strategy and contract_strategy and route_strategy != contract_strategy
            ):
                emit(
                    findings,
                    "warning",
                    "visual-contract-scene-mismatch",
                    str(plan_path),
                    (
                        f"slide_id={slide_id!r} scene_route and visual_contract disagree "
                        f"(route={route_scene}/{route_strategy}, contract={contract_scene}/{contract_strategy})."
                    ),
                )
        if isinstance(execution_policy, dict):
            policy_scene = execution_policy.get("scene_type")
            policy_strategy = execution_policy.get("generation_strategy")
            contract_scene = contract.get("scene_type")
            contract_strategy = contract.get("generation_strategy")
            if (policy_scene and contract_scene and policy_scene != contract_scene) or (
                policy_strategy and contract_strategy and policy_strategy != contract_strategy
            ):
                emit(
                    findings,
                    "warning",
                    "visual-contract-policy-mismatch",
                    str(plan_path),
                    (
                        f"slide_id={slide_id!r} execution_policy and visual_contract disagree "
                        f"(policy={policy_scene}/{policy_strategy}, contract={contract_scene}/{contract_strategy})."
                    ),
                )

        metrics["visual_contract_valid"] += 1

    total = int(metrics["visual_contract_total"])
    if total > 0:
        metrics["visual_contract_coverage"] = round(float(metrics["visual_contract_valid"]) / float(total), 3)
    policy_total = int(metrics["execution_policy_total"])
    if policy_total > 0:
        metrics["execution_policy_coverage"] = round(float(metrics["execution_policy_valid"]) / float(policy_total), 3)
    return metrics


def choose_svg_dir(project_dir: Path, requested: str | None) -> Path:
    if requested:
        return project_dir / requested
    for name in ("svg_final", "svg_output"):
        candidate = project_dir / name
        if candidate.exists() and any(candidate.glob("slide_*.svg")):
            return candidate
    return project_dir / "svg_output"


def choose_structural_contract_svg_dir(project_dir: Path, svg_dir: Path) -> Path:
    """Structural contract evidence lives in authoring SVGs before finalize strips custom attrs."""
    authoring_dir = project_dir / "svg_output"
    if authoring_dir == svg_dir:
        return svg_dir
    if authoring_dir.exists() and any(authoring_dir.glob("slide_*.svg")):
        return authoring_dir
    return svg_dir


def _contains_point(box: list[float], x: float, y: float) -> bool:
    bx, by, bw, bh = box
    return bx <= x <= bx + bw and by <= y <= by + bh


def _contains_cjk(text: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def _estimated_text_units(text: str) -> float:
    base = visual_width(text)
    if not _contains_cjk(text):
        return base
    cjk_units = sum(1.0 if "\u4e00" <= ch <= "\u9fff" else 0.55 for ch in text)
    return max(base, cjk_units)


def _latest_semantic_pptx(project_dir: Path, export_mode: str) -> Path | None:
    exports_dir = project_dir / "exports"
    if not exports_dir.exists():
        return None
    candidates = sorted(exports_dir.glob(f"*--{export_mode}--v*--*.pptx"), key=lambda path: path.stat().st_mtime)
    return candidates[-1] if candidates else None


def text_lines(elem: ET.Element) -> list[str]:
    tspans = [child for child in elem if local_name(child.tag) == "tspan"]
    if tspans:
        lines: list[str] = []
        current_parts: list[str] = []
        root_text = (elem.text or "").strip()
        if root_text:
            current_parts.append(root_text)
        for child in tspans:
            starts_new_line = any(child.get(attr) for attr in ("x", "y", "dy"))
            if starts_new_line and current_parts:
                line = "".join(current_parts).strip()
                if line:
                    lines.append(line)
                current_parts = []
            child_text = (child.text or "").strip()
            if child_text:
                current_parts.append(child_text)
            child_tail = (child.tail or "").strip()
            if child_tail:
                current_parts.append(child_tail)
        line = "".join(current_parts).strip()
        if line:
            lines.append(line)
        return lines
    value = (elem.text or "").strip()
    return [value] if value else []


def parse_translate(transform: str | None) -> tuple[float, float]:
    if not transform:
        return 0.0, 0.0
    dx = 0.0
    dy = 0.0
    for match in TRANSLATE_RE.finditer(transform):
        dx += float(match.group(1))
        dy += float(match.group(2) or 0.0)
    return dx, dy


def iter_text_elements_with_offsets(
    elem: ET.Element,
    offset_x: float = 0.0,
    offset_y: float = 0.0,
):
    dx, dy = parse_translate(elem.get("transform"))
    next_x = offset_x + dx
    next_y = offset_y + dy
    if local_name(elem.tag) == "text":
        yield elem, next_x, next_y
    for child in elem:
        yield from iter_text_elements_with_offsets(child, next_x, next_y)


def estimate_text_box(
    elem: ET.Element, offset_x: float = 0.0, offset_y: float = 0.0
) -> tuple[float, float, float, float] | None:
    tspans = [child for child in elem if local_name(child.tag) == "tspan"]
    if tspans and any(any(child.get(attr) for attr in ("x", "y", "dy")) for child in tspans):
        anchor = elem.get("text-anchor", "start")
        inherited_font_size = number(elem.get("font-size"), 18)
        current_x = number(elem.get("x")) + offset_x
        current_y = number(elem.get("y")) + offset_y
        boxes: list[tuple[float, float, float, float]] = []
        root_text = (elem.text or "").strip()
        if root_text:
            boxes.append(_estimate_text_fragment_box(root_text, current_x, current_y, inherited_font_size, anchor))
        for child in tspans:
            if child.get("x") is not None:
                current_x = number(child.get("x")) + offset_x
            if child.get("y") is not None:
                current_y = number(child.get("y")) + offset_y
            elif child.get("dy") is not None:
                current_y += number(child.get("dy"))
            font_size = number(child.get("font-size"), inherited_font_size)
            child_text = (child.text or "").strip()
            if child_text:
                boxes.append(_estimate_text_fragment_box(child_text, current_x, current_y, font_size, anchor))
        if boxes:
            return (
                min(box[0] for box in boxes),
                min(box[1] for box in boxes),
                max(box[2] for box in boxes),
                max(box[3] for box in boxes),
            )

    lines = text_lines(elem)
    if not lines:
        return None
    x = number(elem.get("x")) + offset_x
    y = number(elem.get("y")) + offset_y
    font_size = number(elem.get("font-size"), 18)
    anchor = elem.get("text-anchor", "start")
    width = max(
        1.0,
        max(_estimated_text_units(line) * font_size * (1.0 if _contains_cjk(line) else 0.78) for line in lines),
    )
    line_height = font_size * 0.84
    height = max(line_height, len(lines) * line_height)
    top = y - font_size * 0.72
    if anchor == "middle":
        left = x - width / 2
    elif anchor == "end":
        left = x - width
    else:
        left = x
    return left, top, left + width, top + height


def _estimate_text_fragment_box(
    text: str,
    x: float,
    y: float,
    font_size: float,
    anchor: str,
) -> tuple[float, float, float, float]:
    width = max(1.0, _estimated_text_units(text) * font_size * (1.0 if _contains_cjk(text) else 0.78))
    top = y - font_size * 0.72
    height = font_size * 0.84
    if anchor == "middle":
        left = x - width / 2
    elif anchor == "end":
        left = x - width
    else:
        left = x
    return left, top, left + width, top + height


def validate_svg_file(
    svg_file: Path,
    findings: list[Finding],
    canvas: CanvasContext,
    safe_area_profile: str = "legacy",
    slide_plan_blocks: list[dict[str, Any]] | None = None,
    safe_edge_whitelist: set[str] | None = None,
    profile: str = "presentation",
) -> dict[str, int]:
    metrics = {
        "texts": 0,
        "unsupported_nodes": 0,
        "overflow_risk_high_count": 0,
        "overflow_risk_medium_count": 0,
        "overflow_risk_low_count": 0,
    }
    try:
        tree = ET.parse(svg_file)
    except (ET.ParseError, OSError) as exc:
        emit(findings, "error", "invalid-svg", str(svg_file), f"Could not parse SVG: {exc}")
        return metrics

    root = tree.getroot()
    if local_name(root.tag) != "svg":
        emit(findings, "error", "invalid-svg-root", str(svg_file), "Root element must be <svg>.")
        return metrics

    width = number(root.get("width"), canvas.width)
    height = number(root.get("height"), canvas.height)
    view_box = root.get("viewBox")
    if not math.isclose(width, canvas.width) or not math.isclose(height, canvas.height):
        emit(
            findings,
            "error",
            "invalid-svg-size",
            str(svg_file),
            f"Expected {canvas.width:g}x{canvas.height:g}, got {width:g}x{height:g}.",
        )
    if view_box != canvas.viewbox:
        emit(
            findings,
            "warning",
            "invalid-viewbox",
            str(svg_file),
            f"Expected viewBox {canvas.viewbox!r}, got {view_box!r}.",
        )

    safe_profile = SAFE_AREA_PROFILES.get(safe_area_profile, SAFE_AREA_PROFILES["legacy"])
    safe_left = canvas.safe_x - safe_profile["pad_x"]
    safe_right = canvas.safe_x + canvas.safe_w + safe_profile["pad_x"]
    safe_top = canvas.safe_y - safe_profile["pad_top"]
    safe_bottom = canvas.safe_y + canvas.safe_h + safe_profile["pad_bottom"]

    text_blocks: list[dict[str, Any]] = []
    for elem in root.iter():
        tag = local_name(elem.tag)
        if tag in BAD_TAGS:
            metrics["unsupported_nodes"] += 1
            emit(findings, "error", "unsupported-svg-node", str(svg_file), f"Unsupported SVG node: <{tag}>.")

    for elem, offset_x, offset_y in iter_text_elements_with_offsets(root):
        metrics["texts"] += 1
        box = estimate_text_box(elem, offset_x, offset_y)
        if box is None:
            continue
        lines = text_lines(elem)
        font_size = number(elem.get("font-size"), 18)
        text_blocks.append(
            {
                "box": box,
                "font_size": font_size,
                "lines": lines,
                "text": " ".join(lines),
            }
        )
        left, top, right, bottom = box
        if left < -2 or right > canvas.width + 2 or top < -2 or bottom > canvas.height + 2:
            emit(
                findings,
                "error",
                "text-outside-canvas",
                str(svg_file),
                (
                    "Text bbox appears outside canvas: "
                    f"left={left:.1f}, top={top:.1f}, right={right:.1f}, bottom={bottom:.1f}."
                ),
            )
        elif left < safe_left or right > safe_right or top < safe_top or bottom > safe_bottom:
            cy = (top + bottom) / 2.0
            in_whitelist = False
            for region in safe_edge_whitelist or set():
                ranges = canvas.safe_edge_ranges().get(region)
                if ranges and ranges[0] <= cy <= ranges[1]:
                    in_whitelist = True
                    break
            emit(
                findings,
                "warning",
                "text-near-safe-edge-whitelist" if in_whitelist else "text-near-safe-edge",
                str(svg_file),
                (
                    "Text is near or outside safe area: "
                    f"left={left:.1f}, top={top:.1f}, right={right:.1f}, bottom={bottom:.1f}."
                ),
            )

    for left_index, left_item in enumerate(text_blocks):
        left_box = left_item["box"]
        left_text = str(left_item.get("text") or "")
        if not _contains_cjk(left_text):
            continue
        for right_item in text_blocks[left_index + 1 :]:
            right_box = right_item["box"]
            right_text = str(right_item.get("text") or "")
            if not _contains_cjk(right_text):
                continue
            overlap_w = min(left_box[2], right_box[2]) - max(left_box[0], right_box[0])
            overlap_h = min(left_box[3], right_box[3]) - max(left_box[1], right_box[1])
            if overlap_w <= 12 or overlap_h <= 4:
                continue
            emit(
                findings,
                "error",
                "text-overlap-risk",
                str(svg_file),
                (
                    "Detected overlapping SVG text boxes. "
                    f"overlap_w={overlap_w:.1f}, overlap_h={overlap_h:.1f}; "
                    f"left={left_box[0]:.1f},{left_box[1]:.1f},{left_box[2]:.1f},{left_box[3]:.1f}; "
                    f"right={right_box[0]:.1f},{right_box[1]:.1f},{right_box[2]:.1f},{right_box[3]:.1f}."
                ),
            )

    for item in text_blocks:
        box = item["box"]
        lines = item["lines"]
        if len(lines) != 1:
            continue
        line = lines[0]
        if not _contains_cjk(line):
            continue
        cx = (box[0] + box[2]) / 2.0
        cy = (box[1] + box[3]) / 2.0
        for block in slide_plan_blocks or []:
            b = block.get("box")
            if not (isinstance(b, list) and len(b) == 4):
                continue
            if not _contains_point(b, cx, cy):
                continue
            block_w = float(b[2])
            if block_w > 360:
                break
            font_size = float(item["font_size"])
            max_units = max(1.0, (max(40.0, block_w - 16.0)) / max(font_size * 0.92, 1.0))
            line_units = _estimated_text_units(line)
            risk_level = _classify_overflow_risk(line_units / max_units, profile)
            if risk_level:
                metrics[f"overflow_risk_{risk_level}_count"] += 1
                emit(
                    findings,
                    "warning",
                    f"overflow-risk-{risk_level}",
                    str(svg_file),
                    (
                        "Detected narrow-column overflow risk for CJK single-line text "
                        f"(box_w={block_w:.1f}, line_units={line_units:.1f}, max_units={max_units:.1f})."
                    ),
                )
            break

    return metrics


def validate_svg_set(
    project_dir: Path,
    canvas: CanvasContext,
    svg_dir: Path,
    slides: list[dict[str, Any]],
    findings: list[Finding],
    slide_plan: dict[int, list[dict[str, Any]]] | None = None,
    slide_id: int | None = None,
    enforce_blueprint_sync: bool = False,
    safe_area_profile: str = "legacy",
    safe_edge_whitelist: set[str] | None = None,
    profile: str = "presentation",
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "svg_dir": str(svg_dir),
        "svg_files": 0,
        "svg_texts": 0,
        "unsupported_nodes": 0,
        "overflow_risk_high_count": 0,
        "overflow_risk_medium_count": 0,
        "overflow_risk_low_count": 0,
    }
    if not svg_dir.exists():
        emit(findings, "error", "missing-svg-dir", str(svg_dir), "SVG output directory does not exist.")
        return metrics

    svg_files = sorted(svg_dir.glob("slide_*.svg"))
    expected = {
        f"slide_{int(slide.get('id', idx)):02d}.svg"
        for idx, slide in enumerate(slides, start=1)
        if isinstance(slide, dict)
    }
    if slide_id is not None:
        target_name = f"slide_{slide_id:02d}.svg"
        expected = {target_name}
        svg_files = [path for path in svg_files if path.name == target_name]
    metrics["svg_files"] = len(svg_files)

    actual = {path.name for path in svg_files}
    if slide_id is None and len(actual) != len(slides):
        emit(
            findings,
            "error" if enforce_blueprint_sync else "warning",
            "svg-count-mismatch",
            str(svg_dir),
            (
                f"Expected {len(slides)} SVG files, found {len(actual)}. "
                "Enable --enforce-blueprint-sync to block this in CI/delivery."
            ),
        )

    for name in sorted(expected - actual):
        emit(
            findings,
            "error" if (enforce_blueprint_sync or slide_id is not None) else "warning",
            "missing-svg-file",
            str(svg_dir / name),
            "Expected SVG file is missing.",
        )
    if slide_id is None:
        for name in sorted(actual - expected):
            emit(
                findings,
                "error" if enforce_blueprint_sync else "warning",
                "extra-svg-file",
                str(svg_dir / name),
                "Extra SVG file does not map to a blueprint slide id.",
            )

    parsed_slide_plan = slide_plan if slide_plan is not None else parse_slide_plan(project_dir, findings)
    for svg_file in svg_files:
        sid_match = re.match(r"slide_(\d+)\.svg$", svg_file.name)
        sid = int(sid_match.group(1)) if sid_match else -1
        file_metrics = validate_svg_file(
            svg_file,
            findings,
            canvas,
            safe_area_profile=safe_area_profile,
            slide_plan_blocks=parsed_slide_plan.get(sid, []),
            safe_edge_whitelist=safe_edge_whitelist,
            profile=profile,
        )
        metrics["svg_texts"] += file_metrics["texts"]
        metrics["unsupported_nodes"] += file_metrics["unsupported_nodes"]
        metrics["overflow_risk_high_count"] += int(file_metrics.get("overflow_risk_high_count", 0))
        metrics["overflow_risk_medium_count"] += int(file_metrics.get("overflow_risk_medium_count", 0))
        metrics["overflow_risk_low_count"] += int(file_metrics.get("overflow_risk_low_count", 0))
        quality = check_svg_file(svg_file, canvas_key=canvas.key)
        for item in quality.errors:
            message = str(item)
            code = (
                "svg-contrast-illegible"
                if "contrast too low" in message.lower()
                and "ink likely invisible or illegible" in message.lower()
                else "svg-quality"
            )
            emit(findings, "error", code, str(svg_file), message)
        for item in quality.warnings:
            emit(findings, "warning", "svg-quality", str(svg_file), item)
    return metrics


def _slide_id_from_svg_name(svg_file: Path) -> int | None:
    match = re.match(r"slide_(\d+)\.svg$", svg_file.name)
    if not match:
        return None
    return int(match.group(1))


def _collect_text_nodes(root: ET.Element) -> list[dict[str, Any]]:
    parent_by_child = {child: parent for parent in root.iter() for child in parent}
    group_ids = {
        elem: f"group-{index}"
        for index, elem in enumerate(
            (candidate for candidate in root.iter() if local_name(candidate.tag) == "g"),
            start=1,
        )
    }
    nodes: list[dict[str, Any]] = []
    for elem, offset_x, offset_y in iter_text_elements_with_offsets(root):
        lines = text_lines(elem)
        if not lines:
            continue
        text_value = " ".join(lines).strip()
        font_size = number(elem.get("font-size"), 18)
        line_height = font_size * 0.84
        max_line_width = 0.0
        area = 0.0
        for line in lines:
            width = max(1.0, visual_width(line) * font_size * 0.78)
            max_line_width = max(max_line_width, width)
            area += width * line_height
        x = number(elem.get("x")) + offset_x
        anchor = (elem.get("text-anchor") or "start").strip().lower()
        estimated_box = estimate_text_box(elem, offset_x, offset_y)
        if estimated_box is None:
            if anchor == "middle":
                left = x - max_line_width / 2.0
            elif anchor == "end":
                left = x - max_line_width
            else:
                left = x
            top = number(elem.get("y")) + offset_y - font_size * 0.72
            right = left + max_line_width
            bottom = top + (line_height * len(lines))
        else:
            left, top, right, bottom = estimated_box
            max_line_width = max(max_line_width, right - left)
        center_x = (left + right) / 2.0

        group_elem: ET.Element | None = None
        parent = parent_by_child.get(elem)
        while parent is not None:
            if local_name(parent.tag) == "g":
                group_elem = parent
                break
            parent = parent_by_child.get(parent)

        nodes.append(
            {
                "x": x,
                "y": number(elem.get("y")) + offset_y,
                "text_anchor": anchor,
                "font_size": font_size,
                "font_weight": number(elem.get("font-weight"), 400.0),
                "fill": (elem.get("fill") or "").strip(),
                "text": text_value,
                "chars": len(text_value),
                "line_count": len(lines),
                "max_line_width": max_line_width,
                "estimated_area": area,
                "left": left,
                "right": right,
                "top": top,
                "bottom": bottom,
                "center_x": center_x,
                "group_id": group_ids.get(group_elem) if group_elem is not None else None,
                "group_element": group_elem,
            }
        )
    return nodes


def _semantic_text_count_for_visual_quality(project_dir: Path, svg_file: Path, fallback_count: int) -> int:
    if svg_file.parent.name != "svg_final":
        return fallback_count
    source = project_dir / "svg_output" / svg_file.name
    if not source.exists():
        return fallback_count
    try:
        root = ET.parse(source).getroot()
    except (ET.ParseError, OSError):
        return fallback_count
    semantic_nodes = _collect_text_nodes(root)
    return len(semantic_nodes) or fallback_count


def _estimate_hierarchy_depth(texts: list[dict[str, Any]]) -> int:
    buckets: set[int] = set()
    for node in texts:
        try:
            size = float(node["font_size"])
        except (KeyError, TypeError, ValueError):
            continue
        bucket = int(round(size / 2.0) * 2)
        buckets.add(bucket)
    return len(buckets)


def _median_number(values: Iterable[float], default: float = 0.0) -> float:
    ordered = sorted(float(value) for value in values)
    if not ordered:
        return default
    middle = len(ordered) // 2
    if len(ordered) % 2:
        return ordered[middle]
    return (ordered[middle - 1] + ordered[middle]) / 2.0


def _normalize_visual_text(value: str) -> str:
    return "".join(char.lower() for char in value if char.isalnum())


def _load_visual_quality_slide_contexts(project_dir: Path) -> dict[int, dict[str, Any]]:
    """Read existing slide facts without introducing another persisted state source."""
    loaded = load_json_object(project_dir / "blueprint.json", encoding="utf-8-sig")
    if loaded.status != "ok" or not isinstance(loaded.payload, dict):
        return {}
    slides = loaded.payload.get("slides")
    if not isinstance(slides, list):
        return {}
    visual_plans = _load_slide_visual_plan_contract(project_dir)
    contexts: dict[int, dict[str, Any]] = {}
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        sid = slide.get("id")
        if not isinstance(sid, int):
            continue
        visual_plan = visual_plans.get(sid)
        contexts[sid] = {
            "slide": slide,
            "visual_plan": visual_plan,
            "usable": isinstance(visual_plan, dict),
        }
    return contexts


def _is_confirmed_cover_context(context: dict[str, Any] | None) -> bool:
    if not context or not bool(context.get("usable")):
        return False
    slide = context.get("slide")
    if not isinstance(slide, dict):
        return False
    page_type = str(slide.get("page_type") or "").strip().lower()
    layout_tag = str(slide.get("layout_tag") or "").strip().lower()
    return page_type == "cover" or layout_tag.startswith("cover-")


def _title_nodes_for_context(
    texts: list[dict[str, Any]],
    context: dict[str, Any] | None,
) -> tuple[list[dict[str, Any]], bool]:
    slide = context.get("slide") if isinstance(context, dict) else None
    expected_title = str(slide.get("title") or "") if isinstance(slide, dict) else ""
    normalized_title = _normalize_visual_text(expected_title)
    matched = [
        node
        for node in texts
        if normalized_title
        and len(_normalize_visual_text(str(node.get("text") or ""))) >= 4
        and _normalize_visual_text(str(node.get("text") or "")) in normalized_title
    ]
    if matched:
        joined = _normalize_visual_text(
            " ".join(
                str(node.get("text") or "")
                for node in sorted(matched, key=lambda item: (float(item.get("y", 0.0)), float(item.get("x", 0.0))))
            )
        )
        return matched, bool(joined and (joined in normalized_title or normalized_title in joined))
    max_font = max((float(node.get("font_size", 0.0)) for node in texts), default=0.0)
    return [node for node in texts if float(node.get("font_size", 0.0)) >= max_font * 0.9], False


def _evaluate_headline_and_dominance(
    texts: list[dict[str, Any]],
    context: dict[str, Any] | None,
    canvas: CanvasContext | None,
) -> tuple[bool, bool, int, dict[str, Any]]:
    legacy_headline = any(
        float(node.get("y", 0.0)) <= 170.0 and float(node.get("font_size", 0.0)) >= 26.0
        for node in texts
    )
    max_font = max((float(node.get("font_size", 0.0)) for node in texts), default=0.0)
    legacy_dominant = [
        node
        for node in texts
        if float(node.get("font_size", 0.0)) >= max_font * 0.9
        and float(node.get("y", 0.0)) <= 220.0
        and int(node.get("chars", 0)) <= 80
    ]
    if not _is_confirmed_cover_context(context):
        return legacy_headline, bool(legacy_dominant), len(legacy_dominant), {"mode": "legacy"}

    title_nodes, title_match = _title_nodes_for_context(texts, context)
    title_ids = {id(node) for node in title_nodes}
    body_nodes = [
        node
        for node in texts
        if id(node) not in title_ids
        and float(node.get("font_size", 0.0)) > 12.0
        and not str(node.get("text") or "").strip().isdigit()
    ]
    body_median = _median_number(
        (float(node.get("font_size", 0.0)) for node in body_nodes),
        default=18.0,
    )
    title_size = max((float(node.get("font_size", 0.0)) for node in title_nodes), default=0.0)
    title_y_values = [float(node.get("y", 0.0)) for node in title_nodes]
    unique_cluster = bool(title_nodes) and len(title_nodes) <= 3
    if len(title_y_values) > 1:
        unique_cluster = unique_cluster and (max(title_y_values) - min(title_y_values)) <= title_size * 1.6
    canvas_width = float(canvas.width) if canvas is not None else 1280.0
    canvas_height = float(canvas.height) if canvas is not None else 720.0
    in_safe_canvas = bool(title_nodes) and all(
        float(node.get("left", 0.0)) >= 0.0
        and float(node.get("right", 0.0)) <= canvas_width
        and float(node.get("top", 0.0)) >= 0.0
        and float(node.get("bottom", 0.0)) <= canvas_height
        for node in title_nodes
    )
    prominent = (
        title_size >= 26.0
        and title_size >= body_median * 1.45
        and (title_size - body_median) >= 8.0
    )
    emphasized = any(float(node.get("font_weight", 400.0)) >= 600.0 for node in title_nodes)
    meaningful = sum(int(node.get("chars", 0)) for node in title_nodes) >= 4
    strong_readable_anchor = bool(
        prominent and unique_cluster and in_safe_canvas and meaningful and (title_match or emphasized)
    )
    return (
        strong_readable_anchor,
        strong_readable_anchor,
        1 if strong_readable_anchor else 0,
        {
            "mode": "cover",
            "title_size": round(title_size, 2),
            "body_median": round(body_median, 2),
            "title_match": title_match,
            "unique_cluster": unique_cluster,
        },
    )


def _evaluate_hierarchy(
    texts: list[dict[str, Any]],
    context: dict[str, Any] | None,
    policy: Any,
    canvas: CanvasContext | None,
) -> tuple[bool, dict[str, Any]]:
    max_font = max((float(node.get("font_size", 0.0)) for node in texts), default=0.0)
    if not context or not bool(context.get("usable")):
        return max_font >= float(policy.min_heading_font_px), {"mode": "legacy", "max_font": max_font}

    title_nodes, title_match = _title_nodes_for_context(texts, context)
    title_ids = {id(node) for node in title_nodes}
    body_nodes = [
        node
        for node in texts
        if id(node) not in title_ids
        and float(node.get("font_size", 0.0)) > 12.0
        and int(node.get("chars", 0)) > 0
    ]
    if not title_nodes or not body_nodes:
        return max_font >= float(policy.min_heading_font_px), {"mode": "legacy", "max_font": max_font}

    title_size = max(float(node.get("font_size", 0.0)) for node in title_nodes)
    body_median = _median_number(float(node.get("font_size", 0.0)) for node in body_nodes)
    relative_clear = title_size >= body_median * 1.45 and (title_size - body_median) >= 8.0
    depth = _estimate_hierarchy_depth(texts)
    depth_clear = depth >= 3
    title_weight = max(float(node.get("font_weight", 400.0)) for node in title_nodes)
    body_weight = _median_number(
        (float(node.get("font_weight", 400.0)) for node in body_nodes),
        default=400.0,
    )
    title_fills = {str(node.get("fill") or "").strip().lower() for node in title_nodes}
    body_fills = {str(node.get("fill") or "").strip().lower() for node in body_nodes}
    emphasis_clear = (title_weight - body_weight) >= 100.0 or bool(title_fills - body_fills)
    canvas_height = float(canvas.height) if canvas is not None else 720.0
    cover_context = _is_confirmed_cover_context(context)
    position_clear = cover_context or min(float(node.get("y", 0.0)) for node in title_nodes) <= canvas_height * 0.28
    title_y_values = [float(node.get("y", 0.0)) for node in title_nodes]
    unique_clear = len(title_nodes) <= 3 and (
        len(title_y_values) <= 1 or (max(title_y_values) - min(title_y_values)) <= title_size * 1.6
    )
    supporting_signals = sum(
        (
            depth_clear,
            bool(position_clear and title_match),
            emphasis_clear,
            unique_clear,
        )
    )
    hierarchy_clear = bool(relative_clear and supporting_signals >= 2)
    return hierarchy_clear, {
        "mode": "relative",
        "title_size": round(title_size, 2),
        "body_median": round(body_median, 2),
        "depth": depth,
        "supporting_signals": supporting_signals,
    }


def _semantic_visual_quality_source(
    project_dir: Path,
    svg_file: Path,
    rendered_root: ET.Element,
) -> tuple[ET.Element, list[dict[str, Any]], str]:
    if svg_file.parent.name == "svg_final":
        source = project_dir / "svg_output" / svg_file.name
        if source.exists():
            try:
                source_root = ET.parse(source).getroot()
            except (ET.ParseError, OSError):
                pass
            else:
                source_nodes = _collect_text_nodes(source_root)
                if source_nodes:
                    return source_root, source_nodes, str(source)
    return rendered_root, _collect_text_nodes(rendered_root), str(svg_file)


def _iter_elements_with_offsets(
    elem: ET.Element,
    offset_x: float = 0.0,
    offset_y: float = 0.0,
):
    dx, dy = parse_translate(elem.get("transform"))
    next_x = offset_x + dx
    next_y = offset_y + dy
    yield elem, next_x, next_y
    for child in elem:
        yield from _iter_elements_with_offsets(child, next_x, next_y)


def _visible_container_bboxes(
    root: ET.Element,
    *,
    canvas_width: float,
    canvas_height: float,
) -> list[tuple[float, float, float, float]]:
    canvas_area = max(1.0, canvas_width * canvas_height)
    boxes: list[tuple[float, float, float, float]] = []
    for elem, offset_x, offset_y in _iter_elements_with_offsets(root):
        name = local_name(elem.tag)
        fill = _path_style_value(elem, "fill").strip().lower()
        if fill in {"", "none", "transparent"} or number(elem.get("opacity"), 1.0) <= 0.0:
            continue
        bbox: tuple[float, float, float, float] | None = None
        if name == "rect":
            x = number(elem.get("x")) + offset_x
            y = number(elem.get("y")) + offset_y
            width = number(elem.get("width"))
            height = number(elem.get("height"))
            if width >= 120.0 and height >= 50.0:
                bbox = (x, y, x + width, y + height)
        elif name == "path":
            raw_bbox = _path_card_like_bbox(elem)
            if raw_bbox is not None:
                bbox = (
                    raw_bbox[0] + offset_x,
                    raw_bbox[1] + offset_y,
                    raw_bbox[2] + offset_x,
                    raw_bbox[3] + offset_y,
                )
        if bbox is None:
            continue
        width = bbox[2] - bbox[0]
        height = bbox[3] - bbox[1]
        if width <= 0.0 or height <= 0.0 or (width * height) / canvas_area > 0.4:
            continue
        boxes.append(bbox)
    return boxes


def _union_bbox(boxes: Iterable[tuple[float, float, float, float]]) -> tuple[float, float, float, float] | None:
    collected = list(boxes)
    if not collected:
        return None
    return (
        min(box[0] for box in collected),
        min(box[1] for box in collected),
        max(box[2] for box in collected),
        max(box[3] for box in collected),
    )


def _clustered_fraction(values: Iterable[float], tolerance: float) -> float:
    ordered = sorted(float(value) for value in values)
    if not ordered:
        return 0.0
    best = 1
    left = 0
    for right, value in enumerate(ordered):
        while value - ordered[left] > tolerance:
            left += 1
        best = max(best, right - left + 1)
    return best / len(ordered)


def _visual_group_diagnostics(
    nodes: list[dict[str, Any]],
    container_boxes: list[tuple[float, float, float, float]],
) -> dict[str, float | int | bool]:
    node_bbox = _union_bbox(
        (
            float(node.get("left", 0.0)),
            float(node.get("top", 0.0)),
            float(node.get("right", 0.0)),
            float(node.get("bottom", 0.0)),
        )
        for node in nodes
    )
    group_bbox = _union_bbox([*(container_boxes or []), *([node_bbox] if node_bbox else [])])
    if node_bbox is None or group_bbox is None:
        return {"valid": False, "alignment_quality": 0.0, "scatter_ratio": 1.0, "coverage": 0.0}
    node_area = max(1.0, (node_bbox[2] - node_bbox[0]) * (node_bbox[3] - node_bbox[1]))
    group_area = max(1.0, (group_bbox[2] - group_bbox[0]) * (group_bbox[3] - group_bbox[1]))
    alignment_quality = max(
        _clustered_fraction((float(node.get("left", 0.0)) for node in nodes), 18.0),
        _clustered_fraction((float(node.get("y", 0.0)) for node in nodes), 14.0),
    )
    coverage = min(1.0, node_area / group_area)
    node_count = len(nodes)
    valid = bool(
        container_boxes
        and 2 <= node_count <= 10
        and coverage >= 0.05
        and alignment_quality >= 0.48
    )
    return {
        "valid": valid,
        "alignment_quality": round(alignment_quality, 4),
        "scatter_ratio": round(1.0 - alignment_quality, 4),
        "coverage": round(coverage, 4),
        "node_count": node_count,
    }


def _evaluate_fragmentation(
    root: ET.Element,
    texts: list[dict[str, Any]],
    node_threshold: int,
    canvas: CanvasContext | None,
    semantic_source: str,
) -> tuple[bool, dict[str, Any]]:
    canvas_width = float(canvas.width) if canvas is not None else 1280.0
    canvas_height = float(canvas.height) if canvas is not None else 720.0
    valid_groups: list[set[int]] = []
    group_diagnostics: list[dict[str, float | int | bool]] = []

    nodes_by_group: dict[str, list[tuple[int, dict[str, Any]]]] = {}
    group_elements: dict[str, ET.Element] = {}
    for index, node in enumerate(texts):
        group_id = node.get("group_id")
        group_element = node.get("group_element")
        if isinstance(group_id, str) and isinstance(group_element, ET.Element):
            nodes_by_group.setdefault(group_id, []).append((index, node))
            group_elements[group_id] = group_element

    assigned: set[int] = set()
    for group_id, indexed_nodes in nodes_by_group.items():
        group_nodes = [node for _index, node in indexed_nodes]
        boxes = _visible_container_bboxes(
            group_elements[group_id],
            canvas_width=canvas_width,
            canvas_height=canvas_height,
        )
        group_evidence = _visual_group_diagnostics(group_nodes, boxes)
        if bool(group_evidence.get("valid")):
            indices = {index for index, _node in indexed_nodes}
            valid_groups.append(indices)
            group_diagnostics.append(group_evidence)
            assigned.update(indices)

    containers = _visible_container_bboxes(root, canvas_width=canvas_width, canvas_height=canvas_height)
    nodes_by_container: dict[int, list[tuple[int, dict[str, Any]]]] = {}
    for index, node in enumerate(texts):
        if index in assigned:
            continue
        center_x = float(node.get("center_x", 0.0))
        center_y = (float(node.get("top", 0.0)) + float(node.get("bottom", 0.0))) / 2.0
        candidates = [
            (container_index, box)
            for container_index, box in enumerate(containers)
            if box[0] <= center_x <= box[2] and box[1] <= center_y <= box[3]
        ]
        if not candidates:
            continue
        container_index, _box = min(
            candidates,
            key=lambda item: (item[1][2] - item[1][0]) * (item[1][3] - item[1][1]),
        )
        nodes_by_container.setdefault(container_index, []).append((index, node))

    for container_index, indexed_nodes in nodes_by_container.items():
        group_nodes = [node for _index, node in indexed_nodes]
        group_evidence = _visual_group_diagnostics(group_nodes, [containers[container_index]])
        if bool(group_evidence.get("valid")):
            indices = {index for index, _node in indexed_nodes}
            valid_groups.append(indices)
            group_diagnostics.append(group_evidence)
            assigned.update(indices)

    grouped_node_count = sum(len(indices) for indices in valid_groups)
    effective_count = len(texts) - grouped_node_count + len(valid_groups)
    ungrouped_count = len(texts) - grouped_node_count
    alignment_quality = _median_number(
        (float(item.get("alignment_quality", 0.0)) for item in group_diagnostics),
        default=0.0,
    )
    scatter_ratio = _median_number(
        (float(item.get("scatter_ratio", 1.0)) for item in group_diagnostics),
        default=1.0,
    )
    container_coverage = grouped_node_count / max(1, len(texts))
    max_group_nodes = max((len(indices) for indices in valid_groups), default=0)
    diagnostics: dict[str, Any] = {
        "semantic_source": semantic_source,
        "raw_text_count": len(texts),
        "effective_fragment_count": effective_count,
        "valid_group_count": len(valid_groups),
        "ungrouped_node_count": ungrouped_count,
        "ungrouped_ratio": round(ungrouped_count / max(1, len(texts)), 4),
        "group_scatter_ratio": round(scatter_ratio, 4),
        "max_group_nodes": max_group_nodes,
        "container_coverage": round(container_coverage, 4),
        "alignment_quality": round(alignment_quality, 4),
    }
    return effective_count <= node_threshold, diagnostics


def _visual_severity_from_gap(gap: float, advisory_threshold: float, warning_threshold: float) -> str | None:
    if gap >= warning_threshold:
        return "warning"
    if gap >= advisory_threshold:
        return "advisory"
    return None


def validate_visual_quality(
    project_dir: Path,
    svg_dir: Path,
    findings: list[Finding],
    slide_id: int | None = None,
    *,
    profile: str = "presentation",
    canvas: CanvasContext | None = None,
) -> VisualMetrics:
    policy = resolve_profile_policy(profile)
    canvas_area_ratio = 1.0
    if canvas is not None and canvas.width > 0 and canvas.height > 0:
        canvas_area_ratio = (canvas.width * canvas.height) / (1280.0 * 720.0)
    char_threshold = max(320, int(round(policy.visual_density_base_chars * canvas_area_ratio)))
    node_threshold = max(12, int(round(policy.visual_density_base_nodes * canvas_area_ratio)))

    metrics: VisualMetrics = {
        "visual_score": 100.0,
        "visual_findings": [],
        "repair_recommendation": [],
        "density_flag": False,
        "hierarchy_flag": False,
        "visual_whitespace_ratio": 0.0,
        "visual_hierarchy_depth_score": 0.0,
        "visual_dominant_point_count": 0,
        "visual_repetition_penalty": 0.0,
        "visual_alignment_quality_score": 0.0,
        "visual_fragmentation_by_slide": {},
    }
    svg_files = sorted(svg_dir.glob("slide_*.svg"))
    if slide_id is not None:
        target = f"slide_{slide_id:02d}.svg"
        svg_files = [path for path in svg_files if path.name == target]
    if not svg_files:
        return metrics

    per_slide_scores: list[float] = []
    per_slide_whitespace: list[float] = []
    per_slide_hierarchy_depth_score: list[float] = []
    per_slide_alignment_score: list[float] = []
    dominant_points_total = 0
    slide_signatures: list[str] = []
    visual_findings: list[dict[str, Any]] = []
    repair_recommendation: list[dict[str, Any]] = []
    density_flag = False
    hierarchy_flag = False
    fragmentation_by_slide: dict[str, dict[str, Any]] = {}
    slide_contexts = _load_visual_quality_slide_contexts(project_dir)

    for svg_file in svg_files:
        slide_score = 100.0
        local_repair: list[str] = []
        local_findings: list[tuple[str, str, str]] = []
        sid = _slide_id_from_svg_name(svg_file)
        try:
            tree = ET.parse(svg_file)
            root = tree.getroot()
        except (ET.ParseError, OSError):
            continue
        texts = _collect_text_nodes(root)
        if not texts:
            slide_score -= 25
            local_findings.append(("warning", "visual-no-text", "Slide has no readable text nodes."))
            local_repair.append("Add a clear headline and at least one supporting text block.")
        else:
            context = slide_contexts.get(sid) if isinstance(sid, int) else None
            char_count = sum(int(node["chars"]) for node in texts)
            text_count = len(texts)
            semantic_root, semantic_texts, semantic_source = _semantic_visual_quality_source(
                project_dir,
                svg_file,
                root,
            )
            semantic_text_count = len(semantic_texts)
            fragmentation_clear, fragmentation_diagnostics = _evaluate_fragmentation(
                semantic_root,
                semantic_texts,
                node_threshold,
                canvas,
                semantic_source,
            )
            fragmentation_by_slide[str(sid if sid is not None else svg_file.name)] = fragmentation_diagnostics
            headline_clear, dominance_clear, dominant_point_count, _headline_diagnostics = (
                _evaluate_headline_and_dominance(texts, context, canvas)
            )
            max_font = max(float(node["font_size"]) for node in texts)
            x_tracks = {round(float(node["x"]) / 8) * 8 for node in texts}
            text_area = sum(float(node.get("estimated_area", 0.0)) for node in texts)
            whitespace_ratio = max(0.0, min(1.0, 1.0 - (text_area / max(1.0, 1280.0 * 720.0))))
            hierarchy_depth = _estimate_hierarchy_depth(texts)
            hierarchy_depth_score = max(0.0, min(1.0, hierarchy_depth / 4.0))
            hierarchy_clear, hierarchy_diagnostics = _evaluate_hierarchy(texts, context, policy, canvas)
            alignment_quality_score = max(0.0, min(1.0, 1.0 - max(0.0, len(x_tracks) - 6) / 12.0))

            conclusion_nodes = [
                node for node in texts if CONCLUSION_ANCHOR_RE.search(str(node.get("text") or ""))
            ]
            if conclusion_nodes and canvas is not None:
                bad_nodes = [
                    node
                    for node in conclusion_nodes
                    if float(node.get("y", 0.0)) < (canvas.footer_start - 12.0)
                ]
                if bad_nodes:
                    slide_score -= 4
                    local_findings.append(
                        (
                            "warning",
                            "visual-anchor-conclusion-misaligned",
                            "Conclusion sentence is not placed in the footer anchor band.",
                        )
                    )
                    local_repair.append("Move final value/conclusion sentence to the bottom footer anchor band.")

            engine_nodes = [
                node
                for node in texts
                if ENGINE_TITLE_RE.search(str(node.get("text") or ""))
                and float(node.get("font_size", 0.0)) >= 24.0
                and int(node.get("chars", 0)) <= 40
            ]
            if engine_nodes and canvas is not None:
                focus = max(engine_nodes, key=lambda item: float(item.get("font_size", 0.0)))
                cx = float(focus.get("center_x", focus.get("x", 0.0)))
                cy = float(focus.get("y", 0.0))
                center_left = canvas.safe_x + canvas.safe_w / 3.0 - 40.0
                center_right = canvas.safe_x + 2.0 * canvas.safe_w / 3.0 + 40.0
                top_band_min = canvas.safe_y + 90.0
                top_band_max = canvas.safe_y + 230.0
                if not (center_left <= cx <= center_right and top_band_min <= cy <= top_band_max):
                    slide_score -= 4
                    local_findings.append(
                        (
                            "warning",
                            "visual-anchor-engine-title-misaligned",
                            "Engine hub title is away from the center-top anchor zone.",
                        )
                    )
                    local_repair.append("Re-anchor the engine title to the center column header zone.")

            per_slide_whitespace.append(round(whitespace_ratio, 4))
            per_slide_hierarchy_depth_score.append(round(hierarchy_depth_score, 4))
            per_slide_alignment_score.append(round(alignment_quality_score, 4))
            dominant_points_total += dominant_point_count
            signature = (
                f"tracks:{','.join(str(item) for item in sorted(x_tracks)[:6])}"
                f"|headline:{1 if headline_clear else 0}"
            )
            slide_signatures.append(signature)

            if char_count > char_threshold:
                gap = char_count - char_threshold
                sev = _visual_severity_from_gap(gap, 120, 260) or "advisory"
                slide_score -= 5 if sev == "advisory" else 8
                density_flag = True
                local_findings.append(
                    (
                        sev,
                        "visual-density-high",
                        f"Text density is high ({char_count} characters, budget={char_threshold}).",
                    )
                )
                local_repair.append("Compress secondary copy and keep each primary bullet within two lines.")
            if not fragmentation_clear:
                effective_count = int(fragmentation_diagnostics["effective_fragment_count"])
                gap = effective_count - node_threshold
                sev = _visual_severity_from_gap(gap, 4, 10) or "advisory"
                slide_score -= 4 if sev == "advisory" else 7
                density_flag = True
                local_findings.append(
                    (
                        sev,
                        "visual-text-fragmented",
                        (
                            "Too many effective semantic text fragments "
                            f"(raw={semantic_text_count}, effective={effective_count}, budget={node_threshold}, "
                            f"groups={fragmentation_diagnostics['valid_group_count']}, "
                            f"scatter={fragmentation_diagnostics['group_scatter_ratio']:.2f}, "
                            f"max_group_nodes={fragmentation_diagnostics['max_group_nodes']}, "
                            f"container_coverage={fragmentation_diagnostics['container_coverage']:.2f}, "
                            f"alignment={fragmentation_diagnostics['alignment_quality']:.2f})."
                        ),
                    )
                )
                local_repair.append("Merge scattered labels into grouped blocks or cards.")
            if not headline_clear:
                slide_score -= 8
                hierarchy_flag = True
                local_findings.append(
                    ("advisory", "visual-headline-weak", "No strong headline detected in title zone.")
                )
                local_repair.append("Promote one headline in the top safe area with larger font weight/size.")
            if not hierarchy_clear:
                heading_gap = max(0.0, float(policy.min_heading_font_px - max_font))
                sev = _visual_severity_from_gap(heading_gap, 6, 12) or "advisory"
                slide_score -= 5 if sev == "advisory" else 8
                hierarchy_flag = True
                local_findings.append(
                    (
                        sev,
                        "visual-hierarchy-flat",
                        (
                            "Rendered title/body hierarchy is weak "
                            f"(max_font={max_font:.1f}px, required={policy.min_heading_font_px:.1f}px, "
                            f"mode={hierarchy_diagnostics.get('mode')}, "
                            f"title={hierarchy_diagnostics.get('title_size', max_font)}, "
                            f"body_median={hierarchy_diagnostics.get('body_median', 0.0)})."
                        ),
                    )
                )
                local_repair.append("Increase contrast between headline and body typography.")
            if len(x_tracks) > 12 and text_count > 12:
                gap = len(x_tracks) - 12
                sev = _visual_severity_from_gap(gap, 3, 6) or "advisory"
                slide_score -= 4 if sev == "advisory" else 6
                local_findings.append(
                    (sev, "visual-alignment-chaos", f"Too many x-alignment tracks ({len(x_tracks)}).")
                )
                local_repair.append("Snap blocks to fewer alignment tracks to strengthen rhythm.")
            if whitespace_ratio < 0.75:
                slide_score -= 3
                local_findings.append(
                    (
                        "advisory",
                        "visual-whitespace-low",
                        f"Whitespace ratio is low ({whitespace_ratio:.2f}, expected>=0.75).",
                    )
                )
                local_repair.append("Increase spacing and remove secondary text blocks to restore breathing room.")
            if not dominance_clear:
                slide_score -= 3
                local_findings.append(
                    (
                        "advisory",
                        "visual-dominance-weak",
                        "No dominant visual point detected in top decision zone.",
                    )
                )
                local_repair.append("Promote one dominant headline/callout in the upper visual zone.")
            if alignment_quality_score < 0.45:
                slide_score -= 2
                local_findings.append(
                    (
                        "advisory",
                        "visual-alignment-weak",
                        f"Alignment quality is weak ({alignment_quality_score:.2f}).",
                    )
                )
                local_repair.append("Reduce alignment tracks and snap text blocks to a tighter column system.")

        slide_score = max(0.0, min(100.0, slide_score))
        per_slide_scores.append(slide_score)

        for severity, code, message in local_findings:
            visual_findings.append(
                {
                    "slide": sid,
                    "severity": severity,
                    "code": code,
                    "message": message,
                }
            )
            emit(findings, severity, code, str(svg_file), message)
        if local_repair:
            repair_recommendation.append(
                {
                    "slide": sid,
                    "actions": local_repair,
                }
            )

    visual_score = round(sum(per_slide_scores) / len(per_slide_scores), 2) if per_slide_scores else 100.0
    repetition_penalty = 0.0
    if len(slide_signatures) >= 3:
        counts: dict[str, int] = {}
        for sig in slide_signatures:
            counts[sig] = counts.get(sig, 0) + 1
        dominant_signature = max(counts.values())
        repetition_penalty = max(0.0, min(1.0, (dominant_signature / len(slide_signatures)) - 0.5))
        if repetition_penalty >= 0.2:
            emit(
                findings,
                "warning",
                "visual-repetition-high",
                str(svg_dir),
                f"Slide composition repetition is high (penalty={repetition_penalty:.2f}).",
            )

    metrics["visual_score"] = visual_score
    metrics["visual_findings"] = visual_findings
    metrics["repair_recommendation"] = repair_recommendation
    metrics["density_flag"] = density_flag
    metrics["hierarchy_flag"] = hierarchy_flag
    metrics["visual_whitespace_ratio"] = (
        round(sum(per_slide_whitespace) / len(per_slide_whitespace), 4) if per_slide_whitespace else 0.0
    )
    metrics["visual_hierarchy_depth_score"] = (
        round(sum(per_slide_hierarchy_depth_score) / len(per_slide_hierarchy_depth_score), 4)
        if per_slide_hierarchy_depth_score
        else 0.0
    )
    metrics["visual_dominant_point_count"] = int(dominant_points_total)
    metrics["visual_repetition_penalty"] = round(repetition_penalty, 4)
    metrics["visual_alignment_quality_score"] = (
        round(sum(per_slide_alignment_score) / len(per_slide_alignment_score), 4)
        if per_slide_alignment_score
        else 0.0
    )
    metrics["visual_fragmentation_by_slide"] = fragmentation_by_slide
    return metrics


def _detect_takeaway_bar(root: ET.Element) -> bool:
    def _path_points(path_d: str) -> list[float]:
        raw = re.findall(r"[-+]?\d*\.?\d+", path_d)
        points: list[float] = []
        for token in raw:
            try:
                points.append(float(token))
            except ValueError:
                continue
        return points

    for elem in root.iter():
        name = local_name(elem.tag)
        if name == "rect":
            y = number(elem.get("y"))
            h = number(elem.get("height"))
            w = number(elem.get("width"))
            if y <= 220 and 3 <= h <= 18 and w >= 360:
                return True
            continue

        if name == "line":
            x1 = number(elem.get("x1"))
            x2 = number(elem.get("x2"))
            y1 = number(elem.get("y1"))
            y2 = number(elem.get("y2"))
            stroke_w = number(elem.get("stroke-width"), 1.0)
            if abs(x2 - x1) >= 360 and abs(y2 - y1) <= 16 and min(y1, y2) <= 220 and stroke_w <= 20:
                return True
            continue

        if name == "path":
            d = str(elem.get("d") or "").strip()
            if not d:
                continue
            points = _path_points(d)
            if len(points) < 4:
                continue
            xs = points[0::2]
            ys = points[1::2]
            if not xs or not ys:
                continue
            x_span = max(xs) - min(xs)
            y_span = max(ys) - min(ys)
            if x_span >= 360 and y_span <= 24 and min(ys) <= 220:
                return True
    return False


_PATH_COMMAND_PARAMS = {
    "M": 2,
    "L": 2,
    "H": 1,
    "V": 1,
    "A": 7,
    "Z": 0,
}
_PATH_TOKEN_RE = re.compile(
    r"[A-Za-z]|[-+]?(?:\d+(?:\.\d*)?|\.\d+)(?:[eE][-+]?\d+)?"
)
_PATH_LINEAR_EDGE_TOLERANCE = 1e-6


def _path_style_value(elem: ET.Element, name: str) -> str:
    value = elem.get(name)
    if value is not None:
        return value.strip()
    style = elem.get("style") or ""
    match = re.search(rf"(?:^|;)\s*{re.escape(name)}\s*:\s*([^;]+)", style, re.IGNORECASE)
    return match.group(1).strip() if match else ""


def _path_card_like_bbox(elem: ET.Element) -> tuple[float, float, float, float] | None:
    """Return an endpoint-based bbox for the small observed Path command subset.

    The first version deliberately supports only M/L/H/V/A/Z and skips paths
    with transforms or any unsupported command. Arc extrema are not solved; the
    observed card arcs are rounded-corner segments whose endpoints provide the
    intended conservative card bounds. A card candidate must also contain both
    a horizontal and a vertical linear edge; arc segments do not provide either.
    """
    if local_name(elem.tag) != "path":
        return None
    path_d = str(elem.get("d") or "").strip()
    if not path_d or elem.get("transform"):
        return None

    tokens = _PATH_TOKEN_RE.findall(path_d)
    if not tokens or "".join(tokens) != re.sub(r"[\s,]+", "", path_d):
        return None

    points: list[tuple[float, float]] = []
    current: tuple[float, float] | None = None
    subpath_start: tuple[float, float] | None = None
    subpath_closed = False
    subpath_count = 0
    has_horizontal_linear_edge = False
    has_vertical_linear_edge = False
    index = 0
    command: str | None = None

    def read_numbers(count: int) -> list[float] | None:
        nonlocal index
        if index + count > len(tokens):
            return None
        values: list[float] = []
        for token in tokens[index : index + count]:
            if token.isalpha():
                return None
            try:
                values.append(float(token))
            except ValueError:
                return None
        index += count
        return values

    while index < len(tokens):
        if tokens[index].isalpha():
            command = tokens[index]
            index += 1
        if command is None or command.upper() not in _PATH_COMMAND_PARAMS:
            return None

        upper = command.upper()
        relative = command.islower()
        if upper == "Z":
            if subpath_start is None or current is None:
                return None
            current = subpath_start
            points.append(current)
            subpath_closed = True
            command = None
            continue

        first_group = True
        while index < len(tokens) and not tokens[index].isalpha():
            values = read_numbers(_PATH_COMMAND_PARAMS[upper])
            if values is None:
                return None

            if upper == "M":
                if not first_group:
                    upper = "L"
                if upper == "M":
                    subpath_count += 1
                    if subpath_count > 1:
                        return None
                    if current is not None and not subpath_closed:
                        return None
                    x, y = values
                    if relative and current is not None:
                        x += current[0]
                        y += current[1]
                    current = (x, y)
                    subpath_start = current
                    subpath_closed = False
                    points.append(current)
                else:
                    x, y = values
                    if relative and current is not None:
                        x += current[0]
                        y += current[1]
                    current = (x, y)
                    points.append(current)
            elif upper == "L":
                x, y = values
                if current is None:
                    return None
                if relative:
                    x += current[0]
                    y += current[1]
                if (
                    math.isclose(y, current[1], rel_tol=0.0, abs_tol=_PATH_LINEAR_EDGE_TOLERANCE)
                    and not math.isclose(x, current[0], rel_tol=0.0, abs_tol=_PATH_LINEAR_EDGE_TOLERANCE)
                ):
                    has_horizontal_linear_edge = True
                if (
                    math.isclose(x, current[0], rel_tol=0.0, abs_tol=_PATH_LINEAR_EDGE_TOLERANCE)
                    and not math.isclose(y, current[1], rel_tol=0.0, abs_tol=_PATH_LINEAR_EDGE_TOLERANCE)
                ):
                    has_vertical_linear_edge = True
                current = (x, y)
                points.append(current)
            elif upper == "H":
                if current is None:
                    return None
                x = values[0] + (current[0] if relative else 0.0)
                if not math.isclose(
                    x, current[0], rel_tol=0.0, abs_tol=_PATH_LINEAR_EDGE_TOLERANCE
                ):
                    has_horizontal_linear_edge = True
                current = (x, current[1])
                points.append(current)
            elif upper == "V":
                if current is None:
                    return None
                y = values[0] + (current[1] if relative else 0.0)
                if not math.isclose(
                    y, current[1], rel_tol=0.0, abs_tol=_PATH_LINEAR_EDGE_TOLERANCE
                ):
                    has_vertical_linear_edge = True
                current = (current[0], y)
                points.append(current)
            elif upper == "A":
                if current is None:
                    return None
                if (
                    values[0] < 0
                    or values[1] < 0
                    or values[2] != 0.0
                    or values[3] not in {0.0, 1.0}
                    or values[4] not in {0.0, 1.0}
                ):
                    return None
                x, y = values[5], values[6]
                if relative:
                    x += current[0]
                    y += current[1]
                current = (x, y)
                points.append(current)
            first_group = False

        if upper != "M" and command.upper() == "M":
            command = command.replace("m", "l").replace("M", "L")

    if (
        not subpath_closed
        or len(points) < 2
        or not has_horizontal_linear_edge
        or not has_vertical_linear_edge
    ):
        return None
    if len({point for point in points[:-1]}) < 4:
        return None
    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    return min(xs), min(ys), max(xs), max(ys)


def _is_card_like_path(
    elem: ET.Element,
    canvas_width: float = 1280.0,
    canvas_height: float = 720.0,
) -> bool:
    if local_name(elem.tag) != "path":
        return False
    if any(
        _path_style_value(elem, name)
        for name in ("marker-start", "marker-mid", "marker-end")
    ):
        return False
    fill = _path_style_value(elem, "fill").strip().lower()
    if fill in {"", "none", "transparent"}:
        return False
    bbox = _path_card_like_bbox(elem)
    if bbox is None:
        return False
    x_min, y_min, x_max, y_max = bbox
    width = x_max - x_min
    height = y_max - y_min
    if width < 160 or height < 80:
        return False
    if width / max(height, 1.0) > 5.5:
        return False
    if (
        width >= canvas_width * 0.9
        and height >= canvas_height * 0.9
    ):
        return False
    return True


def _count_rect_card_like_blocks(root: ET.Element) -> int:
    count = 0
    for elem in root.iter():
        if local_name(elem.tag) != "rect":
            continue
        rx = number(elem.get("rx"))
        w = number(elem.get("width"))
        h = number(elem.get("height"))
        fill = (elem.get("fill") or "").strip().lower()
        if rx < 6:
            continue
        if w < 160 or h < 80:
            continue
        if fill in {"", "none", "transparent"}:
            continue
        if (w / max(h, 1.0)) > 5.5:
            continue
        count += 1
    return count


def _count_path_card_like_blocks(root: ET.Element) -> int:
    count = 0
    for elem in root.iter():
        if _is_card_like_path(elem):
            count += 1
    return count


def _card_like_block_counts(root: ET.Element) -> dict[str, int]:
    rect_count = _count_rect_card_like_blocks(root)
    path_count = _count_path_card_like_blocks(root)
    return {
        "rect": rect_count,
        "path": path_count,
        "combined": rect_count + path_count,
    }


def _count_card_like_blocks(root: ET.Element) -> int:
    """Compatibility wrapper: preserve the original rect-only behavior."""
    return _count_rect_card_like_blocks(root)


def _has_conclusion_headline(root: ET.Element) -> bool:
    for elem in root.iter():
        if local_name(elem.tag) != "text":
            continue
        y = number(elem.get("y"))
        font_size = number(elem.get("font-size"), 18)
        if y <= 170 and font_size >= CONCLUSION_HEADLINE_MIN_FONT:
            lines = text_lines(elem)
            if lines and any(line.strip() for line in lines):
                return True
    return False


def _conclusion_zone_anchors(root: ET.Element) -> list[tuple[float, str]]:
    """Conclusion-level text anchors, sorted by y (AB-05/AB-06).

    Reuses the conclusion-zone recognition of ``_has_conclusion_headline``
    (same font threshold) but scans the full page and requires sentence-length
    text so KPI numerals do not count as conclusion zones. Takeaway bars are
    deliberately not counted as standalone anchors: a top accent bar plus the
    headline is one zone, and counting bars alone would false-positive normal
    headline pages.
    """
    anchors: list[tuple[float, str]] = []
    for elem in root.iter():
        if local_name(elem.tag) != "text":
            continue
        font_size = number(elem.get("font-size"), 18)
        if font_size < CONCLUSION_HEADLINE_MIN_FONT:
            continue
        content = " ".join(line.strip() for line in text_lines(elem) if line.strip())
        if len(content) < CONCLUSION_ZONE_MIN_TEXT_CHARS:
            continue
        anchors.append((number(elem.get("y")), content))
    anchors.sort(key=lambda item: item[0])
    return anchors


def _count_conclusion_zones(root: ET.Element) -> int:
    """Count vertically separated conclusion-level zones on one slide."""
    zones = 0
    prev_y: float | None = None
    for y, _content in _conclusion_zone_anchors(root):
        if prev_y is None or (y - prev_y) > CONCLUSION_ZONE_CLUSTER_GAP:
            zones += 1
        prev_y = y
    return zones


def _conclusion_text_has_qualifier(root: ET.Element) -> bool:
    """True if the on-canvas conclusion carries an explicit qualifier (AB-06).

    Checks conclusion-zone texts when present; falls back to all page text when
    no conclusion-grade text is detected, so weak-headline pages are not
    double-flagged for a missing qualifier.
    """
    anchors = _conclusion_zone_anchors(root)
    if anchors:
        candidates = [content for _y, content in anchors]
    else:
        candidates = [
            " ".join(line.strip() for line in text_lines(elem) if line.strip())
            for elem in root.iter()
            if local_name(elem.tag) == "text"
        ]
    lowered = [item.lower() for item in candidates if item]
    return any(term in item for item in lowered for term in CLAIM_BOUNDARY_QUALIFIER_TERMS)


def _is_consulting_profile(profile: str) -> bool:
    lowered = profile.strip().lower()
    if not lowered:
        return False
    return any(token in lowered for token in CONSULTING_PROFILE_TOKENS)


def _normalize_hex_color(value: str) -> str | None:
    raw = value.strip().lower()
    if not raw or raw in {"none", "transparent", "currentcolor"}:
        return None
    if raw.startswith("url("):
        return None
    if raw.startswith("#"):
        token = raw[1:]
        if len(token) == 3 and re.fullmatch(r"[0-9a-f]{3}", token):
            return "#" + "".join(ch * 2 for ch in token)
        if len(token) == 6 and re.fullmatch(r"[0-9a-f]{6}", token):
            return raw
    return None


def _is_grayscale_hex(hex_color: str) -> bool:
    token = hex_color.lstrip("#")
    if len(token) != 6:
        return False
    return token[0:2] == token[2:4] == token[4:6]


def _hex_rgb_255(hex_color: str) -> tuple[int, int, int] | None:
    token = hex_color.strip().lstrip("#")
    if len(token) != 6 or not re.fullmatch(r"[0-9a-fA-F]{6}", token):
        return None
    return int(token[0:2], 16), int(token[2:4], 16), int(token[4:6], 16)


def _is_pale_surface_hex(hex_color: str) -> bool:
    rgb = _hex_rgb_255(hex_color)
    if rgb is None:
        return False
    # Pale tints are usually card/surface roles, not primary palette roles.
    return min(rgb) >= 210 and (sum(rgb) / 3.0) >= 232 and max(rgb) >= 245


def _collect_primary_colors(
    root: ET.Element,
    *,
    exclude_tokens: set[str],
) -> set[str]:
    colors: set[str] = set()
    ignore_white = "#ffffff" in exclude_tokens
    ignore_gray_scale = "gray_scale" in exclude_tokens
    for elem in root.iter():
        for attr in ("fill", "stroke"):
            raw = elem.get(attr)
            if not isinstance(raw, str):
                continue
            normalized = _normalize_hex_color(raw)
            if not normalized:
                continue
            lowered = normalized.lower()
            if lowered in exclude_tokens:
                continue
            if ignore_white and lowered == "#ffffff":
                continue
            if ignore_gray_scale and _is_grayscale_hex(lowered):
                continue
            if _is_pale_surface_hex(lowered):
                continue
            colors.add(lowered)
    return colors


def _style_support_color_excludes(
    project_dir: Path,
    color_system: dict[str, Any],
) -> set[str]:
    excludes: set[str] = set()

    theme = Theme.from_design_spec(project_dir / "design_spec.md")
    for color in (theme.background, theme.canvas_background, theme.card, theme.text, theme.muted, theme.line, theme.soft):
        normalized = _normalize_hex_color(str(color or ""))
        if normalized:
            excludes.add(normalized.lower())

    support_palette = color_system.get("support_palette")
    if isinstance(support_palette, list):
        for color in support_palette:
            normalized = _normalize_hex_color(str(color or ""))
            if normalized:
                excludes.add(normalized.lower())

    values = _read_design_spec_values(project_dir)
    for key, value in values.items():
        if not key.endswith("_color") and key not in {"card_bg", "background_color", "canvas_background", "line_color"}:
            continue
        normalized = _normalize_hex_color(value)
        if normalized and key not in {"primary_color", "accent_color", "secondary_accent"}:
            excludes.add(normalized.lower())

    return excludes


def _load_icons_index(project_dir: Path, findings: list[Finding]) -> set[str]:
    custom = os.environ.get("PPT_ICON_INDEX_PATH", "").strip()
    if custom:
        index_path = Path(custom)
    else:
        index_path = SCRIPT_DIR.parent / "assets" / "icons" / "icons_index.json"
    loaded = load_json_object(index_path, encoding="utf-8")
    if loaded.status == "missing":
        return set()
    if loaded.status in {"invalid", "read-failed"}:
        emit(findings, "warning", "invalid-icons-index-json", str(index_path), loaded.message)
        return set()
    if loaded.status == "schema_mismatch":
        emit(
            findings,
            "warning",
            "invalid-icons-index-shape",
            str(index_path),
            loaded.message,
        )
        return set()
    payload = loaded.payload if isinstance(loaded.payload, dict) else {}
    return {str(key) for key in payload.keys() if str(key).strip()}


def _extract_icon_refs(svg_text: str) -> set[str]:
    refs = set(ICON_REF_RE.findall(svg_text))
    refs.update(ICON_REF_ATTR_RE.findall(svg_text))
    return {item.strip() for item in refs if item.strip()}


def _collect_chart_palette_violations(root: ET.Element, allowed_palette: set[str]) -> set[str]:
    invalid: set[str] = set()
    for elem in root.iter():
        marker = str(elem.get("data-chart-fragment") or "")
        if not marker.strip():
            continue
        for node in elem.iter():
            for attr in ("fill", "stroke"):
                raw = node.get(attr)
                if not isinstance(raw, str):
                    continue
                normalized = _normalize_hex_color(raw)
                if not normalized:
                    continue
                if normalized in {"#ffffff", "#000000"}:
                    continue
                if normalized not in allowed_palette:
                    invalid.add(normalized)
    return invalid


def validate_icon_and_chart_asset_refs(
    project_dir: Path,
    svg_dir: Path,
    findings: list[Finding],
) -> dict[str, Any]:
    values = _read_design_spec_values(project_dir)
    data_palette = {item.lower() for item in parse_data_palette(values.get("data_palette"))}
    icons_index = _load_icons_index(project_dir, findings)

    icon_refs_checked = 0
    icon_refs_missing = 0
    chart_palette_violations = 0
    for svg_file in sorted(svg_dir.glob("slide_*.svg")):
        text = svg_file.read_text(encoding="utf-8-sig")
        refs = _extract_icon_refs(text)
        for ref in sorted(refs):
            icon_refs_checked += 1
            if icons_index and ref in icons_index:
                continue
            icon_refs_missing += 1
            emit(
                findings,
                "warning",
                "icon-ref-missing",
                str(svg_file),
                f"icon_ref '{ref}' was not found in icons_index.json.",
            )

        if not data_palette:
            continue
        try:
            root = ET.fromstring(text)
        except ET.ParseError:
            continue
        invalid_colors = _collect_chart_palette_violations(root, data_palette)
        if invalid_colors:
            chart_palette_violations += len(invalid_colors)
            emit(
                findings,
                "warning",
                "chart-color-outside-palette",
                str(svg_file),
                "Chart fragment uses colors outside data_palette: " + ", ".join(sorted(invalid_colors)),
            )

    return {
        "icon_ref_checked_count": icon_refs_checked,
        "icon_ref_missing_count": icon_refs_missing,
        "chart_color_outside_palette_count": chart_palette_violations,
        "data_palette_color_count": len(data_palette),
    }


def validate_style_hard_tokens(
    project_dir: Path,
    svg_dir: Path,
    slides: list[dict[str, Any]],
    findings: list[Finding],
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "style_hard_token_checks": {
            "enabled": False,
            "profile": "",
            "checked_slides": 0,
            "violations": {
                "color_limit": 0,
                "consecutive_homogeneous": 0,
                "conclusion_hierarchy": 0,
            },
            "pass_rate": 1.0,
        }
    }
    reference_pack = _load_json_file(project_dir / "reference_pack.json", findings, code="invalid-reference-pack-json")
    style_route = _load_json_file(project_dir / "style_route.json", findings, code="invalid-style-route-json")
    plan_payload = _load_json_file(
        project_dir / "slide_visual_plan.json",
        findings,
        code="invalid-slide-visual-plan-json",
    )

    profile = ""
    if isinstance(reference_pack, dict):
        profile = str(reference_pack.get("style_profile") or "")
    if isinstance(style_route, dict):
        route_profile = str(style_route.get("style_profile") or "")
        if route_profile:
            profile = route_profile
    metrics["style_hard_token_checks"]["profile"] = profile
    if not _is_consulting_profile(profile):
        return metrics

    execution_tokens: dict[str, Any] = {}
    if isinstance(reference_pack, dict):
        raw_tokens = reference_pack.get("execution_tokens")
        if isinstance(raw_tokens, dict):
            execution_tokens = raw_tokens
    if not execution_tokens and isinstance(style_route, dict):
        profile_tokens = style_route.get("profile_tokens")
        if isinstance(profile_tokens, dict):
            hard_tokens = profile_tokens.get("hard_tokens")
            if isinstance(hard_tokens, dict):
                execution_tokens = hard_tokens
    if not execution_tokens:
        return metrics

    metrics["style_hard_token_checks"]["enabled"] = True
    color_system_raw = execution_tokens.get("color_system")
    color_system: dict[str, Any] = color_system_raw if isinstance(color_system_raw, dict) else {}
    if isinstance(style_route, dict):
        profile_tokens = style_route.get("profile_tokens")
        if isinstance(profile_tokens, dict):
            hard_tokens = profile_tokens.get("hard_tokens")
            route_color_system = hard_tokens.get("color_system") if isinstance(hard_tokens, dict) else None
            if isinstance(route_color_system, dict):
                for key in ("support_palette", "exclude_from_primary_count"):
                    if key not in color_system and key in route_color_system:
                        color_system[key] = route_color_system[key]
    rhythm_rules_raw = execution_tokens.get("rhythm_rules")
    rhythm_rules: dict[str, Any] = rhythm_rules_raw if isinstance(rhythm_rules_raw, dict) else {}
    max_primary_colors = int(color_system.get("max_primary_colors_per_slide", 3) or 3)
    raw_excludes = color_system.get("exclude_from_primary_count", ["#FFFFFF", "gray_scale"])
    exclude_tokens = (
        {str(item).strip().lower() for item in raw_excludes}
        if isinstance(raw_excludes, list)
        else {"#ffffff", "gray_scale"}
    )
    if "#ffffff" not in exclude_tokens:
        exclude_tokens.add("#ffffff")
    exclude_tokens.update(_style_support_color_excludes(project_dir, color_system))

    by_slide: dict[int, Path] = {}
    for path in sorted(svg_dir.glob("slide_*.svg")):
        sid = _slide_id_from_svg_name(path)
        if sid is not None:
            by_slide[sid] = path

    checked = 0
    violations_color = 0
    key_layout_tags = {
        "Strategy-Map",
        "Capability-Mapping",
        "Roadmap-MultiPhase",
        "Data-Single-KPI",
        "Data-Three-KPIs",
        "Chart-Bar",
        "Chart-Line",
    }
    weak_conclusion: list[int] = []

    for slide in slides:
        sid = slide.get("id")
        if not isinstance(sid, int):
            continue
        svg_file = by_slide.get(sid)
        if not svg_file:
            continue
        try:
            root = ET.parse(svg_file).getroot()
        except (ET.ParseError, OSError):
            continue
        checked += 1
        primary_colors = _collect_primary_colors(root, exclude_tokens=exclude_tokens)
        if len(primary_colors) > max_primary_colors:
            violations_color += 1
            emit(
                findings,
                "warning",
                "style-hard-token-color-limit",
                str(svg_file),
                (
                    f"Slide {sid} uses {len(primary_colors)} primary colors; "
                    f"consulting hard-token limit is {max_primary_colors}."
                ),
            )
        layout_tag = str(slide.get("layout_tag") or "")
        if layout_tag in key_layout_tags and not _has_conclusion_headline(root):
            weak_conclusion.append(sid)

    if weak_conclusion:
        emit(
            findings,
            "warning",
            "style-hard-token-conclusion-hierarchy-weak",
            str(svg_dir),
            "Core conclusion hierarchy is weak on key slides: "
            + ", ".join(str(item) for item in weak_conclusion[:8]),
        )

    monotony_violation = 0
    max_same_rhythm = int(rhythm_rules.get("max_consecutive_same_rhythm_role", 2) or 2)
    if isinstance(plan_payload, dict):
        plan_slides = plan_payload.get("slides")
        if isinstance(plan_slides, list):
            streak = 0
            prev_signature = ""
            for item in plan_slides:
                if not isinstance(item, dict):
                    continue
                archetype = str(item.get("selected_archetype") or item.get("visual_archetype") or "").strip().lower()
                rhythm_role = str(item.get("rhythm_role") or "").strip().lower()
                signature = f"{archetype}|{rhythm_role}"
                if not signature.strip("|"):
                    continue
                if signature == prev_signature:
                    streak += 1
                else:
                    streak = 1
                    prev_signature = signature
                if streak > max_same_rhythm:
                    monotony_violation = 1
                    break
    if monotony_violation:
        emit(
            findings,
            "warning",
            "style-hard-token-consecutive-homogeneous",
            str(project_dir / "slide_visual_plan.json"),
            "Consecutive pages reuse the same archetype/rhythm signature above consulting hard-token threshold.",
        )

    total_checks = checked + 1 + 1
    total_violations = violations_color + (1 if monotony_violation else 0) + (1 if weak_conclusion else 0)
    pass_rate = 1.0 if total_checks <= 0 else max(0.0, min(1.0, (total_checks - total_violations) / total_checks))
    metrics["style_hard_token_checks"] = {
        "enabled": True,
        "profile": profile,
        "checked_slides": checked,
        "max_primary_colors_per_slide": max_primary_colors,
        "violations": {
            "color_limit": violations_color,
            "consecutive_homogeneous": int(monotony_violation),
            "conclusion_hierarchy": int(bool(weak_conclusion)),
        },
        "pass_rate": round(pass_rate, 3),
    }
    return metrics


def validate_profile_style_consistency(
    project_dir: Path,
    svg_dir: Path,
    slides: list[dict[str, Any]],
    findings: list[Finding],
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "style_profile": "",
        "style_profile_checks": {
            "takeaway_bar_present": False,
            "over_cardization_flag": False,
            "conclusion_first_flag": False,
            "rhythm_monotony_flag": False,
            "conclusion_competition_slides": [],
            "claim_boundary_unqualified_slides": [],
        }
    }
    reference_pack = _load_json_file(project_dir / "reference_pack.json", findings, code="invalid-reference-pack-json")
    plan_payload = _load_json_file(
        project_dir / "slide_visual_plan.json",
        findings,
        code="invalid-slide-visual-plan-json",
    )
    style_route = _load_json_file(project_dir / "style_route.json", findings, code="invalid-style-route-json")

    svg_files = sorted(svg_dir.glob("slide_*.svg"))
    by_slide: dict[int, Path] = {}
    for path in svg_files:
        sid = _slide_id_from_svg_name(path)
        if sid is not None:
            by_slide[sid] = path

    takeaway_required = False
    profile = ""
    if isinstance(reference_pack, dict):
        profile = str(reference_pack.get("style_profile") or "")
        motifs = reference_pack.get("motifs")
        if isinstance(motifs, list) and any("takeaway bar" in str(item).lower() for item in motifs):
            takeaway_required = True
    if isinstance(style_route, dict):
        route_profile = str(style_route.get("style_profile") or "")
        if route_profile:
            profile = route_profile
        if "strategy" in route_profile.lower():
            takeaway_required = True
    metrics["style_profile"] = profile

    takeaway_found = False
    card_heavy_slides = 0
    card_counts: list[int] = []
    rect_card_counts: list[int] = []
    path_card_counts: list[int] = []
    combined_card_counts: list[int] = []
    path_card_observed_slides: list[int] = []
    path_card_heavy_slides = 0
    conclusion_competition_slides: list[int] = []
    claim_boundary_unqualified_slides: list[int] = []
    core_total = 0
    core_without_headline = 0
    core_tags = {
        "Strategy-Map",
        "Capability-Mapping",
        "Roadmap-MultiPhase",
        "Chart-Bar",
        "Chart-Line",
        "Data-Single-KPI",
        "Data-Three-KPIs",
    }

    for slide in slides:
        sid = slide.get("id")
        if not isinstance(sid, int):
            continue
        svg_file = by_slide.get(sid)
        if not svg_file:
            continue
        try:
            root = ET.parse(svg_file).getroot()
        except (ET.ParseError, OSError):
            continue

        if _detect_takeaway_bar(root):
            takeaway_found = True
        card_counts_for_slide = _card_like_block_counts(root)
        card_count = card_counts_for_slide["rect"]
        path_card_count = card_counts_for_slide["path"]
        combined_card_count = card_counts_for_slide["combined"]
        rect_card_counts.append(card_count)
        path_card_counts.append(path_card_count)
        combined_card_counts.append(combined_card_count)
        card_counts.append(card_count)
        if path_card_count > 0:
            path_card_observed_slides.append(sid)
        if path_card_count >= 4:
            path_card_heavy_slides += 1
        if card_count >= 4:
            card_heavy_slides += 1

        tag = str(slide.get("layout_tag", ""))
        if tag in core_tags:
            core_total += 1
            if not _has_conclusion_headline(root):
                core_without_headline += 1

        zone_count = _count_conclusion_zones(root)
        if zone_count >= 2:
            conclusion_competition_slides.append(sid)
            emit(
                findings,
                "advisory",
                "style-conclusion-competition",
                str(svg_file),
                (
                    f"Slide {sid} has {zone_count} competing conclusion-level zones; "
                    "keep a single core conclusion zone per page."
                ),
            )

        boundary = str(slide.get("claim_boundary") or "").strip().lower()
        if boundary in {"assumption", "inference"} and not _conclusion_text_has_qualifier(root):
            claim_boundary_unqualified_slides.append(sid)
            emit(
                findings,
                "advisory",
                "style-claim-boundary-unqualified",
                str(svg_file),
                (
                    f"Slide {sid} claim_boundary is '{boundary}' but the on-canvas conclusion "
                    "has no explicit qualifier (e.g. 预计/计划/假设/拟/待验证); "
                    "do not present it as an established result."
                ),
            )

    metrics["style_profile_checks"]["takeaway_bar_present"] = takeaway_found
    metrics["style_profile_checks"]["conclusion_competition_slides"] = conclusion_competition_slides
    metrics["style_profile_checks"]["claim_boundary_unqualified_slides"] = claim_boundary_unqualified_slides
    if takeaway_required and not takeaway_found:
        emit(
            findings,
            "warning",
            "style-takeaway-bar-missing",
            str(project_dir / "reference_pack.json"),
            "Profile expects a takeaway bar motif, but no qualifying bar was detected in slide SVGs.",
        )

    if card_counts:
        avg_cards = sum(card_counts) / len(card_counts)
        heavy_ratio = card_heavy_slides / len(card_counts)
        avg_rect_cards = sum(rect_card_counts) / len(rect_card_counts)
        avg_path_cards = sum(path_card_counts) / len(path_card_counts)
        avg_combined_cards = sum(combined_card_counts) / len(combined_card_counts)
        path_heavy_ratio = path_card_heavy_slides / len(path_card_counts)
        combined_heavy_slides = sum(1 for count in combined_card_counts if count >= 4)
        combined_heavy_ratio = combined_heavy_slides / len(combined_card_counts)
        metrics["style_profile_checks"]["average_card_blocks"] = round(avg_cards, 2)
        metrics["style_profile_checks"]["average_rect_card_blocks"] = round(avg_rect_cards, 2)
        metrics["style_profile_checks"]["average_path_card_blocks"] = round(avg_path_cards, 2)
        metrics["style_profile_checks"]["average_combined_card_blocks"] = round(avg_combined_cards, 2)
        metrics["style_profile_checks"]["path_card_heavy_slide_ratio"] = round(path_heavy_ratio, 3)
        metrics["style_profile_checks"]["path_card_observed_slides"] = path_card_observed_slides
        # Avoid noisy judgments on tiny decks and only flag clearly card-heavy patterns.
        if len(card_counts) >= 3 and (avg_cards >= 4.5 or heavy_ratio >= (2.0 / 3.0)):
            metrics["style_profile_checks"]["over_cardization_flag"] = True
            emit(
                findings,
                "warning",
                "style-over-cardization",
                str(svg_dir),
                f"Card density is high across the deck (avg={avg_cards:.2f}, heavy_ratio={heavy_ratio:.2f}).",
            )
        elif (
            len(card_counts) >= 3
            and len(path_card_observed_slides) >= 2
            and max(path_card_counts, default=0) >= 2
            and (avg_combined_cards >= 4.5 or combined_heavy_ratio >= (2.0 / 3.0))
        ):
            emit(
                findings,
                "advisory",
                "style-path-cardization-observed",
                str(svg_dir),
                "Path-based card density was observed across the deck. "
                "This is advisory-only while false-positive behavior is being validated.",
            )

    if core_total > 0:
        weak_ratio = core_without_headline / core_total
        metrics["style_profile_checks"]["core_without_conclusion_headline_ratio"] = round(weak_ratio, 3)
        if weak_ratio >= 0.4:
            metrics["style_profile_checks"]["conclusion_first_flag"] = True
            emit(
                findings,
                "warning",
                "style-conclusion-first-weak",
                str(svg_dir),
                "Conclusion-first signal is weak on core strategy/data slides (headline not detected in top zone).",
            )

    rhythm_roles: list[str] = []
    if isinstance(plan_payload, dict):
        plan_slides = plan_payload.get("slides")
        if isinstance(plan_slides, list):
            for item in plan_slides:
                if not isinstance(item, dict):
                    continue
                role = item.get("rhythm_role")
                if isinstance(role, str) and role.strip():
                    rhythm_roles.append(role.strip())
    if len(rhythm_roles) >= 2:
        counts: dict[str, int] = {}
        for role in rhythm_roles:
            counts[role] = counts.get(role, 0) + 1
        dominant = max(counts.values())
        ratio = dominant / len(rhythm_roles)
        metrics["style_profile_checks"]["rhythm_roles"] = counts
        if ratio >= 0.7:
            metrics["style_profile_checks"]["rhythm_monotony_flag"] = True
            emit(
                findings,
                "warning",
                "style-rhythm-monotony",
                str(project_dir / "slide_visual_plan.json"),
                f"Deck rhythm is monotonic (dominant rhythm role ratio={ratio:.2f}).",
            )
    return metrics


def validate_visual_diversity_plan(project_dir: Path, findings: list[Finding]) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "visual_diversity_checked_slides": 0,
        "visual_diversity_warning_count": 0,
        "visual_diversity_report": None,
        "layout_exploration_enabled": False,
        "layout_candidate_count": 0,
        "archetype_switch_count": 0,
        "consecutive_repeat_count": 0,
        "diversity_gate_result": None,
    }
    plan_path = project_dir / "slide_visual_plan.json"
    if not plan_path.exists():
        return metrics
    try:
        report = check_visual_diversity(project_dir)
    except (FileNotFoundError, ValueError, OSError) as exc:
        emit(
            findings,
            "warning",
            "visual-diversity-check-failed",
            str(plan_path),
            f"Could not run visual diversity check: {exc}",
        )
        return metrics

    metrics["visual_diversity_checked_slides"] = report.checked_slides
    metrics["visual_diversity_warning_count"] = len(report.warnings)
    metrics["visual_diversity_report"] = str(report.report_path)
    metrics["layout_exploration_enabled"] = bool(report.layout_exploration_enabled)
    metrics["layout_candidate_count"] = int(report.candidate_count)
    metrics["archetype_switch_count"] = int(report.archetype_switch_count)
    metrics["consecutive_repeat_count"] = int(report.consecutive_repeat_count)
    metrics["diversity_gate_result"] = str(report.diversity_gate_result)
    for warning in report.warnings:
        emit(findings, "warning", "visual-diversity", str(plan_path), warning)
    return metrics


def validate_visual_baseline_delivery_gate(
    project_dir: Path,
    findings: list[Finding],
    *,
    quality_mode: str,
) -> dict[str, Any]:
    metrics: dict[str, Any] = {
        "visual_baseline_grade": None,
        "visual_baseline_gate_applied": False,
        "visual_baseline_gate_skipped_reason": None,
    }
    if quality_mode not in {"release-safe", "premium"}:
        return metrics
    if (project_dir / "workbench_status.json").exists():
        metrics["visual_baseline_gate_skipped_reason"] = "workbench_project"
        return metrics
    grade = baseline_grade_for_project(project_dir.name)
    if grade is None:
        return metrics
    metrics["visual_baseline_grade"] = grade
    metrics["visual_baseline_gate_applied"] = True
    if grade in {"C", "D"}:
        emit(
            findings,
            "warning",
            "visual-baseline-below-b",
            str(project_dir),
            (
                f"Project baseline grade is {grade}; release-safe delivery requires >= B. "
                "Treat this as visual non-delivery until page quality is raised."
            ),
        )
    return metrics


def validate_pptx(pptx_path: Path, expected_slides: int, findings: list[Finding]) -> dict[str, Any]:
    metrics: dict[str, Any] = {"pptx": str(pptx_path), "exists": pptx_path.exists()}
    if not pptx_path.exists():
        emit(findings, "warning", "missing-pptx", str(pptx_path), "PPTX export does not exist; skipping PPTX checks.")
        return metrics
    try:
        from pptx import Presentation
        from pptx.exc import InvalidXmlError, PackageNotFoundError
    except ImportError as exc:
        emit(findings, "warning", "pptx-lib-unavailable", str(pptx_path), f"Could not import python-pptx: {exc}")
        return metrics

    try:
        prs = Presentation(pptx_path)
    except (PackageNotFoundError, InvalidXmlError, zipfile.BadZipFile, OSError, ValueError) as exc:
        emit(findings, "error", "invalid-pptx", str(pptx_path), f"Could not open PPTX: {exc}")
        return metrics

    text_shapes = sum(
        1 for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
    )
    pictures = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    shape_count = sum(len(slide.shapes) for slide in prs.slides)
    metrics.update(
        {
            "slides": len(prs.slides),
            "shapes": shape_count,
            "text_shapes": text_shapes,
            "pictures": pictures,
            "size": pptx_path.stat().st_size,
        }
    )
    if len(prs.slides) != expected_slides:
        emit(
            findings,
            "error",
            "pptx-slide-count-mismatch",
            str(pptx_path),
            f"Expected {expected_slides} slides, found {len(prs.slides)}.",
        )
    if "native" in pptx_path.stem and pictures:
        emit(
            findings,
            "warning",
            "native-pptx-has-pictures",
            str(pptx_path),
            f"Native PPTX contains {pictures} picture shape(s).",
        )
    if "native" in pptx_path.stem and text_shapes == 0:
        emit(findings, "error", "native-pptx-no-text", str(pptx_path), "Native PPTX has no editable text shapes.")
    return metrics


def explicit_pptx_paths(project_dir: Path, pptx: str | Path | Sequence[str | Path]) -> list[Path]:
    raw_paths: list[str | Path]
    if isinstance(pptx, (str, Path)):
        raw_paths = [pptx]
    else:
        raw_paths = list(pptx)

    paths: list[Path] = []
    for raw_path in raw_paths:
        path = Path(raw_path)
        paths.append(path if path.is_absolute() else project_dir / path)
    return paths


def explicit_pptx_metric_key(pptx_path: Path, index: int, seen: set[str]) -> str:
    stem = pptx_path.stem
    if "native" in stem:
        base = "native_pptx"
    elif "raster" in stem or stem.endswith("_svg"):
        base = "raster_pptx"
    else:
        base = "pptx"

    key = base
    if key in seen:
        key = f"{base}_{index}"
    seen.add(key)
    return key


def render_snapshots(
    svg_dir: Path,
    qa_dir: Path,
    findings: list[Finding],
    slide_id: int | None = None,
) -> dict[str, Any]:
    metrics: dict[str, Any] = {"snapshots": 0}
    try:
        from PIL import Image, ImageDraw
        from svg_to_pptx import render_svg_to_png
    except ImportError as exc:
        emit(
            findings,
            "warning",
            "snapshot-deps-unavailable",
            str(svg_dir),
            f"Could not import snapshot dependencies: {exc}",
        )
        return metrics

    snapshot_dir = qa_dir / "snapshots"
    snapshot_dir.mkdir(parents=True, exist_ok=True)
    images: list[Path] = []
    target_name = f"slide_{slide_id:02d}.svg" if slide_id is not None else None
    for svg_file in sorted(svg_dir.glob("slide_*.svg")):
        if target_name and svg_file.name != target_name:
            continue
        png_file = snapshot_dir / f"{svg_file.stem}.png"
        try:
            render_svg_to_png(svg_file, png_file)
        except (OSError, RuntimeError, ValueError, ET.ParseError) as exc:
            emit(findings, "warning", "snapshot-render-failed", str(svg_file), f"Could not render snapshot: {exc}")
            continue
        images.append(png_file)

    metrics["snapshots"] = len(images)
    if not images:
        return metrics

    thumb_w, thumb_h = 512, 288
    label_h = 28
    cols = 2
    rows = (len(images) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + label_h)), "white")
    draw = ImageDraw.Draw(sheet)
    resampling = getattr(Image, "Resampling", None)
    lanczos_filter = getattr(resampling, "LANCZOS", getattr(Image, "LANCZOS", 3))
    for idx, path in enumerate(images):
        img = Image.open(path).convert("RGB").resize((thumb_w, thumb_h), lanczos_filter)
        x = (idx % cols) * thumb_w
        y = (idx // cols) * (thumb_h + label_h)
        sheet.paste(img, (x, y + label_h))
        draw.rectangle([x, y, x + thumb_w, y + label_h], fill=(245, 247, 250))
        draw.text((x + 10, y + 7), f"Slide {idx + 1:02d}", fill=(20, 23, 31))

    sheet_path = qa_dir / "contact-sheet.png"
    sheet.save(sheet_path)
    metrics["contact_sheet"] = str(sheet_path)
    return metrics


def run_qa(
    project_dir: Path,
    svg_dir_name: str | None = None,
    pptx: str | Path | Sequence[str | Path] | None = None,
    snapshots: bool = False,
    strict: bool = False,
    slide_id: int | None = None,
    enforce_blueprint_sync: bool = False,
    safe_area_profile: str = "legacy",
    enable_visual_qa: bool = False,
    profile: str = "presentation",
    quality_mode: str = "dev-fast",
    safe_edge_whitelist: set[str] | None = None,
) -> QaReport:
    project_dir = project_dir.resolve()
    qa_dir = project_dir / "qa"
    canvas = load_canvas_context(project_dir)
    findings: list[Finding] = []
    metrics: dict[str, Any] = {}
    if safe_edge_whitelist is None:
        safe_edge_whitelist = {"header", "footer"}

    blueprint = load_blueprint(project_dir, findings)
    tags = validate_layout_contracts(project_dir, findings)
    slides = validate_blueprint(blueprint, findings, tags)
    metrics.update(validate_blueprint_source_refs(project_dir, slides, findings))
    metrics.update(validate_art_direction_consumption(project_dir, slides, findings))
    metrics.update(
        check_visual_contracts(
            project_dir,
            findings,
            quality_mode=quality_mode,
            profile=profile,
        )
    )
    if slide_id is not None:
        slides = [slide for slide in slides if isinstance(slide.get("id"), int) and slide.get("id") == slide_id]
        if not slides:
            emit(
                findings,
                "error",
                "missing-slide-id",
                "blueprint.json",
                f"Requested --slide {slide_id} does not exist in blueprint.json.",
            )
    metrics["blueprint_slides"] = len(slides)
    metrics["supported_layouts"] = len(tags)
    metrics["layout_contracts"] = len(layout_tags())
    metrics["checked_slide"] = slide_id if slide_id is not None else "all"
    metrics["enforce_blueprint_sync"] = enforce_blueprint_sync
    metrics["safe_area_profile"] = safe_area_profile
    normalized_quality_mode = _normalize_quality_mode(quality_mode)
    metrics["quality_mode"] = normalized_quality_mode
    metrics["enable_visual_qa"] = enable_visual_qa
    metrics["canvas_key"] = canvas.key
    metrics["canvas_width"] = canvas.width
    metrics["canvas_height"] = canvas.height
    metrics["canvas_viewbox"] = canvas.viewbox
    metrics["profile"] = profile
    metrics.update(
        validate_template_consistency(
            project_dir,
            findings,
            quality_mode=normalized_quality_mode,
            profile=profile,
        )
    )
    metrics.update(validate_banned_terms(project_dir, slides, findings))
    metrics.update(validate_content_coverage(slides, findings))
    metrics.update(validate_slide_budget(slides, findings, profile=profile))
    metrics.update(validate_data_palette_presence(project_dir, slides, findings))
    context_profile = _qa_context_profile(project_dir, slide_id=slide_id)
    budget_stage = "executor_per_slide" if slide_id is not None else "checks"
    budget_limit = int(TOKEN_BUDGET_CAPS.get(budget_stage, TOKEN_BUDGET_CAPS["checks"]))
    token_estimate = _estimate_tokens_from_bytes(int(context_profile["context_bytes_estimate"]))
    token_overflow = max(0, token_estimate - budget_limit)
    metrics["token_budget_policy"] = TOKEN_BUDGET_POLICY
    metrics["token_budget_stage"] = budget_stage
    metrics["token_budget_limit"] = budget_limit
    metrics["context_file_count"] = context_profile["context_file_count"]
    metrics["context_bytes_estimate"] = context_profile["context_bytes_estimate"]
    metrics["context_files"] = context_profile["context_files"]
    metrics["context_token_estimate"] = token_estimate
    metrics["token_budget_overflow"] = token_overflow
    metrics["token_budget_warning"] = token_overflow > 0
    if token_overflow > 0:
        emit(
            findings,
            "advisory",
            "token-budget-over-cap",
            str(project_dir),
            (
                "Estimated QA context token usage exceeds stage budget "
                f"(stage={budget_stage}, estimate={token_estimate}, limit={budget_limit}, overflow={token_overflow})."
            ),
        )
    metrics.update(validate_theme_readability(project_dir, findings, profile=profile))
    metrics.update(validate_design_token_guard(project_dir, findings, profile=profile))

    svg_dir = choose_svg_dir(project_dir, svg_dir_name)
    contract_svg_dir = choose_structural_contract_svg_dir(project_dir, svg_dir)
    slide_plan = parse_slide_plan(project_dir, findings)
    metrics["structural_contract_svg_dir"] = str(contract_svg_dir)
    metrics.update(
        check_comparison_matrix_summarybar_svg_contract(
            project_dir,
            slides,
            contract_svg_dir,
            slide_plan,
            findings,
            slide_id=slide_id,
        )
    )
    metrics.update(
        check_core_orbit_relationship_svg_contract(
            project_dir,
            slides,
            contract_svg_dir,
            slide_plan,
            findings,
            slide_id=slide_id,
        )
    )
    metrics.update(
        check_roadmap_lane_milestones_svg_contract(
            project_dir,
            slides,
            contract_svg_dir,
            slide_plan,
            findings,
            slide_id=slide_id,
        )
    )
    metrics.update(
        check_architecture_three_zones_svg_contract(
            project_dir,
            slides,
            contract_svg_dir,
            slide_plan,
            findings,
            slide_id=slide_id,
        )
    )
    metrics.update(validate_planning_consistency(project_dir, slides, svg_dir, findings))
    metrics.update(
        validate_svg_set(
            project_dir,
            canvas,
            svg_dir,
            slides,
            findings,
            slide_plan=slide_plan,
            slide_id=slide_id,
            enforce_blueprint_sync=enforce_blueprint_sync,
            safe_area_profile=safe_area_profile,
            safe_edge_whitelist=safe_edge_whitelist,
            profile=profile,
        )
    )
    metrics.update(validate_style_hard_tokens(project_dir, svg_dir, slides, findings))
    metrics.update(validate_icon_and_chart_asset_refs(project_dir, svg_dir, findings))
    metrics.update(validate_profile_style_consistency(project_dir, svg_dir, slides, findings))
    metrics.update(validate_visual_diversity_plan(project_dir, findings))
    metrics.update(
        validate_visual_baseline_delivery_gate(
            project_dir,
            findings,
            quality_mode=normalized_quality_mode,
        )
    )
    aesthetic_result = run_aesthetic_critic(
        project_dir,
        svg_dir,
        profile=profile,
        quality_mode=normalized_quality_mode,
        slide_id=slide_id,
    )
    metrics["aesthetic_critic_enabled"] = bool(aesthetic_result.get("enabled"))
    metrics["aesthetic_critic_checked_slides"] = int(aesthetic_result.get("checked_slides") or 0)
    metrics["aesthetic_critic_finding_count"] = int(aesthetic_result.get("finding_count") or 0)
    for item in aesthetic_result.get("findings", []):
        if not isinstance(item, dict):
            continue
        findings.append(
            Finding(
                str(item.get("severity") or "advisory"),
                str(item.get("code") or "design-aesthetic-finding"),
                str(item.get("path") or project_dir),
                str(item.get("message") or ""),
            )
        )

    if slide_id is not None and not pptx:
        # Single-slide QA focuses on blueprint/SVG/snapshot iteration speed and does not
        # require full-deck PPTX slide-count parity checks.
        metrics["pptx_checks"] = "skipped-for-single-slide"
    elif pptx:
        explicit_paths = explicit_pptx_paths(project_dir, pptx)
        metrics["checked_pptx_paths"] = [str(path) for path in explicit_paths]
        seen_metric_keys: set[str] = set()
        for index, explicit_path in enumerate(explicit_paths, start=1):
            metric_key = explicit_pptx_metric_key(explicit_path, index, seen_metric_keys)
            metrics[metric_key] = validate_pptx(explicit_path, len(slides), findings)
    else:
        native = project_dir / "exports" / "output-native.pptx"
        raster = project_dir / "exports" / "output.pptx"
        checked_paths: list[str] = []
        if not native.exists():
            semantic_native = _latest_semantic_pptx(project_dir, "native")
            if semantic_native:
                native = semantic_native
        if not raster.exists():
            semantic_raster = _latest_semantic_pptx(project_dir, "raster")
            if semantic_raster:
                raster = semantic_raster
        if native.exists():
            checked_paths.append(str(native))
            metrics["native_pptx"] = validate_pptx(native, len(slides), findings)
        if raster.exists():
            checked_paths.append(str(raster))
            metrics["raster_pptx"] = validate_pptx(raster, len(slides), findings)
        metrics["checked_pptx_paths"] = checked_paths

    visual_metrics: VisualMetrics = {
        "visual_score": None,
        "visual_findings": [],
        "repair_recommendation": [],
        "density_flag": False,
        "hierarchy_flag": False,
        "visual_whitespace_ratio": 0.0,
        "visual_hierarchy_depth_score": 0.0,
        "visual_dominant_point_count": 0,
        "visual_repetition_penalty": 0.0,
        "visual_alignment_quality_score": 0.0,
        "visual_fragmentation_by_slide": {},
    }
    snapshot_metrics, visual_metrics, snapshot_findings, visual_findings_buffer = run_snapshot_and_visual_checks(
        snapshots=snapshots,
        enable_visual_qa=enable_visual_qa,
        svg_dir=svg_dir,
        qa_dir=qa_dir,
        project_dir=project_dir,
        slide_id=slide_id,
        profile=profile,
        canvas=canvas,
        visual_metrics_default=visual_metrics,
        render_snapshots_fn=render_snapshots,
        validate_visual_quality_fn=validate_visual_quality,
    )

    findings.extend(snapshot_findings)
    findings.extend(visual_findings_buffer)
    if snapshot_metrics:
        metrics.update(snapshot_metrics)
    if enable_visual_qa:
        metrics["visual_score"] = visual_metrics["visual_score"]
        metrics["density_flag"] = visual_metrics["density_flag"]
        metrics["hierarchy_flag"] = visual_metrics["hierarchy_flag"]
        metrics["visual_whitespace_ratio"] = visual_metrics.get("visual_whitespace_ratio", 0.0)
        metrics["visual_hierarchy_depth_score"] = visual_metrics.get("visual_hierarchy_depth_score", 0.0)
        metrics["visual_dominant_point_count"] = visual_metrics.get("visual_dominant_point_count", 0)
        metrics["visual_repetition_penalty"] = visual_metrics.get("visual_repetition_penalty", 0.0)
        metrics["visual_alignment_quality_score"] = visual_metrics.get("visual_alignment_quality_score", 0.0)
        metrics["visual_fragmentation_by_slide"] = visual_metrics.get("visual_fragmentation_by_slide", {})
        metrics["visual_finding_count"] = len(visual_metrics["visual_findings"])
        metrics["repair_recommendation_count"] = len(visual_metrics["repair_recommendation"])

    _apply_quality_mode(findings, normalized_quality_mode)
    proposal_critic_codes = _apply_visual_critic_gate(findings, profile=profile)
    errors = sum(1 for item in findings if item.severity == "error")
    warnings = sum(1 for item in findings if item.severity == "warning")
    advisories = sum(1 for item in findings if item.severity == "advisory")
    warning_items = [item for item in findings if item.severity == "warning"]
    release_non_blocking_warning_count = 0
    release_non_blocking_warning_codes: list[str] = []
    budget_non_blocking = [item for item in warning_items if item.code in BUDGET_NON_BLOCKING_WARNING_CODES]
    budget_non_blocking_warning_count = len(budget_non_blocking)
    budget_non_blocking_warning_codes = sorted({item.code for item in budget_non_blocking})
    design_token_non_blocking = [item for item in warning_items if item.code in DESIGN_TOKEN_NON_BLOCKING_WARNING_CODES]
    design_token_non_blocking_warning_count = len(design_token_non_blocking)
    design_token_non_blocking_warning_codes = sorted({item.code for item in design_token_non_blocking})
    asset_non_blocking = [item for item in warning_items if item.code in ASSET_NON_BLOCKING_WARNING_CODES]
    asset_non_blocking_warning_count = len(asset_non_blocking)
    asset_non_blocking_warning_codes = sorted({item.code for item in asset_non_blocking})
    blocking_warning_items = [
        item for item in warning_items if not _is_warning_non_blocking(item.code, normalized_quality_mode)
    ]
    blocking_warnings = len(blocking_warning_items)
    if normalized_quality_mode == "release-safe":
        non_blocking = [item for item in warning_items if item.code in RELEASE_SAFE_NON_BLOCKING_WARNING_CODES]
        release_non_blocking_warning_count = len(non_blocking)
        release_non_blocking_warning_codes = sorted({item.code for item in non_blocking})
    metrics["quality_tiers"] = {
        "blocking": errors,
        "warning": warnings,
        "advisory": advisories,
    }
    metrics["release_safe_non_blocking_warning_count"] = release_non_blocking_warning_count
    metrics["release_safe_non_blocking_warning_codes"] = release_non_blocking_warning_codes
    metrics["budget_non_blocking_warning_count"] = budget_non_blocking_warning_count
    metrics["budget_non_blocking_warning_codes"] = budget_non_blocking_warning_codes
    metrics["design_token_non_blocking_warning_count"] = design_token_non_blocking_warning_count
    metrics["design_token_non_blocking_warning_codes"] = design_token_non_blocking_warning_codes
    metrics["asset_non_blocking_warning_count"] = asset_non_blocking_warning_count
    metrics["asset_non_blocking_warning_codes"] = asset_non_blocking_warning_codes
    metrics["blocking_warning_count"] = blocking_warnings
    metrics["blocking_warning_codes"] = sorted({item.code for item in blocking_warning_items})
    metrics["blocking_visual_warning_count"] = sum(
        1 for item in blocking_warning_items if _is_visual_delivery_code(item.code)
    )
    metrics["blocking_visual_warning_codes"] = sorted(
        {item.code for item in blocking_warning_items if _is_visual_delivery_code(item.code)}
    )
    metrics["visual_critic_gate_profile"] = profile
    metrics["visual_critic_gate_blocking_codes"] = proposal_critic_codes
    metrics["visual_critic_gate_blocking_count"] = len(proposal_critic_codes)
    strict_effective = strict or normalized_quality_mode in {"release-safe", "premium"}
    metrics["strict_effective"] = strict_effective
    ok = errors == 0 and (blocking_warnings == 0 if strict_effective else True)
    report = QaReport(
        str(project_dir),
        ok,
        errors,
        warnings,
        findings,
        metrics,
        visual_score=visual_metrics["visual_score"],
        visual_findings=visual_metrics["visual_findings"],
        repair_recommendation=visual_metrics["repair_recommendation"],
        density_flag=visual_metrics["density_flag"],
        hierarchy_flag=visual_metrics["hierarchy_flag"],
        advisories=advisories,
    )
    write_reports_output(report, qa_dir)
    return report


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Run quality checks for an AI-PPT project.")
    parser.add_argument("project_dir", type=Path, help="Path to projects/<project_name>.")
    parser.add_argument("--svg-dir", help="Relative SVG directory to check. Defaults to svg_final, then svg_output.")
    parser.add_argument("--pptx", help="Relative PPTX path to check.")
    parser.add_argument("--snapshots", action="store_true", help="Render PNG snapshots and a contact sheet.")
    parser.add_argument("--strict", action="store_true", help="Treat warnings as failures.")
    parser.add_argument("--slide", type=int, help="Check only one slide id (for quick single-slide review).")
    parser.add_argument(
        "--enforce-blueprint-sync",
        action="store_true",
        help="Treat blueprint-to-SVG mismatch (missing/extra files) as blocking errors.",
    )
    parser.add_argument(
        "--safe-area-profile",
        choices=("legacy", "presentation"),
        default="legacy",
        help="Safe-area warning profile. presentation is stricter near slide edges.",
    )
    parser.add_argument(
        "--enable-visual-qa",
        action="store_true",
        help="Enable visual quality scoring and repair recommendations.",
    )
    parser.add_argument(
        "--profile",
        choices=("presentation", "print_a4", "proposal_consulting"),
        default="presentation",
        help="Governance profile for readability and visual-density thresholds.",
    )
    parser.add_argument(
        "--quality-mode",
        choices=QUALITY_MODES,
        default="dev-fast",
        help=(
            "Severity profile: dev-fast keeps momentum; premium blocks warnings; "
            "release-safe keeps selected hard-token warnings non-blocking."
        ),
    )
    parser.add_argument(
        "--safe-edge-whitelist",
        default="header,footer",
        help="Comma-separated regions allowed near safe edge (header,footer).",
    )
    args = parser.parse_args(argv)

    report = run_qa(
        args.project_dir,
        args.svg_dir,
        args.pptx,
        args.snapshots,
        args.strict,
        args.slide,
        enforce_blueprint_sync=args.enforce_blueprint_sync,
        safe_area_profile=args.safe_area_profile,
        enable_visual_qa=args.enable_visual_qa,
        profile=args.profile,
        quality_mode=args.quality_mode,
        safe_edge_whitelist={token.strip() for token in args.safe_edge_whitelist.split(",") if token.strip()},
    )
    print(f"QA {'passed' if report.ok else 'failed'}: errors={report.errors}, warnings={report.warnings}")
    print(Path(report.project) / "qa" / "report.md")
    return 0 if report.ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
