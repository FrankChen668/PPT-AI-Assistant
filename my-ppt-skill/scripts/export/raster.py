from __future__ import annotations

import os
import re
import shutil
import subprocess
import tempfile
import uuid
from contextlib import contextmanager
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

PIXEL_W = 1280
PIXEL_H = 720
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5


@contextmanager
def writable_temp_dir(prefix: str):
    base = Path(tempfile.gettempdir())
    last_error: Exception | None = None

    for _ in range(10):
        tmp_dir = base / f"{prefix}{uuid.uuid4().hex}"
        try:
            tmp_dir.mkdir(parents=True, exist_ok=False)
            probe = tmp_dir / ".write_probe"
            probe.write_text("ok", encoding="utf-8")
            probe.unlink()
            break
        except Exception as exc:
            last_error = exc
            shutil.rmtree(tmp_dir, ignore_errors=True)
    else:
        raise RuntimeError(f"Could not create writable temp dir under {base}: {last_error}")

    try:
        yield tmp_dir
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def browser_candidates() -> list[Path]:
    candidates: list[Path] = []

    env_browser = os.environ.get("AI_PPT_BROWSER")
    if env_browser:
        candidates.append(Path(env_browser))

    for name in ("msedge", "chrome", "chromium", "google-chrome"):
        found = shutil.which(name)
        if found:
            candidates.append(Path(found))

    candidates.extend(
        [
            Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
            Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
        ]
    )

    unique: list[Path] = []
    seen: set[str] = set()
    for candidate in candidates:
        key = str(candidate).lower()
        if key not in seen and candidate.exists():
            unique.append(candidate)
            seen.add(key)
    return unique


def render_with_cairosvg(svg_file: Path, png_file: Path) -> Exception | None:
    try:
        import cairosvg

        cairosvg.svg2png(
            url=str(svg_file),
            write_to=str(png_file),
            output_width=PIXEL_W,
            output_height=PIXEL_H,
        )
        return None
    except Exception as exc:
        return exc


def render_with_svglib_pdf(svg_file: Path, png_file: Path) -> Exception | None:
    try:
        import fitz
        from reportlab.graphics import renderPDF
        from svglib.svglib import svg2rlg

        png_file = png_file.resolve()
        png_file.parent.mkdir(parents=True, exist_ok=True)

        with writable_temp_dir("ai_ppt_svg_pdf_") as tmp_dir:
            pdf_file = tmp_dir / "render.pdf"
            drawing = svg2rlg(str(svg_file))
            if drawing is None:
                return RuntimeError("svglib could not parse SVG")
            renderPDF.drawToFile(drawing, str(pdf_file))

            doc = fitz.open(str(pdf_file))
            try:
                page = doc[0]
                scale_x = PIXEL_W / page.rect.width
                scale_y = PIXEL_H / page.rect.height
                pix = page.get_pixmap(matrix=fitz.Matrix(scale_x, scale_y), alpha=False)
                pix.save(str(png_file))
            finally:
                doc.close()

        return None
    except Exception as exc:
        return exc


def render_with_browser(svg_file: Path, png_file: Path) -> Exception | None:
    png_file = png_file.resolve()
    png_file.parent.mkdir(parents=True, exist_ok=True)
    browsers = browser_candidates()
    if not browsers:
        return FileNotFoundError("No local browser found. Install Microsoft Edge or Chrome, or set AI_PPT_BROWSER.")

    last_error: Exception | None = None
    with writable_temp_dir("ai_ppt_svg_browser_") as tmp_dir:
        wrapper = tmp_dir / "render.html"
        svg_markup = svg_file.read_text(encoding="utf-8-sig")
        svg_markup = re.sub(r"^\s*<\?xml[^>]*>\s*", "", svg_markup)
        wrapper.write_text(
            "\n".join(
                [
                    "<!doctype html>",
                    "<html>",
                    "<head>",
                    '<meta charset="utf-8">',
                    f'<base href="{svg_file.parent.resolve().as_uri()}/">',
                    "<style>",
                    (
                        "html, body { margin: 0; width: 1280px; height: 720px; "
                        "overflow: hidden; background: transparent; }"
                    ),
                    "svg { display: block; width: 1280px; height: 720px; }",
                    "</style>",
                    "</head>",
                    f"<body>{svg_markup}</body>",
                    "</html>",
                ]
            ),
            encoding="utf-8",
        )

        screenshot_path = tmp_dir / "screenshot.png"
        browser_profile = tmp_dir / "browser-profile"
        browser_cache = tmp_dir / "browser-cache"
        browser_crashes = tmp_dir / "browser-crashes"
        browser_profile.mkdir()
        browser_cache.mkdir()
        browser_crashes.mkdir()
        for browser in browsers:
            cmd = [
                str(browser),
                "--headless=new",
                "--disable-gpu",
                "--disable-breakpad",
                "--disable-crash-reporter",
                "--disable-crashpad",
                "--disable-dev-shm-usage",
                "--hide-scrollbars",
                "--no-sandbox",
                "--no-first-run",
                "--no-default-browser-check",
                "--noerrdialogs",
                "--allow-file-access-from-files",
                "--force-device-scale-factor=1",
                f"--user-data-dir={browser_profile}",
                f"--disk-cache-dir={browser_cache}",
                f"--crash-dumps-dir={browser_crashes}",
                f"--window-size={PIXEL_W},{PIXEL_H + 120}",
                f"--screenshot={screenshot_path}",
                wrapper.resolve().as_uri(),
            ]
            try:
                result = subprocess.run(
                    cmd,
                    check=False,
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                    timeout=45,
                )
                if result.returncode == 0 and screenshot_path.exists() and screenshot_path.stat().st_size > 0:
                    from PIL import Image

                    image = Image.open(screenshot_path)
                    image.crop((0, 0, PIXEL_W, PIXEL_H)).save(png_file)
                    return None
                last_error = RuntimeError(
                    f"{browser} exited with {result.returncode}: {result.stderr or result.stdout}"
                )
            except Exception as exc:
                last_error = exc

    return last_error or RuntimeError("Browser SVG rendering failed")


def render_svg_to_png(svg_file: Path, png_file: Path) -> None:
    cairo_error = render_with_cairosvg(svg_file, png_file)
    if cairo_error is None:
        return

    browser_error = render_with_browser(svg_file, png_file)
    if browser_error is None:
        return

    svglib_error = render_with_svglib_pdf(svg_file, png_file)
    if svglib_error is None:
        return

    raise RuntimeError(
        "Could not render SVG to PNG. "
        f"CairoSVG error: {cairo_error}. svglib/PDF error: {svglib_error}. Browser error: {browser_error}"
    )


def convert_raster(svg_files: list[Path], out_path: Path) -> None:
    prs = new_presentation()
    blank = prs.slide_layouts[6]

    with writable_temp_dir("svg_to_pptx_") as tmp_dir:
        for idx, svg_file in enumerate(svg_files, start=1):
            png_file = tmp_dir / f"slide_{idx:02d}.png"
            render_svg_to_png(svg_file, png_file)

            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(png_file),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

    prs.save(out_path)
