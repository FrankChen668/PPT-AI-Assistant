from __future__ import annotations

import base64
import io
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .pptx_style import (
    add_alpha,
    add_outer_shadow,
    effective_opacity,
    element_context,
    has_shadow,
    marker_end,
    number,
    pt_from_px,
    px,
    set_shape_fill,
    set_shape_line,
    svg_color,
)
from .raster import PIXEL_W, SLIDE_H_IN, SLIDE_W_IN
from .svg_utils import (
    extract_inheritable_styles,
    local_name,
    parse_points,
    path_points,
    resolve_attr,
    svg_attr,
    text_lines,
)

XLINK_HREF = "{http://www.w3.org/1999/xlink}href"


class ConversionStats(dict):
    """Tiny metrics bag for exporter observability."""

    def __init__(self) -> None:
        super().__init__(
            converted=0,
            skipped=0,
            unsupported=0,
            images=0,
        )


def _merge_inherited_styles(base: dict[str, str | float], overrides: dict[str, str]) -> dict[str, str | float]:
    merged = dict(base)
    # Multiply opacity channels when both present.
    for key in ("opacity", "fill-opacity", "stroke-opacity"):
        if key in overrides:
            existing = merged.get(key)
            try:
                if existing is not None:
                    merged[key] = str(float(existing) * float(overrides[key]))
                else:
                    merged[key] = overrides[key]
            except Exception:
                merged[key] = overrides[key]
    for k, v in overrides.items():
        if k in ("opacity", "fill-opacity", "stroke-opacity"):
            continue
        merged[k] = v
    return merged


def parse_transform(transform_str: str) -> tuple[float, float, float, float]:
    """Extract translate + scale from a simple SVG transform string."""
    if not transform_str:
        return 0.0, 0.0, 1.0, 1.0

    dx, dy = 0.0, 0.0
    sx, sy = 1.0, 1.0

    m = re.search(r"translate\(\s*([-\d.]+)(?:[\s,]+([-\d.]+))?\s*\)", transform_str)
    if m:
        dx = float(m.group(1))
        dy = float(m.group(2)) if m.group(2) is not None else 0.0

    m = re.search(r"scale\(\s*([-\d.]+)(?:[\s,]+([-\d.]+))?\s*\)", transform_str)
    if m:
        sx = float(m.group(1))
        sy = float(m.group(2)) if m.group(2) is not None else sx

    return dx, dy, sx, sy


def _xform_xy(x: float, y: float, inherited: dict[str, str | float]) -> tuple[float, float]:
    dx = float(inherited.get("_dx", 0.0))
    dy = float(inherited.get("_dy", 0.0))
    sx = float(inherited.get("_sx", 1.0))
    sy = float(inherited.get("_sy", 1.0))
    return x * sx + dx, y * sy + dy


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def add_rect(slide, elem: ET.Element, inherited: dict[str, str | float]) -> None:
    x = number(svg_attr(elem, "x"))
    y = number(svg_attr(elem, "y"))
    width = number(svg_attr(elem, "width"))
    height = number(svg_attr(elem, "height"))
    if width <= 0 or height <= 0:
        return
    x, y = _xform_xy(x, y, inherited)
    sx = float(inherited.get("_sx", 1.0))
    sy = float(inherited.get("_sy", 1.0))
    width *= sx
    height *= sy
    rx = number(svg_attr(elem, "rx"))
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, px(x), px(y), px(width), px(height))
    set_shape_fill(shape, resolve_attr(elem, inherited, "fill", "#000000"), effective_opacity(elem, inherited, "fill"))
    set_shape_line(
        shape,
        resolve_attr(elem, inherited, "stroke", "none"),
        resolve_attr(elem, inherited, "stroke-width", "1"),
        effective_opacity(elem, inherited, "stroke"),
    )
    if has_shadow(elem):
        add_outer_shadow(shape)


def add_circle(slide, elem: ET.Element, inherited: dict[str, str | float]) -> None:
    cx = number(svg_attr(elem, "cx"))
    cy = number(svg_attr(elem, "cy"))
    r = number(svg_attr(elem, "r"))
    if r <= 0:
        return
    cx, cy = _xform_xy(cx, cy, inherited)
    sx = float(inherited.get("_sx", 1.0))
    sy = float(inherited.get("_sy", 1.0))
    r *= (sx + sy) / 2.0
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - r), px(cy - r), px(2 * r), px(2 * r))
    set_shape_fill(shape, resolve_attr(elem, inherited, "fill", "#000000"), effective_opacity(elem, inherited, "fill"))
    set_shape_line(
        shape,
        resolve_attr(elem, inherited, "stroke", "none"),
        resolve_attr(elem, inherited, "stroke-width", "1"),
        effective_opacity(elem, inherited, "stroke"),
    )


def add_ellipse(slide, elem: ET.Element, inherited: dict[str, str | float]) -> None:
    cx = number(svg_attr(elem, "cx"))
    cy = number(svg_attr(elem, "cy"))
    rx = number(svg_attr(elem, "rx"))
    ry = number(svg_attr(elem, "ry"))
    if rx <= 0 or ry <= 0:
        return
    cx, cy = _xform_xy(cx, cy, inherited)
    sx = float(inherited.get("_sx", 1.0))
    sy = float(inherited.get("_sy", 1.0))
    rx *= sx
    ry *= sy
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - rx), px(cy - ry), px(2 * rx), px(2 * ry))
    set_shape_fill(shape, resolve_attr(elem, inherited, "fill", "#000000"), effective_opacity(elem, inherited, "fill"))
    set_shape_line(
        shape,
        resolve_attr(elem, inherited, "stroke", "none"),
        resolve_attr(elem, inherited, "stroke-width", "1"),
        effective_opacity(elem, inherited, "stroke"),
    )


def add_line(slide, elem: ET.Element, inherited: dict[str, str | float]) -> None:
    x1 = number(svg_attr(elem, "x1"))
    y1 = number(svg_attr(elem, "y1"))
    x2 = number(svg_attr(elem, "x2"))
    y2 = number(svg_attr(elem, "y2"))
    x1, y1 = _xform_xy(x1, y1, inherited)
    x2, y2 = _xform_xy(x2, y2, inherited)
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
    set_shape_line(
        shape,
        resolve_attr(elem, inherited, "stroke", "#000000"),
        resolve_attr(elem, inherited, "stroke-width", "1"),
        effective_opacity(elem, inherited, "stroke"),
        marker_end(elem),
    )


def add_poly_shape(
    slide,
    points: list[tuple[float, float]],
    closed: bool,
    elem: ET.Element,
    inherited: dict[str, str | float],
) -> None:
    if len(points) < 2:
        return

    fill = resolve_attr(elem, inherited, "fill", "none")
    if closed and svg_color(fill) is not None:
        points = [_xform_xy(x, y, inherited) for x, y in points]
        builder = slide.shapes.build_freeform(px(points[0][0]), px(points[0][1]))
        builder.add_line_segments([(px(x), px(y)) for x, y in points[1:]], close=True)
        shape = builder.convert_to_shape()
        set_shape_fill(shape, fill, effective_opacity(elem, inherited, "fill"))
        set_shape_line(
            shape,
            resolve_attr(elem, inherited, "stroke", "none"),
            resolve_attr(elem, inherited, "stroke-width", "1"),
            effective_opacity(elem, inherited, "stroke"),
        )
        return

    for start, end in zip(points, points[1:]):
        start = _xform_xy(start[0], start[1], inherited)
        end = _xform_xy(end[0], end[1], inherited)
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px(start[0]),
            px(start[1]),
            px(end[0]),
            px(end[1]),
        )
        set_shape_line(
            shape,
            resolve_attr(elem, inherited, "stroke", "#000000"),
            resolve_attr(elem, inherited, "stroke-width", "1"),
            effective_opacity(elem, inherited, "stroke"),
            marker_end(elem) and end == points[-1],
        )
    if closed:
        start = _xform_xy(points[-1][0], points[-1][1], inherited)
        end = _xform_xy(points[0][0], points[0][1], inherited)
        shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(start[0]), px(start[1]), px(end[0]), px(end[1]))
        set_shape_line(
            shape,
            resolve_attr(elem, inherited, "stroke", "#000000"),
            resolve_attr(elem, inherited, "stroke-width", "1"),
            effective_opacity(elem, inherited, "stroke"),
        )


def add_path(slide, elem: ET.Element, inherited: dict[str, str | float]) -> None:
    points, closed = path_points(svg_attr(elem, "d"))
    add_poly_shape(slide, points, closed, elem, inherited)


def add_polyline(slide, elem: ET.Element, inherited: dict[str, str | float], closed: bool) -> None:
    add_poly_shape(slide, parse_points(svg_attr(elem, "points")), closed, elem, inherited)


def text_box_geometry(
    elem: ET.Element,
    inherited: dict[str, str | float],
    values: list[str],
    font_size: float,
) -> tuple[float, float, float, float, int]:
    x = number(svg_attr(elem, "x"))
    y = number(svg_attr(elem, "y"))
    x, y = _xform_xy(x, y, inherited)
    anchor = resolve_attr(elem, inherited, "text-anchor", "start")
    max_chars = max((len(value) for value in values), default=1)
    estimated_width = max(100, min(1160, max_chars * font_size * 0.78 + 48))

    if anchor == "middle":
        left = x - estimated_width / 2
        align = PP_ALIGN.CENTER
    elif anchor == "end":
        left = x - estimated_width
        align = PP_ALIGN.RIGHT
    else:
        left = x
        align = PP_ALIGN.LEFT

    left = max(0, min(left, PIXEL_W - estimated_width))
    top = max(0, y - font_size * 0.92)
    height = max(font_size * 1.35, len(values) * font_size * 1.35 + 8)
    return left, top, estimated_width, height, align


def add_text(slide, elem: ET.Element, inherited: dict[str, str | float]) -> None:
    values = text_lines(elem)
    if not values:
        return

    font_size = number(resolve_attr(elem, inherited, "font-size"), 18)
    left, top, width, height, align = text_box_geometry(elem, inherited, values, font_size)
    shape = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = False

    color = svg_color(resolve_attr(elem, inherited, "fill", "#000000"))
    weight = number(resolve_attr(elem, inherited, "font-weight"), 400)
    family = (resolve_attr(elem, inherited, "font-family") or "Arial").split(",")[0].strip().strip('"')
    opacity = effective_opacity(elem, inherited, "fill")

    for idx, value in enumerate(values):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = value
        run.font.name = family
        run.font.size = pt_from_px(font_size)
        run.font.bold = weight >= 600
        if color is not None:
            run.font.color.rgb = color
            solid_fill = run._r.find(f".//{qn('a:solidFill')}")
            if solid_fill is not None:
                add_alpha(solid_fill, opacity)


def _image_href(elem: ET.Element) -> str | None:
    return elem.get("href") or elem.get(XLINK_HREF)


def add_image(slide, elem: ET.Element, inherited: dict[str, str | float], svg_dir: Path) -> bool:
    href = _image_href(elem)
    if not href:
        return False

    x = number(svg_attr(elem, "x"))
    y = number(svg_attr(elem, "y"))
    w = number(svg_attr(elem, "width"))
    h = number(svg_attr(elem, "height"))
    if w <= 0 or h <= 0:
        return False
    x, y = _xform_xy(x, y, inherited)
    sx = float(inherited.get("_sx", 1.0))
    sy = float(inherited.get("_sy", 1.0))
    w *= sx
    h *= sy

    href = href.strip()
    if href.startswith("data:image/"):
        m = re.match(r"data:image/(\w+);base64,(.+)", href, re.DOTALL)
        if not m:
            return False
        raw = base64.b64decode(m.group(2).encode("ascii"), validate=False)
        stream = io.BytesIO(raw)
        slide.shapes.add_picture(stream, px(x), px(y), px(w), px(h))
        return True

    image_path = (svg_dir / href).resolve() if not Path(href).is_absolute() else Path(href)
    if not image_path.exists():
        return False
    slide.shapes.add_picture(str(image_path), px(x), px(y), px(w), px(h))
    return True


def add_svg_element(
    slide,
    elem: ET.Element,
    svg_dir: Path,
    stats: ConversionStats,
    inherited: dict[str, str | float] | None = None,
) -> None:
    inherited = inherited or {"_opacity": 1.0, "_dx": 0.0, "_dy": 0.0, "_sx": 1.0, "_sy": 1.0}
    context = element_context(elem, inherited)
    tag = local_name(elem.tag)
    if tag in {"defs", "title"}:
        stats["skipped"] += 1
        return
    if tag == "g":
        dx, dy, sx, sy = parse_transform(elem.get("transform", ""))
        child_ctx = dict(context)
        child_ctx = _merge_inherited_styles(child_ctx, extract_inheritable_styles(elem))
        child_ctx["_dx"] = float(child_ctx.get("_dx", 0.0)) + dx * float(child_ctx.get("_sx", 1.0))
        child_ctx["_dy"] = float(child_ctx.get("_dy", 0.0)) + dy * float(child_ctx.get("_sy", 1.0))
        child_ctx["_sx"] = float(child_ctx.get("_sx", 1.0)) * sx
        child_ctx["_sy"] = float(child_ctx.get("_sy", 1.0)) * sy
        for child in elem:
            add_svg_element(slide, child, svg_dir, stats, child_ctx)
        stats["converted"] += 1
    elif tag == "rect":
        add_rect(slide, elem, context)
        stats["converted"] += 1
    elif tag == "circle":
        add_circle(slide, elem, context)
        stats["converted"] += 1
    elif tag == "ellipse":
        add_ellipse(slide, elem, context)
        stats["converted"] += 1
    elif tag == "line":
        add_line(slide, elem, context)
        stats["converted"] += 1
    elif tag == "polyline":
        add_polyline(slide, elem, context, closed=False)
        stats["converted"] += 1
    elif tag == "polygon":
        add_polyline(slide, elem, context, closed=True)
        stats["converted"] += 1
    elif tag == "path":
        add_path(slide, elem, context)
        stats["converted"] += 1
    elif tag == "text":
        add_text(slide, elem, context)
        stats["converted"] += 1
    elif tag == "image":
        ok = add_image(slide, elem, context, svg_dir)
        if ok:
            stats["images"] += 1
            stats["converted"] += 1
        else:
            stats["unsupported"] += 1
    else:
        stats["unsupported"] += 1


def add_svg_native(slide, svg_file: Path, stats: ConversionStats) -> None:
    root = ET.parse(svg_file).getroot()
    if local_name(root.tag) != "svg":
        raise ValueError(f"{svg_file.name} root element is not <svg>")
    for child in root:
        add_svg_element(slide, child, svg_file.parent, stats)


def convert_native(svg_files: list[Path], out_path: Path) -> ConversionStats:
    prs = new_presentation()
    blank = prs.slide_layouts[6]
    stats = ConversionStats()
    for svg_file in svg_files:
        slide = prs.slides.add_slide(blank)
        add_svg_native(slide, svg_file, stats)
    prs.save(out_path)
    return stats

