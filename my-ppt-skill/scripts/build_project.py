#!/usr/bin/env python3
"""Run the standard AI-PPT project pipeline end to end."""

from __future__ import annotations

import argparse
import hashlib
import json
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from time import perf_counter
from typing import Any

SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPT_DIR))

from auto_copyfit_svg import CopyfitStats, run_auto_copyfit  # noqa: E402
from build_config.fallback_policy import load_fallback_policy  # noqa: E402
from check_svg_encoding import format_issues, validate_project_svg_output  # noqa: E402
from copyfit_contract import ContractReport, validate_copyfit_contract  # noqa: E402
from doctor_export import DoctorReport, run_export_doctor  # noqa: E402
from export.error_utils import (  # noqa: E402
    classify_export_error,
    fallback_action_hint,
    format_pptx_write_error,
    looks_like_locked_pptx_error,
)
from export.name_utils import (  # noqa: E402
    build_output_plan,
    is_stable_target,
    next_semantic_target,
    planned_semantic_target,
)
from finalize_svg import finalize_project  # noqa: E402
from generate_slide_plan import (  # noqa: E402
    BudgetReport,
    SlidePlanStats,
    ensure_semantic_slide_plan,
    evaluate_budget_policy,
    generate_slide_plan,
)
from lint_svg_layout import LayoutLintReport, run_layout_lint  # noqa: E402
from normalize_images import NormalizeStats, normalize_project_images_white_bg  # noqa: E402
from optimize_delivery_svg import DeliveryOptimizeStats, run_delivery_optimize  # noqa: E402
from pipeline import context_observability as context_obs  # noqa: E402
from pipeline.build_phases import handle_authoring_phase, handle_finalize_skip_qa_phase  # noqa: E402
from pipeline.build_types import NativeTextVerifyReport  # noqa: E402
from pipeline.stages import resolve_phase_options  # noqa: E402
from pipeline.style_draft_gate import ensure_style_draft_selection, validate_style_draft_selection  # noqa: E402
from planning_quality import evaluate_planning_quality  # noqa: E402
from preflight_svg import PreflightReport, run_preflight  # noqa: E402
from qa_project import QaReport, run_qa  # noqa: E402
from quality.checks.repair import (  # noqa: E402
    append_repair_audit,
    apply_deterministic_repair,
    collect_repair_slide_ids,
    write_repair_round_note,
)
from quality.checks.reporting import extract_quality_tiers  # noqa: E402
from render_svg import render_project  # noqa: E402
from render_theme import visual_width  # noqa: E402
from svg_to_pptx import (  # noqa: E402
    convert,
    load_native_conversion_report,
    native_conversion_report_path,
)

INCREMENTAL_CACHE_VERSION = 1
PAGE_SUMMARY_CACHE_VERSION = 1
SESSION_CHECKPOINT_VERSION = 1
TOKEN_BUDGET_POLICY = "default-v1"
TOKEN_BUDGET_CAPS = {
    "strategist": 3000,
    "designer": 4000,
    "executor_per_slide": 2500,
    "checks": 8000,
}
SESSION_CHECKPOINT_DEFAULT_PAGES = 4
SESSION_CHECKPOINT_MIN_PAGES = 3
SESSION_CHECKPOINT_MAX_PAGES = 5


def _project_relative_path(project_dir: Path, path: Path) -> str:
    try:
        return str(path.relative_to(project_dir))
    except ValueError:
        return str(path)


def _blueprint_slide_count(project_dir: Path) -> int:
    blueprint = project_dir / "blueprint.json"
    if not blueprint.exists():
        return 0
    try:
        data = json.loads(blueprint.read_text(encoding="utf-8-sig"))
    except Exception:
        return 0
    slides = data.get("slides")
    return len(slides) if isinstance(slides, list) else 0


def _blueprint_hash(project_dir: Path) -> str | None:
    blueprint = project_dir / "blueprint.json"
    if not blueprint.exists():
        return None
    return hashlib.sha256(blueprint.read_bytes()).hexdigest()


def _git_sha(project_dir: Path) -> str | None:
    result = subprocess.run(
        ["git", "-C", str(project_dir), "rev-parse", "HEAD"],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    if result.returncode != 0:
        return None
    return result.stdout.strip() or None


def _template_key(project_dir: Path) -> str | None:
    binding = project_dir / "template_binding.json"
    if not binding.exists():
        return None
    try:
        data = json.loads(binding.read_text(encoding="utf-8"))
    except Exception:
        return None
    for key in ("template_key", "layout_id", "template_id"):
        value = data.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    return None


def _safe_int(value: Any, default: int = 0) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _safe_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _template_lookup_metrics(project_dir: Path) -> dict[str, Any]:
    defaults = {
        "template_lookup_mode": "unknown",
        "template_reference_files_loaded": 0,
        "template_reference_files_skipped": 0,
        "template_lazy_load_hit_ratio": 0.0,
    }
    path = project_dir / "reference_pack.json"
    if not path.exists():
        return defaults
    try:
        payload = json.loads(path.read_text(encoding="utf-8-sig"))
    except Exception:
        return defaults
    if not isinstance(payload, dict):
        return defaults
    loaded_default = _safe_int(payload.get("template_reference_files_loaded", 0), 0)
    mode = payload.get("template_lookup_mode")
    return {
        "template_lookup_mode": str(mode).strip() if isinstance(mode, str) and mode.strip() else "unknown",
        "template_reference_files_loaded": loaded_default,
        "template_reference_files_skipped": _safe_int(payload.get("template_reference_files_skipped", 0), 0),
        "template_lazy_load_hit_ratio": _safe_float(
            payload.get("template_lazy_load_hit_ratio", 1.0 if loaded_default > 0 else 0.0),
            1.0 if loaded_default > 0 else 0.0,
        ),
    }


def _sha256_file(path: Path) -> str | None:
    if not path.exists() or not path.is_file():
        return None
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _collect_incremental_source_hashes(project_dir: Path, changed_slides: list[int] | None = None) -> dict[str, Any]:
    files = {
        "design_spec_hash": project_dir / "design_spec.md",
        "blueprint_hash": project_dir / "blueprint.json",
        "style_route_hash": project_dir / "style_route.json",
        "art_direction_hash": project_dir / "art_direction.md",
        "reference_pack_hash": project_dir / "reference_pack.json",
        "slide_visual_plan_hash": project_dir / "slide_visual_plan.json",
        "slide_plan_hash": project_dir / "slide_plan.json",
    }
    payload: dict[str, Any] = {}
    for key, path in files.items():
        file_hash = _sha256_file(path)
        if file_hash:
            payload[key] = file_hash

    svg_dir = project_dir / "svg_output"
    if changed_slides:
        slide_hashes: dict[str, str] = {}
        for slide_id in sorted(set(changed_slides)):
            slide_file = svg_dir / f"slide_{slide_id:02d}.svg"
            slide_hash = _sha256_file(slide_file)
            if slide_hash:
                slide_hashes[str(slide_id)] = slide_hash
        payload["changed_slide_hashes"] = slide_hashes
    else:
        digest = hashlib.sha256()
        for svg_file in sorted(svg_dir.glob("slide_*.svg")):
            digest.update(svg_file.name.encode("utf-8"))
            digest.update((hashlib.sha256(svg_file.read_bytes()).hexdigest()).encode("utf-8"))
        payload["svg_output_digest"] = digest.hexdigest()
    return payload


def _incremental_context_profile(project_dir: Path, changed_slides: list[int] | None = None) -> dict[str, Any]:
    files: list[Path] = [
        project_dir / "design_spec.md",
        project_dir / "blueprint.json",
    ]
    optional = [
        project_dir / "style_route.json",
        project_dir / "slide_visual_plan.json",
        project_dir / "art_direction.md",
    ]
    files.extend(path for path in optional if path.exists())
    if changed_slides:
        files.extend(
            (project_dir / "svg_output" / f"slide_{slide_id:02d}.svg")
            for slide_id in sorted(set(changed_slides))
        )
    total_bytes = 0
    existing_files: list[str] = []
    for path in files:
        if not path.exists() or not path.is_file():
            continue
        total_bytes += path.stat().st_size
        existing_files.append(_project_relative_path(project_dir, path))
    return {
        "context_file_count": len(existing_files),
        "context_bytes_estimate": total_bytes,
        "context_files": existing_files,
    }


def _estimate_tokens_from_bytes(size_bytes: int) -> int:
    return max(0, int(round(size_bytes / 4.0)))


def _resolve_token_budget_stage(phase: str, changed_slides: list[int]) -> str:
    if phase == "authoring":
        return "executor_per_slide" if len(changed_slides) == 1 else "designer"
    return "checks"


def _evaluate_token_budget(
    *,
    phase: str,
    changed_slides: list[int],
    context_profile: dict[str, Any],
) -> dict[str, Any]:
    stage = _resolve_token_budget_stage(phase, changed_slides)
    limit = int(TOKEN_BUDGET_CAPS.get(stage, TOKEN_BUDGET_CAPS["checks"]))
    estimate = _estimate_tokens_from_bytes(int(context_profile.get("context_bytes_estimate", 0)))
    overflow = max(0, estimate - limit)
    return {
        "token_budget_policy": TOKEN_BUDGET_POLICY,
        "token_budget_stage": stage,
        "token_budget_limit": limit,
        "context_token_estimate": estimate,
        "token_budget_overflow": overflow,
        "token_budget_warning": overflow > 0,
    }


def _incremental_cache_path(project_dir: Path) -> Path:
    return project_dir / "exports" / "build_cache.json"


def _read_incremental_cache(project_dir: Path) -> dict[str, Any]:
    cache_path = _incremental_cache_path(project_dir)
    if not cache_path.exists():
        return {"version": INCREMENTAL_CACHE_VERSION, "history": []}
    try:
        payload = json.loads(cache_path.read_text(encoding="utf-8"))
    except Exception:
        return {"version": INCREMENTAL_CACHE_VERSION, "history": []}
    if not isinstance(payload, dict):
        return {"version": INCREMENTAL_CACHE_VERSION, "history": []}
    payload.setdefault("version", INCREMENTAL_CACHE_VERSION)
    payload.setdefault("history", [])
    return payload


def _write_incremental_cache(project_dir: Path, last_success: dict[str, Any]) -> None:
    cache_path = _incremental_cache_path(project_dir)
    cache_path.parent.mkdir(parents=True, exist_ok=True)
    payload = _read_incremental_cache(project_dir)
    history = payload.get("history")
    if not isinstance(history, list):
        history = []
    history.append(
        {
            "timestamp": datetime.now().isoformat(timespec="seconds"),
            "phase": last_success.get("phase"),
            "source_hashes": last_success.get("source_hashes", {}),
            "output_files": last_success.get("output_files", []),
            "qa_ok": bool(last_success.get("qa_ok", False)),
        }
    )
    payload["history"] = history[-40:]
    payload["version"] = INCREMENTAL_CACHE_VERSION
    payload["updated_at"] = datetime.now().isoformat(timespec="seconds")
    payload["last_success"] = last_success
    cache_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _incremental_cache_hit(
    project_dir: Path,
    *,
    phase: str,
    source_hashes: dict[str, Any],
) -> tuple[bool, str, dict[str, Any] | None]:
    payload = _read_incremental_cache(project_dir)
    last = payload.get("last_success")
    if not isinstance(last, dict):
        return False, "cache-miss:no-last-success", None
    if str(last.get("phase")) != str(phase):
        return False, "cache-miss:phase-changed", last
    if last.get("source_hashes") != source_hashes:
        return False, "cache-miss:source-hash-changed", last
    output_files = last.get("output_files", [])
    if isinstance(output_files, list):
        for rel in output_files:
            if not isinstance(rel, str):
                continue
            if not (project_dir / rel).exists():
                return False, "cache-miss:output-missing", last
    return True, "cache-hit:source-unchanged", last


def _page_summary_cache_path(project_dir: Path) -> Path:
    return project_dir / "exports" / "page_summary_cache.json"


def _read_page_summary_cache(project_dir: Path) -> dict[str, Any]:
    path = _page_summary_cache_path(project_dir)
    if not path.exists():
        return {"version": PAGE_SUMMARY_CACHE_VERSION, "entries": {}}
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {"version": PAGE_SUMMARY_CACHE_VERSION, "entries": {}}
    if not isinstance(payload, dict):
        return {"version": PAGE_SUMMARY_CACHE_VERSION, "entries": {}}
    entries = payload.get("entries")
    if not isinstance(entries, dict):
        entries = {}
    payload["version"] = PAGE_SUMMARY_CACHE_VERSION
    payload["entries"] = entries
    return payload


def _slide_plan_by_id(project_dir: Path) -> dict[int, dict[str, Any]]:
    plan_path = project_dir / "slide_visual_plan.json"
    if not plan_path.exists():
        return {}
    try:
        payload = json.loads(plan_path.read_text(encoding="utf-8-sig"))
    except Exception:
        return {}
    slides = payload.get("slides") if isinstance(payload, dict) else None
    if not isinstance(slides, list):
        return {}
    by_id: dict[int, dict[str, Any]] = {}
    for item in slides:
        if not isinstance(item, dict):
            continue
        raw_id = item.get("slide_id", item.get("id"))
        if not isinstance(raw_id, int) or raw_id <= 0:
            continue
        by_id[raw_id] = item
    return by_id


def _hash_json(value: Any) -> str:
    canonical = json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(canonical.encode("utf-8")).hexdigest()


def _slide_svg_hash(project_dir: Path, slide_id: int) -> str | None:
    for folder in ("svg_output", "svg_final"):
        path = project_dir / folder / f"slide_{slide_id:02d}.svg"
        digest = _sha256_file(path)
        if digest:
            return digest
    return None




def _style_route_payload(project_dir: Path) -> dict[str, Any] | None:
    path = project_dir / "style_route.json"
    if not path.exists():
        return None
    try:
        payload = json.loads(path.read_text(encoding="utf-8-sig"))
    except Exception:
        return None
    return payload if isinstance(payload, dict) else None


def _layout_exploration_assessment(project_dir: Path) -> dict[str, Any]:
    enabled = False
    candidate_count = 2
    archetype_switch_count = 0
    consecutive_repeat_count = 0
    diversity_gate_result = "unknown"

    route_payload = _style_route_payload(project_dir)
    if isinstance(route_payload, dict):
        route_exploration = route_payload.get("layout_exploration")
        if isinstance(route_exploration, dict):
            enabled = bool(route_exploration.get("enabled", enabled))
            try:
                candidate_count = int(route_exploration.get("candidate_count", candidate_count))
            except (TypeError, ValueError):
                candidate_count = 2

    plan_path = project_dir / "slide_visual_plan.json"
    if plan_path.exists():
        try:
            plan_payload = json.loads(plan_path.read_text(encoding="utf-8-sig"))
        except Exception:
            plan_payload = {}
        if isinstance(plan_payload, dict):
            plan_exploration = plan_payload.get("layout_exploration")
            if isinstance(plan_exploration, dict):
                enabled = bool(plan_exploration.get("enabled", enabled))
                try:
                    candidate_count = int(plan_exploration.get("candidate_count", candidate_count))
                except (TypeError, ValueError):
                    candidate_count = 2
            slides = plan_payload.get("slides")
            if isinstance(slides, list):
                selected = [
                    str(item.get("selected_archetype") or item.get("visual_archetype") or "").strip()
                    for item in slides
                    if isinstance(item, dict)
                ]
                selected = [item for item in selected if item]
                if selected:
                    diversity_gate_result = "pass"
                previous = ""
                for current in selected:
                    if previous:
                        if current == previous:
                            consecutive_repeat_count += 1
                            diversity_gate_result = "warn"
                        else:
                            archetype_switch_count += 1
                    previous = current

    return {
        "layout_exploration_enabled": enabled,
        "candidate_count": candidate_count,
        "archetype_switch_count": archetype_switch_count,
        "consecutive_repeat_count": consecutive_repeat_count,
        "diversity_gate_result": diversity_gate_result,
    }




def _append_manifest(project_dir: Path, record: dict[str, Any]) -> None:
    exports_dir = project_dir / "exports"
    exports_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = exports_dir / "manifest.json"
    payload: dict[str, Any] = {"records": []}
    if manifest_path.exists():
        try:
            loaded = json.loads(manifest_path.read_text(encoding="utf-8"))
            if isinstance(loaded, dict) and isinstance(loaded.get("records"), list):
                payload = loaded
            elif isinstance(loaded, list):
                payload = {"records": loaded}
        except Exception:
            payload = {"records": []}
    payload["records"].append(record)
    manifest_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _update_manifest_last_record(project_dir: Path, updates: dict[str, Any]) -> None:
    manifest_path = project_dir / "exports" / "manifest.json"
    if not manifest_path.exists():
        return
    try:
        payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    except Exception:
        return
    records = payload.get("records")
    if not isinstance(records, list) or not records:
        return
    tail = records[-1]
    if not isinstance(tail, dict):
        return
    tail.update(updates)
    manifest_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _native_conversion_manifest_fields(report_path: Path, report: dict[str, Any]) -> dict[str, Any]:
    summary_value = report.get("summary")
    summary: dict[str, Any] = summary_value if isinstance(summary_value, dict) else {}
    fields: dict[str, Any] = {
        "native_conversion_report": str(report_path).replace("\\", "/"),
        "native_conversion_status": str(report.get("delivery_status") or "unknown"),
        "native_conversion_delivery_approved": report.get("delivery_approved") is True,
        "native_conversion_manual_review_required": report.get("manual_review_required") is True,
    }
    manifest_names = {
        "source_slide_count": "native_conversion_source_slide_count",
        "pptx_slide_count": "native_conversion_pptx_slide_count",
        "page_count_match": "native_conversion_page_count_match",
        "visible_elements": "native_conversion_visible_element_count",
        "converted_elements": "native_conversion_converted_element_count",
        "conversion_errors": "native_conversion_error_count",
        "unsupported_elements": "native_conversion_unsupported_element_count",
        "source_text_nodes": "native_conversion_source_text_node_count",
        "converted_text_nodes": "native_conversion_converted_text_node_count",
        "images_succeeded": "native_conversion_image_success_count",
        "images_failed": "native_conversion_image_failure_count",
        "critical_elements": "native_conversion_critical_element_count",
        "converted_critical_elements": "native_conversion_converted_critical_element_count",
        "critical_text_nodes": "native_conversion_critical_text_node_count",
        "converted_critical_text_nodes": "native_conversion_converted_critical_text_node_count",
        "critical_images": "native_conversion_critical_image_count",
        "converted_critical_images": "native_conversion_converted_critical_image_count",
        "markers_converted": "native_conversion_marker_converted_count",
        "unsupported_markers": "native_conversion_unsupported_marker_count",
        "blank_slides": "native_conversion_blank_slide_count",
        "hard_blocker_count": "native_conversion_hard_blocker_count",
        "quality_note_count": "native_conversion_quality_note_count",
    }
    for source_key, manifest_key in manifest_names.items():
        fields[manifest_key] = summary.get(source_key)
    fields["native_conversion_unsupported_tags"] = list(summary.get("unsupported_tags") or [])
    return fields


def _planning_manifest_fields(report_path: Path, report: dict[str, Any]) -> dict[str, Any]:
    summary_value = report.get("summary")
    summary: dict[str, Any] = summary_value if isinstance(summary_value, dict) else {}
    return {
        "planning_status": str(report.get("planning_status") or "unknown"),
        "planning_report": str(report_path).replace("\\", "/"),
        "planning_error_count": int(summary.get("planning_error_count") or 0),
        "planning_warning_count": int(summary.get("planning_warning_count") or 0),
        "planning_note_count": int(summary.get("planning_note_count") or 0),
        "slide_type_distribution": dict(summary.get("slide_type_distribution") or {}),
        "visual_archetype_distribution": dict(summary.get("visual_archetype_distribution") or {}),
        "repeated_title_count": int(summary.get("repeated_title_count") or 0),
        "repeated_conclusion_count": int(summary.get("repeated_conclusion_count") or 0),
        "content_overload_slide_count": int(summary.get("content_overload_slide_count") or 0),
        "planning_manual_review_required": report.get("manual_review_required") is True,
        "planning_delivery_approved": report.get("delivery_approved") is True,
    }


def _delivery_status_tuple(
    *,
    artifact_created: bool,
    qa_report: QaReport | None,
    native_conversion_report: dict[str, Any] | None = None,
    planning_report: dict[str, Any] | None = None,
) -> tuple[bool, str, list[str]]:
    if not artifact_created:
        return False, "artifact_failed", ["no_artifact_output"]
    if qa_report is None:
        return False, "artifact_only", ["qa_not_executed"]
    native_summary_value = native_conversion_report.get("summary") if isinstance(native_conversion_report, dict) else None
    native_summary: dict[str, Any] = native_summary_value if isinstance(native_summary_value, dict) else {}
    native_hard_blockers = int(native_summary.get("hard_blocker_count") or 0)
    native_quality_notes = int(native_summary.get("quality_note_count") or 0)
    planning_summary_value = planning_report.get("summary") if isinstance(planning_report, dict) else None
    planning_summary: dict[str, Any] = planning_summary_value if isinstance(planning_summary_value, dict) else {}
    planning_hard_blockers = int(planning_summary.get("planning_error_count") or 0)
    planning_quality_notes = int(planning_summary.get("planning_note_count") or 0)
    if planning_hard_blockers > 0:
        return False, "blocked", [f"planning_hard_blockers:{planning_hard_blockers}"]
    if native_hard_blockers > 0:
        return False, "blocked", [f"native_conversion_hard_blockers:{native_hard_blockers}"]
    if qa_report.ok:
        if native_quality_notes > 0 or planning_quality_notes > 0:
            quality_note_reasons: list[str] = []
            if native_quality_notes > 0:
                quality_note_reasons.append(f"native_conversion_quality_notes:{native_quality_notes}")
            if planning_quality_notes > 0:
                quality_note_reasons.append(f"planning_quality_notes:{planning_quality_notes}")
            return True, "downloadable_with_notes", quality_note_reasons
        if _qa_manual_review_required(qa_report):
            return False, "artifact_created_but_visual_review_required", ["visual_review_required"]
        return True, "approved", []
    reasons: list[str] = []
    if int(qa_report.errors) > 0:
        reasons.append(f"qa_errors:{int(qa_report.errors)}")
    if int(qa_report.warnings) > 0:
        reasons.append(f"qa_warnings:{int(qa_report.warnings)}")
    if not reasons:
        reasons.append("qa_failed")
    return False, "artifact_created_but_qa_failed", reasons


def _qa_manual_review_required(qa_report: QaReport) -> bool:
    visual_score = qa_report.visual_score
    if visual_score is None:
        return False
    blocking_visual_warnings = int(qa_report.metrics.get("blocking_visual_warning_count", 0))
    high_score_density_advisory_only = bool(
        qa_report.density_flag
        and int(qa_report.warnings) == 0
        and visual_score >= 95
    )
    visual_ready = bool(
        visual_score >= 85
        and (not qa_report.density_flag or high_score_density_advisory_only)
        and not qa_report.hierarchy_flag
        and blocking_visual_warnings == 0
    )
    return not visual_ready


def _print_delivery_status(
    *,
    artifact_created: bool,
    delivery_approved: bool,
    delivery_status: str,
    delivery_failure_reasons: list[str],
    qa_report_path: str,
    qa_report: QaReport | None = None,
) -> None:
    reason_text = ", ".join(delivery_failure_reasons) if delivery_failure_reasons else "none"
    print(
        f"Artifact created: {str(bool(artifact_created)).lower()}; "
        f"Delivery approved: {str(bool(delivery_approved)).lower()}; "
        f"Delivery status: {delivery_status}; "
        f"QA report: {qa_report_path}; "
        f"Failure reasons: {reason_text}"
    )
    if qa_report is None:
        print("Readiness status: qa_not_available; visual readiness was not evaluated.")
        return
    blocking_visual_warnings = int(qa_report.metrics.get("blocking_visual_warning_count", 0))
    visual_ready = bool(
        qa_report.visual_score is not None
        and qa_report.visual_score >= 85
        and (
            not qa_report.density_flag
            or (
                int(qa_report.warnings) == 0
                and qa_report.visual_score >= 95
            )
        )
        and not qa_report.hierarchy_flag
        and blocking_visual_warnings == 0
    )
    manual_review_required = _qa_manual_review_required(qa_report)
    readiness_status = (
        "ready"
        if delivery_approved and visual_ready and not manual_review_required
        else "visual_review_required"
    )
    if not delivery_approved and delivery_status != "artifact_created_but_visual_review_required":
        readiness_status = "blocked"
    meaning = (
        "Engineering export passed, but visual delivery still requires review or polish."
        if readiness_status == "visual_review_required"
        else "Engineering and visual delivery gates are both ready."
        if readiness_status == "ready"
        else "Delivery has blocking QA findings or export failures."
    )
    print(
        f"Readiness status: {readiness_status}; "
        f"Visual ready: {str(visual_ready).lower()}; "
        f"Manual review required: {str(manual_review_required).lower()}; "
        f"Visual score: {qa_report.visual_score}; "
        f"Warnings: {qa_report.warnings}; "
        f"Advisories: {qa_report.advisories}; "
        f"Meaning: {meaning}"
    )


def _raise_on_invalid_svg_encoding(project_dir: Path) -> None:
    issues = validate_project_svg_output(project_dir)
    if not issues:
        return
    raise RuntimeError(format_issues(project_dir, issues))


def _contains_cjk(text: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def _count_cjk(text: str) -> int:
    return sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff")


def _verify_native_text_layout(
    pptx_paths: list[Path],
    *,
    narrow_width_pt: float = 300.0,
    risk_ratio: float = 1.06,
    min_cjk_chars: int = 10,
) -> NativeTextVerifyReport:
    if not pptx_paths:
        return NativeTextVerifyReport(True, 0, 0, 0, 0, [])
    try:
        from pptx import Presentation
    except Exception:
        return NativeTextVerifyReport(True, len(pptx_paths), 0, 0, 0, [])

    checked_lines = 0
    narrow_lines = 0
    risk_lines = 0
    findings: list[dict[str, Any]] = []
    for pptx_path in pptx_paths:
        if not pptx_path.exists():
            continue
        prs = Presentation(pptx_path)
        for slide_index, slide in enumerate(prs.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if not hasattr(shape, "text_frame") or shape.text_frame is None:
                    continue
                frame = shape.text_frame
                frame_width_pt = float(getattr(shape, "width", 0)) / 12700.0
                if frame_width_pt <= 0:
                    continue
                lines: list[tuple[str, float]] = []
                for paragraph in frame.paragraphs:
                    text = (paragraph.text or "").strip()
                    if not text:
                        continue
                    font_pt = 18.0
                    if paragraph.runs:
                        run_size = paragraph.runs[0].font.size
                        if run_size is not None and hasattr(run_size, "pt"):
                            font_pt = float(run_size.pt)
                    para_size = paragraph.font.size
                    if para_size is not None and hasattr(para_size, "pt"):
                        font_pt = float(para_size.pt)
                    for line in text.splitlines():
                        line_text = line.strip()
                        if not line_text:
                            continue
                        lines.append((line_text, font_pt))
                for line_text, font_pt in lines:
                    checked_lines += 1
                    if frame_width_pt > narrow_width_pt:
                        continue
                    if not _contains_cjk(line_text):
                        continue
                    if _count_cjk(line_text) < min_cjk_chars:
                        continue
                    narrow_lines += 1
                    estimated_width_pt = visual_width(line_text) * font_pt * 0.92
                    if estimated_width_pt > frame_width_pt * risk_ratio:
                        risk_lines += 1
                        findings.append(
                            {
                                "pptx": str(pptx_path),
                                "slide": slide_index,
                                "shape": shape_index,
                                "estimated_width_pt": round(estimated_width_pt, 2),
                                "frame_width_pt": round(frame_width_pt, 2),
                                "text_preview": line_text[:48],
                            }
                        )
    return NativeTextVerifyReport(
        ok=risk_lines == 0,
        checked_pptx=len(pptx_paths),
        checked_lines=checked_lines,
        narrow_column_lines=narrow_lines,
        risk_lines=risk_lines,
        findings=findings,
    )


def _evaluate_budget_policy_if_available(
    project_dir: Path,
    profile: str,
    *,
    svg_only_export: bool,
) -> BudgetReport | None:
    """Run budget policy when blueprint exists; allow explicit SVG-only export bypass."""
    blueprint_path = project_dir / "blueprint.json"
    if not blueprint_path.exists():
        if not svg_only_export:
            raise FileNotFoundError(
                f"Missing blueprint.json: {blueprint_path}. "
                "Use --svg-only-export only for explicit compatibility/export paths with pre-authored SVG."
            )
        print(
            "Budget policy skipped: missing blueprint.json. "
            "Proceeding because --svg-only-export is enabled for explicit SVG-only export paths.",
            file=sys.stderr,
        )
        return None
    return evaluate_budget_policy(project_dir, profile=profile)


def build_project(
    project_dir: Path,
    phase: str = "finalize",
    mode: str = "native",
    artifact_name: str = "stable",
    clean: bool = False,
    snapshots: bool = False,
    strict: bool = False,
    skip_render: bool = False,
    force_render: bool = False,
    auto_slide_plan: bool = False,
    auto_slide_plan_overwrite: bool = False,
    normalize_images_white_bg: bool = False,
    guard_profile: str = "balanced",
    auto_copyfit: bool = False,
    delivery_ready: bool = False,
    enable_layout_lint: bool = False,
    layout_lint_strict: bool = False,
    layout_quality_profile: str = "standard",
    enforce_blueprint_sync: bool = False,
    safe_area_profile: str = "legacy",
    safe_edge_whitelist: str = "header,footer",
    enable_visual_qa: bool = False,
    deterministic_repair: bool = False,
    auto_repair_failed_slides: bool = False,
    max_repair_rounds: int = 2,
    enforce_copyfit_contract: bool = True,
    native_text_verify: bool = True,
    native_auto_refit_rounds: int = 0,
    enable_preflight: bool = True,
    enable_doctor: bool = False,
    skip_qa: bool = False,
    profile: str = "presentation",
    quality_mode: str = "dev-fast",
    svg_only_export: bool = False,
    incremental: bool = False,
    changed_slides: list[int] | None = None,
) -> QaReport | None:
    started_at = datetime.now()
    started_perf = perf_counter()
    project_dir = project_dir.resolve()
    normalized_changed_slides = sorted({int(item) for item in (changed_slides or []) if int(item) > 0})
    repair_enabled = bool(deterministic_repair or auto_repair_failed_slides)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M")
    fallback_policy = load_fallback_policy(SCRIPT_DIR / "build_config" / "fallback_policy.json")
    fallback_events: list[dict[str, Any]] = []
    failure_error_code = "none"
    failure_action_hint = "none"
    slide_plan_stats: SlidePlanStats | None = None
    normalize_stats: NormalizeStats | None = None
    copyfit_stats: CopyfitStats | None = None
    delivery_stats: DeliveryOptimizeStats | None = None
    preflight_report: PreflightReport | None = None
    doctor_report: DoctorReport | None = None
    budget_report: BudgetReport | None = None
    planning_report: dict[str, Any] | None = None
    stage_timing: dict[str, float] = {
        "render_sec": 0.0,
        "preflight_sec": 0.0,
        "layout_lint_sec": 0.0,
        "finalize_sec": 0.0,
        "export_sec": 0.0,
        "qa_sec": 0.0,
        "snapshot_sec": 0.0,
    }
    safe_edge_whitelist_set = {token.strip() for token in safe_edge_whitelist.split(",") if token.strip()}
    source_hashes = _collect_incremental_source_hashes(project_dir, normalized_changed_slides or None)
    context_profile = context_obs.incremental_context_profile(
        project_dir,
        normalized_changed_slides or None,
        project_relative_path_fn=_project_relative_path,
    )
    token_budget_assessment = context_obs.evaluate_token_budget(
        phase=phase,
        changed_slides=normalized_changed_slides,
        context_profile=context_profile,
        policy=TOKEN_BUDGET_POLICY,
        caps=TOKEN_BUDGET_CAPS,
    )
    token_budget_assessment["token_budget_overflow_source"] = (
        str(token_budget_assessment.get("token_budget_stage", ""))
        if bool(token_budget_assessment.get("token_budget_warning"))
        else ""
    )
    page_summary_assessment = context_obs.update_page_summary_cache(
        project_dir,
        normalized_changed_slides,
        project_relative_path_fn=_project_relative_path,
        sha256_file_fn=_sha256_file,
        page_summary_cache_version=PAGE_SUMMARY_CACHE_VERSION,
    )
    template_lookup_assessment = context_obs.template_lookup_metrics(project_dir)
    checkpoint_assessment = context_obs.update_session_checkpoints(
        project_dir,
        changed_slides=normalized_changed_slides,
        token_budget_assessment=token_budget_assessment,
        project_relative_path_fn=_project_relative_path,
        page_summary_cache_version=PAGE_SUMMARY_CACHE_VERSION,
        session_checkpoint_version=SESSION_CHECKPOINT_VERSION,
        pages_per_chunk_default=SESSION_CHECKPOINT_DEFAULT_PAGES,
        pages_per_chunk_min=SESSION_CHECKPOINT_MIN_PAGES,
        pages_per_chunk_max=SESSION_CHECKPOINT_MAX_PAGES,
    )
    layout_exploration_assessment = _layout_exploration_assessment(project_dir)
    if token_budget_assessment["token_budget_warning"]:
        print(
            "warning: context token estimate exceeds budget cap "
            f"(stage={token_budget_assessment['token_budget_stage']}, "
            f"estimate={token_budget_assessment['context_token_estimate']}, "
            f"limit={token_budget_assessment['token_budget_limit']}).",
            file=sys.stderr,
        )

    existing_svg = sorted((project_dir / "svg_output").glob("slide_*.svg"))
    if not skip_render and not clean and existing_svg:
        if not force_render:
            recommended = f"python scripts/build_project.py {project_dir} --skip-render --snapshots"
            raise RuntimeError(
                "existing svg_output/slide_*.svg detected; render step may overwrite AI-authored slides. "
                "Use --skip-render (recommended) or pass --force-render to proceed intentionally.\n"
                f"example: {recommended}"
            )
        print(
            "warning: existing svg_output/slide_*.svg detected and --force-render is set; continuing render step.",
            file=sys.stderr,
        )

    if incremental:
        hit, reason, cached = _incremental_cache_hit(
            project_dir,
            phase=phase,
            source_hashes=source_hashes,
        )
        if hit:
            elapsed = round(perf_counter() - started_perf, 4)
            cached_files: list[str] = []
            if isinstance(cached, dict):
                raw = cached.get("output_files")
                if isinstance(raw, list):
                    cached_files = [str(item) for item in raw if isinstance(item, str)]
            cache_manifest_record: dict[str, Any] = {
                "project_name": project_dir.name,
                "timestamp": timestamp,
                "phase": phase,
                "mode": mode,
                "artifact_name_mode": artifact_name,
                "output_files": cached_files,
                "slide_count": _blueprint_slide_count(project_dir),
                "skip_render": skip_render,
                "strict": strict,
                "snapshots": snapshots,
                "profile": profile,
                "quality_mode": quality_mode,
                "svg_only_export": svg_only_export,
                "incremental_mode": True,
                "incremental_cache_hit": True,
                "incremental_skip_reason": reason,
                "source_hashes": source_hashes,
                "changed_slides": normalized_changed_slides,
                "build_started_at": started_at.isoformat(timespec="seconds"),
                "build_finished_at": datetime.now().isoformat(timespec="seconds"),
                "build_duration_sec": elapsed,
                "context_file_count": context_profile["context_file_count"],
                "context_bytes_estimate": context_profile["context_bytes_estimate"],
                "context_files": context_profile["context_files"],
                "fallback_policy_version": fallback_policy.version,
                "fallback_event_count": 0,
                "fallback_events": [],
                "failure_error_code": "none",
                "failure_action_hint": "none",
            }
            cache_manifest_record.update(
                context_obs.normalize_stage_metrics(
                    phase=phase,
                    stage_timing=stage_timing,
                    build_duration_sec=elapsed,
                )
            )
            cache_manifest_record.update(token_budget_assessment)
            cache_manifest_record.update(page_summary_assessment)
            cache_manifest_record.update(template_lookup_assessment)
            cache_manifest_record.update(checkpoint_assessment)
            cache_manifest_record.update(layout_exploration_assessment)
            _append_manifest(project_dir, cache_manifest_record)
            print(f"Wrote {project_dir / 'exports' / 'manifest.json'}")
            print(f"Incremental cache hit: {reason}. Skipped unchanged stage work.")
            return None

    budget_report = _evaluate_budget_policy_if_available(
        project_dir,
        profile=profile,
        svg_only_export=svg_only_export,
    )
    if budget_report is not None:
        print(
            "Budget policy: "
            f"profile={budget_report.profile}, "
            f"checked_slides={budget_report.checked_slides}, "
            f"overloaded_slides={budget_report.overloaded_slides}"
        )
        print(budget_report.report_path)

    if not svg_only_export:
        slide_plan_stats = ensure_semantic_slide_plan(project_dir)
        print(
            "Semantic slide plan: "
            f"slides_selected={slide_plan_stats.slides_selected}, "
            f"slides_with_copyfit_blocks={slide_plan_stats.slides_planned}, "
            f"blocks_total={slide_plan_stats.blocks_total}"
        )

    if auto_slide_plan:
        slide_plan_stats = generate_slide_plan(
            project_dir,
            overwrite=auto_slide_plan_overwrite,
        )
        print(
            "Auto slide plan: "
            f"slides_selected={slide_plan_stats.slides_selected}, "
            f"slides_planned={slide_plan_stats.slides_planned}, "
            f"blocks_total={slide_plan_stats.blocks_total}"
        )

    if not svg_only_export:
        planning_report = evaluate_planning_quality(
            project_dir,
            overloaded_slide_ids=budget_report.overloaded_slides if budget_report is not None else [],
        )
        planning_summary = planning_report.get("summary", {})
        planning_report_path = project_dir / "qa" / "planning-report.json"
        print(
            "Planning quality: "
            f"status={planning_report.get('planning_status')}, "
            f"errors={planning_summary.get('planning_error_count', 0)}, "
            f"notes={planning_summary.get('planning_note_count', 0)}"
        )
        print(planning_report_path)
        if planning_report.get("planning_status") == "blocked":
            overloaded = planning_summary.get("content_overload_slide_count", 0)
            prefix = "Budget policy check failed before rendering/export. " if overloaded else ""
            raise RuntimeError(
                prefix
                + "Planning quality gate failed before rendering/export. "
                + "Fix blueprint/slide_plan semantics and rerun.\n"
                + (
                    f"overloaded_slides={sorted(budget_report.overloaded_slides)}\n"
                    if budget_report is not None and budget_report.overloaded_slides
                    else ""
                )
                + f"report: {planning_report_path}"
            )

    if not skip_render:
        render_started = perf_counter()
        render_project(project_dir, output_dir="svg_output", clean=clean)
        stage_timing["render_sec"] = round(perf_counter() - render_started, 4)

    if normalize_images_white_bg:
        normalize_stats = normalize_project_images_white_bg(project_dir)
        print(
            "Image normalize (white bg): "
            f"scanned={normalize_stats.scanned}, "
            f"converted={normalize_stats.converted}, "
            f"skipped={normalize_stats.skipped}"
        )

    if enable_preflight:
        preflight_started = perf_counter()
        preflight_report = run_preflight(
            project_dir,
            svg_dir_name="svg_output",
            slide_ids=set(normalized_changed_slides) if normalized_changed_slides else None,
        )
        stage_timing["preflight_sec"] = round(perf_counter() - preflight_started, 4)
        print(
            "Preflight "
            f"{'passed' if preflight_report.ok else 'failed'}: "
            f"errors={preflight_report.errors}, warnings={preflight_report.warnings}"
        )
        print(preflight_report.report_md)
        if not preflight_report.ok:
            raise RuntimeError(
                "Preflight failed before layout/finalize. "
                "Fix svg_output/ and rerun.\n"
                f"report: {preflight_report.report_md}"
            )

    contract_report: ContractReport | None = None
    if enforce_copyfit_contract:
        contract_report = validate_copyfit_contract(project_dir)
        should_try_dense_autoplan = (
            not contract_report.ok
            and bool(contract_report.dense_slide_ids)
            and any(
                finding.code in {"missing-slide-plan", "missing-dense-slide-blocks"}
                for finding in contract_report.findings
            )
        )
        if should_try_dense_autoplan:
            if not (project_dir / "slide_plan.json").exists() or auto_slide_plan_overwrite:
                slide_plan_stats = generate_slide_plan(
                    project_dir,
                    overwrite=auto_slide_plan_overwrite,
                )
                print(
                    "Auto slide plan (dense contract): "
                    f"slides_selected={slide_plan_stats.slides_selected}, "
                    f"slides_planned={slide_plan_stats.slides_planned}, "
                    f"blocks_total={slide_plan_stats.blocks_total}"
                )
                contract_report = validate_copyfit_contract(project_dir)
        print(
            "Copyfit contract "
            f"{'passed' if contract_report.ok else 'failed'}: "
            f"errors={contract_report.errors}, warnings={contract_report.warnings}"
        )
        print(contract_report.report_path)
        if not contract_report.ok:
            raise RuntimeError(
                "Copyfit contract check failed before export. "
                "Dense slides require valid slide_plan.json blocks.\n"
                f"report: {contract_report.report_path}"
            )

    strict_delivery_gate = phase == "finalize" and (strict or quality_mode in {"release-safe", "premium"})
    if strict_delivery_gate:
        auto_selected_note = ensure_style_draft_selection(project_dir, auto_select_first=True)
        if auto_selected_note:
            print(auto_selected_note)
        style_gate_issues = validate_style_draft_selection(project_dir)
        if style_gate_issues:
            details = "\n".join(f"- {issue}" for issue in style_gate_issues)
            raise RuntimeError(
                "Style drafts gate failed before export. "
                "Low-confidence style routing requires valid style_drafts selection.\n"
                f"{details}"
            )

    if auto_copyfit:
        copyfit_stats = run_auto_copyfit(
            project_dir,
            profile=guard_profile,
            svg_dir_name="svg_output",
            safe_area_profile=safe_area_profile,
        )
        print(
            "Auto copyfit: "
            f"changed_files={copyfit_stats.changed_files}, "
            f"adjusted_text_nodes={copyfit_stats.adjusted_text_nodes}, "
            f"total_size_reduction={copyfit_stats.total_size_reduction:.1f}"
        )
        print(copyfit_stats.report_path)

    if delivery_ready:
        delivery_stats = run_delivery_optimize(project_dir, svg_dir_name="svg_output", profile=guard_profile)
        print(
            "Delivery optimize: "
            f"changed_files={delivery_stats.changed_files}, "
            f"removed_footer_nodes={delivery_stats.removed_footer_nodes}, "
            f"normalized_font_nodes={delivery_stats.normalized_font_nodes}, "
            f"divider_adjusted={delivery_stats.divider_adjusted}"
        )
        print(delivery_stats.report_path)

    layout_lint_report: LayoutLintReport | None = None
    if enable_layout_lint:
        layout_lint_started = perf_counter()
        if normalized_changed_slides:
            scoped_reports: list[LayoutLintReport] = []
            for slide_id in normalized_changed_slides:
                scoped_reports.append(
                    run_layout_lint(
                        project_dir,
                        svg_dir_name="svg_output",
                        strict=layout_lint_strict,
                        quality_profile=layout_quality_profile,
                        safe_area_profile=safe_area_profile,
                        safe_edge_whitelist=safe_edge_whitelist_set,
                        slide_id=slide_id,
                        profile=profile,
                        quality_mode=quality_mode,
                    )
                )
            layout_lint_report = LayoutLintReport(
                project=str(project_dir),
                ok=all(item.ok for item in scoped_reports),
                errors=sum(item.errors for item in scoped_reports),
                warnings=sum(item.warnings for item in scoped_reports),
                findings=[finding for item in scoped_reports for finding in item.findings],
                metrics={
                    "checked_slide": normalized_changed_slides,
                    "svg_files": sum(int(item.metrics.get("svg_files", 0)) for item in scoped_reports),
                    "quality_warning_count": sum(
                        int(item.metrics.get("quality_warning_count", 0)) for item in scoped_reports
                    ),
                    "quality_error_count": sum(
                        int(item.metrics.get("quality_error_count", 0)) for item in scoped_reports
                    ),
                    "quality_advisory_count": sum(
                        int(item.metrics.get("quality_advisory_count", 0)) for item in scoped_reports
                    ),
                    "overflow_risk_high_count": sum(
                        int(item.metrics.get("overflow_risk_high_count", 0)) for item in scoped_reports
                    ),
                    "overflow_risk_medium_count": sum(
                        int(item.metrics.get("overflow_risk_medium_count", 0)) for item in scoped_reports
                    ),
                    "overflow_risk_low_count": sum(
                        int(item.metrics.get("overflow_risk_low_count", 0)) for item in scoped_reports
                    ),
                },
                advisories=sum(item.advisories for item in scoped_reports),
            )
        else:
            layout_lint_report = run_layout_lint(
                project_dir,
                svg_dir_name="svg_output",
                strict=layout_lint_strict,
                quality_profile=layout_quality_profile,
                safe_area_profile=safe_area_profile,
                safe_edge_whitelist=safe_edge_whitelist_set,
                profile=profile,
                quality_mode=quality_mode,
            )
        print(
            "Layout lint "
            f"{'passed' if layout_lint_report.ok else 'failed'}: "
            f"errors={layout_lint_report.errors}, warnings={layout_lint_report.warnings}"
        )
        print(project_dir / "qa" / "layout-lint-report.md")
        if not layout_lint_report.ok:
            raise RuntimeError(
                "Layout lint failed before export. "
                "Fix svg_output/ and rerun, or disable lint for this build.\n"
                f"report: {project_dir / 'qa' / 'layout-lint-report.md'}"
            )
        stage_timing["layout_lint_sec"] = round(perf_counter() - layout_lint_started, 4)

    authoring_handled = handle_authoring_phase(
        phase=phase,
        project_dir=project_dir,
        timestamp=timestamp,
        mode=mode,
        artifact_name=artifact_name,
        skip_render=skip_render,
        strict=strict,
        snapshots=snapshots,
        enable_layout_lint=enable_layout_lint,
        layout_lint_strict=layout_lint_strict,
        layout_quality_profile=layout_quality_profile,
        enforce_blueprint_sync=enforce_blueprint_sync,
        safe_area_profile=safe_area_profile,
        profile=profile,
        quality_mode=quality_mode,
        svg_only_export=svg_only_export,
        safe_edge_whitelist_set=safe_edge_whitelist_set,
        incremental=incremental,
        source_hashes=source_hashes,
        normalized_changed_slides=normalized_changed_slides,
        started_at=started_at,
        started_perf=started_perf,
        context_profile=context_profile,
        fallback_policy_version=fallback_policy.version,
        token_budget_assessment=token_budget_assessment,
        page_summary_assessment=page_summary_assessment,
        template_lookup_assessment=template_lookup_assessment,
        checkpoint_assessment=checkpoint_assessment,
        layout_exploration_assessment=layout_exploration_assessment,
        stage_timing=stage_timing,
        preflight_report=preflight_report,
        layout_lint_report=layout_lint_report,
        append_manifest=_append_manifest,
        write_incremental_cache=_write_incremental_cache,
        blueprint_slide_count=_blueprint_slide_count,
        extract_quality_tiers=extract_quality_tiers,
        project_relative_path=_project_relative_path,
        normalize_stage_metrics=context_obs.normalize_stage_metrics,
    )
    if authoring_handled:
        return None

    if enable_doctor:
        doctor_report = run_export_doctor(
            project_dir,
            expected_cwd=SCRIPT_DIR.parent,
            check_output_lock=True,
        )
        print(
            "Export doctor "
            f"{'passed' if doctor_report.ok else 'failed'}: "
            f"errors={doctor_report.errors}, warnings={doctor_report.warnings}"
        )
        print(doctor_report.report_md)
        if not doctor_report.ok:
            raise RuntimeError(
                "Export doctor failed before finalize/export. "
                "Fix doctor findings and rerun.\n"
                f"report: {doctor_report.report_md}"
            )

    _raise_on_invalid_svg_encoding(project_dir)

    finalize_started = perf_counter()
    try:
        warning_count = finalize_project(project_dir)
    except SystemExit as exc:
        code = exc.code if isinstance(exc.code, int) else 1
        if code != 0:
            raise RuntimeError(f"Finalize failed with exit code {code}.") from exc
        warning_count = 0
    stage_timing["finalize_sec"] = round(perf_counter() - finalize_started, 4)
    if warning_count:
        print(f"Finalize completed with {warning_count} warning(s).", file=sys.stderr)
        if strict:
            raise RuntimeError(f"Finalize produced {warning_count} warning(s) in strict mode.")

    export_started = perf_counter()
    output_plan = build_output_plan(project_dir, mode, artifact_name, timestamp)
    output_files: list[Path] = []
    written_targets: set[Path] = set()
    locked_fallbacks: list[dict[str, str]] = []
    locked_fallback_attempts = 0
    planned_targets = {target.resolve() for targets in output_plan.values() for target in targets}
    native_conversion_report: dict[str, Any] | None = None
    for export_mode, targets in output_plan.items():
        planned_semantic = planned_semantic_target(targets, export_mode)
        for target in targets:
            target_resolved = target.resolve()
            if target_resolved in written_targets:
                continue
            try:
                converted = convert(
                    project_dir,
                    "svg_final",
                    target,
                    mode=export_mode,
                    timeout_sec=fallback_policy.export_subprocess_timeout_sec,
                )
            except Exception as exc:
                error_code = classify_export_error(exc)
                failure_error_code = error_code
                failure_action_hint = fallback_action_hint(error_code)
                if (
                    looks_like_locked_pptx_error(exc)
                    and is_stable_target(target, export_mode)
                    and fallback_policy.fallback_on_locked_stable_target
                    and locked_fallback_attempts < max(0, fallback_policy.max_locked_fallback_attempts)
                ):
                    locked_fallback_attempts += 1
                    fallback_target = planned_semantic
                    if fallback_target is None or fallback_target.resolve() in written_targets:
                        fallback_target = next_semantic_target(
                            project_dir,
                            export_mode,
                            timestamp,
                            reserved=planned_targets | written_targets,
                        )
                    fallback_events.append(
                        {
                            "stage": "export",
                            "error_code": error_code,
                            "trigger": "stable_target_locked",
                            "action": "retry_with_semantic_target",
                            "requested_target": _project_relative_path(project_dir, target),
                            "fallback_target": _project_relative_path(project_dir, fallback_target),
                            "result": "retrying",
                        }
                    )
                    print(
                        f"warning: {target.name} is locked; retrying export as {fallback_target.name}",
                        file=sys.stderr,
                    )
                    try:
                        converted = convert(
                            project_dir,
                            "svg_final",
                            fallback_target,
                            mode=export_mode,
                            timeout_sec=fallback_policy.export_subprocess_timeout_sec,
                        )
                    except Exception as fallback_exc:
                        fallback_events.append(
                            {
                                "stage": "export",
                                "error_code": classify_export_error(fallback_exc),
                                "trigger": "semantic_retry",
                                "action": "failed",
                                "requested_target": _project_relative_path(project_dir, target),
                                "fallback_target": _project_relative_path(project_dir, fallback_target),
                                "result": "failed",
                            }
                        )
                        failure_error_code = classify_export_error(fallback_exc)
                        failure_action_hint = fallback_action_hint(failure_error_code)
                        if looks_like_locked_pptx_error(fallback_exc):
                            raise RuntimeError(format_pptx_write_error(fallback_target, fallback_exc)) from fallback_exc
                        raise
                    written_path = converted if isinstance(converted, Path) else fallback_target
                    output_files.append(written_path.resolve())
                    written_targets.add(fallback_target.resolve())
                    fallback_events.append(
                        {
                            "stage": "export",
                            "error_code": error_code,
                            "trigger": "semantic_retry",
                            "action": "fallback_succeeded",
                            "requested_target": _project_relative_path(project_dir, target),
                            "fallback_target": _project_relative_path(project_dir, written_path),
                            "result": "success",
                        }
                    )
                    locked_fallbacks.append(
                        {
                            "export_mode": export_mode,
                            "requested_target": _project_relative_path(project_dir, target),
                            "fallback_target": _project_relative_path(project_dir, written_path),
                        }
                    )
                    print(f"Wrote {written_path}")
                    continue
                fallback_events.append(
                    {
                        "stage": "export",
                        "error_code": error_code,
                        "trigger": "export_exception",
                        "action": "no_fallback",
                        "requested_target": _project_relative_path(project_dir, target),
                        "result": "failed",
                    }
                )
                if looks_like_locked_pptx_error(exc):
                    raise RuntimeError(format_pptx_write_error(target, exc)) from exc
                raise
            written_path = converted if isinstance(converted, Path) else target
            output_files.append(written_path.resolve())
            written_targets.add(target_resolved)
            print(f"Wrote {written_path}")

    native_verify_report: NativeTextVerifyReport | None = None
    if native_text_verify:
        native_outputs = [
            path for path in output_files if "--native--" in path.name or path.name == "output-native.pptx"
        ]
        native_verify_report = _verify_native_text_layout(native_outputs)
        print(
            "Native text verify: "
            f"checked_pptx={native_verify_report.checked_pptx}, "
            f"checked_lines={native_verify_report.checked_lines}, "
            f"narrow_column_lines={native_verify_report.narrow_column_lines}, "
            f"risk_lines={native_verify_report.risk_lines}"
        )
        if not native_verify_report.ok and native_auto_refit_rounds > 0:
            for round_index in range(1, int(native_auto_refit_rounds) + 1):
                print(
                    f"Native text verify refit round {round_index}: "
                    "applying strict copyfit and re-exporting native outputs."
                )
                run_auto_copyfit(
                    project_dir,
                    profile="strict",
                    svg_dir_name="svg_output",
                    safe_area_profile=safe_area_profile,
                )
                finalize_project(project_dir)
                refreshed: list[Path] = []
                for path in native_outputs:
                    refreshed.append(
                        convert(
                            project_dir,
                            "svg_final",
                            path,
                            mode="native",
                            timeout_sec=fallback_policy.export_subprocess_timeout_sec,
                        )
                    )
                native_outputs = [path.resolve() for path in refreshed]
                output_files = [path for path in output_files if path not in native_outputs] + native_outputs
                native_verify_report = _verify_native_text_layout(native_outputs)
                print(
                    "Native text verify after refit: "
                    f"checked_lines={native_verify_report.checked_lines}, "
                    f"narrow_column_lines={native_verify_report.narrow_column_lines}, "
                    f"risk_lines={native_verify_report.risk_lines}"
                )
                if native_verify_report.ok:
                    break
        if not native_verify_report.ok:
            sample = native_verify_report.findings[:3]
            raise RuntimeError(
                "Native text verification failed: narrow-column CJK lines likely overflow in exported PPTX. "
                f"risk_lines={native_verify_report.risk_lines}, sample={sample}"
            )
    if "native" in output_plan:
        native_conversion_report = load_native_conversion_report(project_dir)
    stage_timing["export_sec"] = round(perf_counter() - export_started, 4)

    output_file_entries = [_project_relative_path(project_dir, path) for path in output_files]
    artifact_created = len(output_file_entries) > 0
    destructive_reasons: list[str] = []
    if auto_copyfit:
        destructive_reasons.append("auto_copyfit")
    if delivery_ready:
        destructive_reasons.append("delivery_ready")
    if int(native_auto_refit_rounds) > 0:
        destructive_reasons.append("native_auto_refit")
    destructive_path = bool(destructive_reasons)
    pre_qa_status = "artifact_only" if skip_qa else "pending_qa"
    pre_qa_elapsed = round(perf_counter() - started_perf, 4)
    pre_qa_stage_metrics = context_obs.normalize_stage_metrics(
        phase=phase,
        stage_timing=stage_timing,
        build_duration_sec=pre_qa_elapsed,
    )

    manifest_record: dict[str, Any] = {
        "project_name": project_dir.name,
        "timestamp": timestamp,
        "phase": phase,
        "mode": mode,
        "artifact_name_mode": artifact_name,
        "output_files": output_file_entries,
        "slide_count": _blueprint_slide_count(project_dir),
        "skip_render": skip_render,
        "strict": strict,
        "snapshots": snapshots,
        "layout_lint_enabled": enable_layout_lint,
        "layout_lint_strict": layout_lint_strict,
        "layout_quality_profile": layout_quality_profile,
        "enforce_blueprint_sync": enforce_blueprint_sync,
        "safe_area_profile": safe_area_profile,
        "profile": profile,
        "quality_mode": quality_mode,
        "svg_only_export": svg_only_export,
        "safe_edge_whitelist": sorted(safe_edge_whitelist_set),
        "enable_visual_qa": enable_visual_qa,
        "deterministic_repair": repair_enabled,
        "auto_repair_failed_slides": repair_enabled,
        "max_repair_rounds": max_repair_rounds,
        "enforce_copyfit_contract": enforce_copyfit_contract,
        "native_text_verify": native_text_verify,
        "native_auto_refit_rounds": native_auto_refit_rounds,
        "enable_preflight": enable_preflight,
        "enable_doctor": enable_doctor,
        "auto_slide_plan": auto_slide_plan,
        "auto_slide_plan_overwrite": auto_slide_plan_overwrite,
        "normalize_images_white_bg": normalize_images_white_bg,
        "guard_profile": guard_profile,
        "auto_copyfit": auto_copyfit,
        "delivery_ready": delivery_ready,
        "destructive_path": destructive_path,
        "destructive_reasons": destructive_reasons,
        "incremental_mode": incremental,
        "incremental_cache_hit": False,
        "source_hashes": source_hashes,
        "changed_slides": normalized_changed_slides,
        "build_started_at": started_at.isoformat(timespec="seconds"),
        "context_file_count": context_profile["context_file_count"],
        "context_bytes_estimate": context_profile["context_bytes_estimate"],
        "context_files": context_profile["context_files"],
        "fallback_policy_version": fallback_policy.version,
    }
    manifest_record.update(token_budget_assessment)
    manifest_record.update(page_summary_assessment)
    manifest_record.update(template_lookup_assessment)
    manifest_record.update(checkpoint_assessment)
    manifest_record.update(layout_exploration_assessment)
    manifest_record.update(stage_timing)
    manifest_record.update(pre_qa_stage_metrics)
    if slide_plan_stats is not None:
        manifest_record["auto_slide_plan_slides_planned"] = slide_plan_stats.slides_planned
        manifest_record["auto_slide_plan_blocks_total"] = slide_plan_stats.blocks_total
    if normalize_stats is not None:
        manifest_record["normalized_images_scanned"] = normalize_stats.scanned
        manifest_record["normalized_images_converted"] = normalize_stats.converted
        manifest_record["normalized_images_skipped"] = normalize_stats.skipped
    if copyfit_stats is not None:
        manifest_record["copyfit_changed_files"] = copyfit_stats.changed_files
        manifest_record["copyfit_adjusted_text_nodes"] = copyfit_stats.adjusted_text_nodes
        manifest_record["copyfit_total_size_reduction"] = round(copyfit_stats.total_size_reduction, 2)
        manifest_record["copyfit_report"] = _project_relative_path(project_dir, copyfit_stats.report_path)
    if delivery_stats is not None:
        manifest_record["delivery_changed_files"] = delivery_stats.changed_files
        manifest_record["delivery_removed_footer_nodes"] = delivery_stats.removed_footer_nodes
        manifest_record["delivery_normalized_font_nodes"] = delivery_stats.normalized_font_nodes
        manifest_record["delivery_divider_adjusted"] = delivery_stats.divider_adjusted
        manifest_record["delivery_report"] = _project_relative_path(project_dir, delivery_stats.report_path)
    if layout_lint_report is not None:
        manifest_record["layout_lint_errors"] = layout_lint_report.errors
        manifest_record["layout_lint_warnings"] = layout_lint_report.warnings
        manifest_record["layout_lint_advisories"] = layout_lint_report.advisories
        manifest_record["layout_quality_warning_count"] = int(
            layout_lint_report.metrics.get("quality_warning_count", 0)
        )
        manifest_record["layout_quality_error_count"] = int(
            layout_lint_report.metrics.get("quality_error_count", 0)
        )
        manifest_record["layout_quality_advisory_count"] = int(
            layout_lint_report.metrics.get("quality_advisory_count", 0)
        )
        manifest_record["overflow_risk_high_count"] = int(layout_lint_report.metrics.get("overflow_risk_high_count", 0))
        manifest_record["overflow_risk_medium_count"] = int(
            layout_lint_report.metrics.get("overflow_risk_medium_count", 0)
        )
        manifest_record["overflow_risk_low_count"] = int(layout_lint_report.metrics.get("overflow_risk_low_count", 0))
        manifest_record["overflow_risk_blocking_count"] = (
            manifest_record["overflow_risk_high_count"] if quality_mode in {"release-safe", "premium"} else 0
        )
        manifest_record["layout_quality_segment"] = layout_quality_profile
        manifest_record["layout_lint_quality_tiers"] = extract_quality_tiers(
            layout_lint_report.metrics,
            fallback_errors=layout_lint_report.errors,
            fallback_warnings=layout_lint_report.warnings,
        )
    else:
        manifest_record["layout_lint_quality_tiers"] = {"blocking": 0, "warning": 0, "advisory": 0}
    if preflight_report is not None:
        manifest_record["preflight_errors"] = preflight_report.errors
        manifest_record["preflight_warnings"] = preflight_report.warnings
        manifest_record["preflight_report"] = _project_relative_path(project_dir, preflight_report.report_md)
    if doctor_report is not None:
        manifest_record["doctor_errors"] = doctor_report.errors
        manifest_record["doctor_warnings"] = doctor_report.warnings
        manifest_record["doctor_report"] = _project_relative_path(project_dir, doctor_report.report_md)
    if contract_report is not None:
        manifest_record["copyfit_contract_errors"] = contract_report.errors
        manifest_record["copyfit_contract_warnings"] = contract_report.warnings
        manifest_record["copyfit_contract_dense_slides"] = contract_report.dense_slide_ids
        manifest_record["copyfit_contract_covered_dense_slides"] = contract_report.covered_dense_slide_ids
        manifest_record["copyfit_contract_report"] = _project_relative_path(project_dir, contract_report.report_path)
    if native_verify_report is not None:
        manifest_record["native_text_verify_checked_pptx"] = native_verify_report.checked_pptx
        manifest_record["native_text_verify_checked_lines"] = native_verify_report.checked_lines
        manifest_record["native_text_verify_narrow_lines"] = native_verify_report.narrow_column_lines
        manifest_record["native_text_verify_risk_lines"] = native_verify_report.risk_lines
    if native_conversion_report is not None:
        manifest_record.update(
            _native_conversion_manifest_fields(
                Path(_project_relative_path(project_dir, native_conversion_report_path(project_dir))),
                native_conversion_report,
            )
        )
    if budget_report is not None:
        manifest_record["budget_profile"] = budget_report.profile
        manifest_record["budget_checked_slides"] = budget_report.checked_slides
        manifest_record["budget_overloaded_slides"] = budget_report.overloaded_slides
        manifest_record["budget_report"] = _project_relative_path(project_dir, budget_report.report_path)
    if planning_report is not None:
        manifest_record.update(
            _planning_manifest_fields(
                Path("qa/planning-report.json"),
                planning_report,
            )
        )
    git_sha = _git_sha(project_dir)
    if git_sha:
        manifest_record["git_sha"] = git_sha
    template_key = _template_key(project_dir)
    if template_key:
        manifest_record["template_key"] = template_key
    blueprint_hash = _blueprint_hash(project_dir)
    if blueprint_hash:
        manifest_record["blueprint_hash"] = blueprint_hash
    if locked_fallbacks:
        manifest_record["locked_pptx_fallback_count"] = len(locked_fallbacks)
        manifest_record["locked_pptx_fallbacks"] = locked_fallbacks

    manifest_record["artifact_created"] = artifact_created
    manifest_record["delivery_approved"] = False
    manifest_record["delivery_status"] = pre_qa_status
    manifest_record["delivery_failure_reasons"] = ["qa_not_executed"] if skip_qa else []
    manifest_record["fallback_event_count"] = len(fallback_events)
    manifest_record["fallback_events"] = fallback_events
    manifest_record["failure_error_code"] = failure_error_code
    manifest_record["failure_action_hint"] = failure_action_hint
    manifest_record["build_finished_at"] = datetime.now().isoformat(timespec="seconds")
    manifest_record["build_duration_sec"] = pre_qa_elapsed
    manifest_record.update(
        context_obs.infer_stage_failure(
            delivery_approved=False,
            delivery_status=pre_qa_status,
            delivery_failure_reasons=manifest_record["delivery_failure_reasons"],
        )
    )
    manifest_record["quality_tiers"] = manifest_record.get(
        "layout_lint_quality_tiers", {"blocking": 0, "warning": 0, "advisory": 0}
    )
    _append_manifest(project_dir, manifest_record)
    print(f"Wrote {project_dir / 'exports' / 'manifest.json'}")

    skip_qa_handled = handle_finalize_skip_qa_phase(
        phase=phase,
        skip_qa=skip_qa,
        artifact_created=artifact_created,
        project_dir=project_dir,
        stage_timing=stage_timing,
        started_perf=started_perf,
        output_file_entries=output_file_entries,
        incremental=incremental,
        source_hashes=source_hashes,
        failure_error_code=failure_error_code,
        failure_action_hint=failure_action_hint,
        update_manifest_last_record=_update_manifest_last_record,
        print_delivery_status=_print_delivery_status,
        write_incremental_cache=_write_incremental_cache,
        normalize_stage_metrics=context_obs.normalize_stage_metrics,
        infer_stage_failure=context_obs.infer_stage_failure,
        delivery_status_tuple=_delivery_status_tuple,
        fallback_action_hint=fallback_action_hint,
    )
    if skip_qa_handled:
        return None

    qa_started = perf_counter()
    report = run_qa(
        project_dir,
        svg_dir_name="svg_final",
        snapshots=snapshots,
        strict=strict,
        pptx=output_file_entries,
        slide_id=normalized_changed_slides[0] if len(normalized_changed_slides) == 1 else None,
        enforce_blueprint_sync=enforce_blueprint_sync,
        safe_area_profile=safe_area_profile,
        enable_visual_qa=enable_visual_qa,
        profile=profile,
        quality_mode=quality_mode,
        safe_edge_whitelist=safe_edge_whitelist_set,
    )
    stage_timing["qa_sec"] = round(perf_counter() - qa_started, 4)
    stage_timing["snapshot_sec"] = round(float(report.metrics.get("snapshot_render_sec", 0.0)), 4)
    qa_quality_tiers = extract_quality_tiers(
        report.metrics,
        fallback_errors=report.errors,
        fallback_warnings=report.warnings,
    )
    post_qa_elapsed = round(perf_counter() - started_perf, 4)
    post_qa_stage_metrics = context_obs.normalize_stage_metrics(
        phase=phase,
        stage_timing=stage_timing,
        build_duration_sec=post_qa_elapsed,
    )
    _update_manifest_last_record(
        project_dir,
        {
            "qa_errors": report.errors,
            "qa_warnings": report.warnings,
            "qa_advisories": report.advisories,
            "qa_quality_tiers": qa_quality_tiers,
            "quality_tiers": qa_quality_tiers,
            "qa_visual_score": report.visual_score,
            "qa_visual_findings_count": len(report.visual_findings or []),
            "qa_repair_recommendation_count": len(report.repair_recommendation or []),
            "qa_visual_diversity_warning_count": int(report.metrics.get("visual_diversity_warning_count", 0)),
            "layout_exploration_enabled": bool(report.metrics.get("layout_exploration_enabled", False)),
            "candidate_count": int(report.metrics.get("layout_candidate_count", 0)),
            "archetype_switch_count": int(report.metrics.get("archetype_switch_count", 0)),
            "consecutive_repeat_count": int(report.metrics.get("consecutive_repeat_count", 0)),
            "diversity_gate_result": str(report.metrics.get("diversity_gate_result", "unknown")),
            "qa_density_flag": report.density_flag,
            "qa_hierarchy_flag": report.hierarchy_flag,
            "qa_sec": stage_timing["qa_sec"],
            "snapshot_sec": stage_timing["snapshot_sec"],
            "build_finished_at": datetime.now().isoformat(timespec="seconds"),
            "build_duration_sec": post_qa_elapsed,
            **post_qa_stage_metrics,
        },
    )
    delivery_approved, delivery_status, delivery_failure_reasons = _delivery_status_tuple(
        artifact_created=artifact_created,
        qa_report=report,
        native_conversion_report=native_conversion_report,
        planning_report=planning_report,
    )
    _update_manifest_last_record(
        project_dir,
        {
            "artifact_created": artifact_created,
            "delivery_ready": delivery_approved,
            "delivery_approved": delivery_approved,
            "delivery_status": delivery_status,
            "delivery_failure_reasons": delivery_failure_reasons,
            **context_obs.infer_stage_failure(
                delivery_approved=delivery_approved,
                delivery_status=delivery_status,
                delivery_failure_reasons=delivery_failure_reasons,
            ),
            "failure_error_code": (
                "none"
                if delivery_approved
                else (failure_error_code if failure_error_code != "none" else "qa_failed")
            ),
            "failure_action_hint": (
                "none"
                if delivery_approved
                else (failure_action_hint if failure_action_hint != "none" else "review_qa_report")
            ),
        },
    )
    _print_delivery_status(
        artifact_created=artifact_created,
        delivery_approved=delivery_approved,
        delivery_status=delivery_status,
        delivery_failure_reasons=delivery_failure_reasons,
        qa_report_path=str(project_dir / "qa" / "report.md"),
        qa_report=report,
    )
    if not enable_visual_qa:
        if incremental:
            _write_incremental_cache(
                project_dir,
                {
                    "phase": phase,
                    "source_hashes": source_hashes,
                    "output_files": output_file_entries,
                    "qa_ok": report.ok,
                },
            )
        return report

    if not repair_enabled:
        advisory_slides = collect_repair_slide_ids(report)
        if advisory_slides:
            note_path = write_repair_round_note(project_dir, 1, report, mode="advisory")
            print(f"Visual QA advisory note: slides={advisory_slides}; note={note_path}")
        if incremental:
            _write_incremental_cache(
                project_dir,
                {
                    "phase": phase,
                    "source_hashes": source_hashes,
                    "output_files": output_file_entries,
                    "qa_ok": report.ok,
                },
            )
        return report

    current_report: QaReport | None = report
    rounds = max(0, int(max_repair_rounds))
    for round_index in range(1, rounds + 1):
        failed_slides = collect_repair_slide_ids(current_report)
        if not failed_slides:
            break
        repaired = apply_deterministic_repair(project_dir, failed_slides, render_project)
        note_path = write_repair_round_note(
            project_dir,
            round_index,
            current_report,
            mode="deterministic",
            replacement_source="render_svg::__repair_tmp_svg_output",
            replaced_slide_ids=repaired,
        )
        if not repaired:
            print(f"warning: auto repair round {round_index} found no replaceable slide artifacts.", file=sys.stderr)
            break
        audit_path = append_repair_audit(
            project_dir,
            round_index=round_index,
            requested_slide_ids=failed_slides,
            repaired_slide_ids=repaired,
            replacement_source="render_svg::__repair_tmp_svg_output",
        )
        print(
            f"Deterministic repair round {round_index}: repaired slides {repaired}; "
            f"note={note_path}; audit={audit_path}"
        )
        current_report = build_project(
            project_dir,
            mode=mode,
            artifact_name=artifact_name,
            clean=False,
            snapshots=snapshots,
            strict=strict,
            skip_render=True,
            force_render=force_render,
            auto_slide_plan=False,
            auto_slide_plan_overwrite=False,
            guard_profile=guard_profile,
            auto_copyfit=auto_copyfit,
            delivery_ready=delivery_ready,
            enable_layout_lint=enable_layout_lint,
            layout_lint_strict=layout_lint_strict,
            layout_quality_profile=layout_quality_profile,
            enforce_blueprint_sync=enforce_blueprint_sync,
            safe_area_profile=safe_area_profile,
            safe_edge_whitelist=safe_edge_whitelist,
            enable_visual_qa=enable_visual_qa,
            deterministic_repair=False,
            auto_repair_failed_slides=False,
            max_repair_rounds=0,
            enforce_copyfit_contract=enforce_copyfit_contract,
            native_text_verify=native_text_verify,
            native_auto_refit_rounds=native_auto_refit_rounds,
            enable_preflight=enable_preflight,
            enable_doctor=False,
            skip_qa=False,
            profile=profile,
            quality_mode=quality_mode,
            svg_only_export=svg_only_export,
        )
        if current_report is None:
            break
    if incremental:
        final_report = current_report if current_report is not None else report
        _write_incremental_cache(
            project_dir,
            {
                "phase": phase,
                "source_hashes": source_hashes,
                "output_files": output_file_entries,
                "qa_ok": bool(final_report.ok) if final_report is not None else bool(report.ok),
            },
        )
    return current_report


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build an AI-PPT project from blueprint to PPTX and QA report.")
    parser.add_argument("project_dir", type=Path, help="Path to projects/<project_name>.")
    parser.add_argument(
        "--phase",
        choices=("full", "authoring", "finalize"),
        default="finalize",
        help=(
            "Pipeline phase: finalize=default full export/QA; "
            "authoring=preflight/lint only; full=legacy compatibility path."
        ),
    )
    parser.add_argument(
        "--mode",
        choices=("raster", "native", "both"),
        default="native",
        help=(
            "PPTX export: native=editable output-native.pptx only (default); "
            "raster=PNG embed output.pptx; both=two files."
        ),
    )
    parser.add_argument(
        "--artifact-name",
        choices=("stable", "semantic", "both"),
        default="stable",
        help="Output naming strategy: stable aliases, semantic names, or both.",
    )
    parser.add_argument("--clean", action="store_true", help="Delete svg_output before rendering.")
    parser.add_argument("--snapshots", action="store_true", help="Render QA snapshots and contact sheet.")
    parser.add_argument("--strict", action="store_true", help="Treat QA warnings as failures.")
    parser.add_argument(
        "--skip-render",
        action="store_true",
        help="Use existing svg_output instead of rendering blueprint.json (recommended for AI-authored slides).",
    )
    parser.add_argument(
        "--force-render",
        action="store_true",
        help="Allow render step even when existing svg_output/slide_*.svg are present.",
    )
    parser.add_argument(
        "--auto-slide-plan",
        action="store_true",
        help="Generate slide_plan.json draft from blueprint.json before render/contract gates.",
    )
    parser.add_argument(
        "--auto-slide-plan-overwrite",
        action="store_true",
        help="Overwrite existing slide_plan.json when using --auto-slide-plan.",
    )
    parser.add_argument(
        "--normalize-images-white-bg",
        action="store_true",
        help="Normalize project images/*.png with alpha to fixed white background before finalize/export.",
    )
    parser.add_argument(
        "--guard-profile",
        choices=("preserve-design", "balanced", "strict"),
        default="balanced",
        help="Shared strategy profile for auto-copyfit and delivery-ready mechanisms.",
    )
    parser.add_argument(
        "--auto-copyfit",
        action=argparse.BooleanOptionalAction,
        default=None,
        help="Enable/disable auto-shrink guard; default is off (non-destructive finalize).",
    )
    parser.add_argument(
        "--delivery-ready",
        action=argparse.BooleanOptionalAction,
        default=None,
        help="Enable/disable delivery cleanup; default is off (non-destructive finalize).",
    )
    parser.add_argument(
        "--enable-layout-lint",
        action="store_true",
        help="Run SVG layout lint gate after State 3 and before finalize/export.",
    )
    parser.add_argument(
        "--layout-lint-strict",
        action="store_true",
        help="Treat layout lint warnings as failures (only used with --enable-layout-lint).",
    )
    parser.add_argument(
        "--layout-quality-profile",
        choices=("standard", "strict"),
        default="standard",
        help="Severity profile for layout quality metrics (standard warns, strict may escalate).",
    )
    parser.add_argument(
        "--enforce-blueprint-sync",
        action="store_true",
        help="Treat blueprint-to-SVG mismatch (missing/extra files) as blocking errors in QA.",
    )
    parser.add_argument(
        "--safe-area-profile",
        choices=("legacy", "presentation"),
        default="legacy",
        help="Safe-area warning profile for layout lint/QA. presentation is stricter near slide edges.",
    )
    parser.add_argument(
        "--safe-edge-whitelist",
        default="header,footer",
        help="Comma-separated regions allowed near safe edge warnings (header,footer).",
    )
    parser.add_argument(
        "--enable-visual-qa",
        action="store_true",
        help="Enable visual QA scoring/recommendations in qa/report.* output.",
    )
    parser.add_argument(
        "--deterministic-repair",
        action="store_true",
        help="Apply deterministic single-slide replacement from visual QA recommendations.",
    )
    parser.add_argument(
        "--auto-repair-failed-slides",
        action="store_true",
        help="Deprecated alias for --deterministic-repair.",
    )
    parser.add_argument(
        "--max-repair-rounds",
        type=int,
        default=2,
        help="Maximum deterministic repair rounds when repair is enabled.",
    )
    parser.add_argument(
        "--enforce-copyfit-contract",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Enforce dense-slide copyfit contract from slide_plan.json before export.",
    )
    parser.add_argument(
        "--native-text-verify",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Run native PPTX text overflow verification for narrow CJK columns.",
    )
    parser.add_argument(
        "--native-auto-refit-rounds",
        type=int,
        default=0,
        help="Auto refit/re-export rounds when native text verification fails.",
    )
    parser.add_argument(
        "--preflight",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Enable fast preflight checks (UTF-8/XML/forbidden features/mojibake signals).",
    )
    parser.add_argument(
        "--doctor",
        action=argparse.BooleanOptionalAction,
        default=False,
        help="Run export doctor checks before finalize/export.",
    )
    parser.add_argument("--skip-qa", action="store_true", help="Build outputs without writing QA reports.")
    parser.add_argument(
        "--profile",
        choices=("presentation", "print_a4", "proposal_consulting"),
        default="presentation",
        help="Governance profile for budget/readability gates.",
    )
    parser.add_argument(
        "--quality-mode",
        choices=("dev-fast", "release-safe", "premium"),
        default="dev-fast",
        help="QA severity mode: release-safe/premium block warnings; premium escalates selected advisories.",
    )
    parser.add_argument(
        "--svg-only-export",
        action="store_true",
        help="Allow build/export without blueprint.json for explicit compatibility paths using pre-authored SVG only.",
    )
    parser.add_argument(
        "--incremental",
        action="store_true",
        help="Enable stage-level incremental cache. Unchanged source hashes skip repeated stage work.",
    )
    parser.add_argument(
        "--changed-slide",
        action="append",
        type=int,
        default=[],
        help="Changed slide id (repeatable). Used for scoped preflight/lint/QA in iterative workflows.",
    )
    args = parser.parse_args(argv)

    phase_options = resolve_phase_options(args)
    if phase_options.used_legacy_repair_flag:
        print("warning: --auto-repair-failed-slides is deprecated; use --deterministic-repair", file=sys.stderr)

    try:
        report = build_project(
            args.project_dir,
            phase=args.phase,
            mode=args.mode,
            artifact_name=args.artifact_name,
            clean=args.clean,
            snapshots=args.snapshots,
            strict=args.strict,
            skip_render=args.skip_render,
            force_render=args.force_render,
            auto_slide_plan=args.auto_slide_plan,
            auto_slide_plan_overwrite=args.auto_slide_plan_overwrite,
            normalize_images_white_bg=args.normalize_images_white_bg,
            guard_profile=args.guard_profile,
            auto_copyfit=phase_options.auto_copyfit,
            delivery_ready=phase_options.delivery_ready,
            enable_layout_lint=phase_options.enable_layout_lint,
            layout_lint_strict=args.layout_lint_strict,
            layout_quality_profile=args.layout_quality_profile,
            enforce_blueprint_sync=args.enforce_blueprint_sync,
            safe_area_profile=phase_options.safe_area_profile,
            safe_edge_whitelist=args.safe_edge_whitelist,
            enable_visual_qa=args.enable_visual_qa,
            deterministic_repair=phase_options.deterministic_repair,
            auto_repair_failed_slides=args.auto_repair_failed_slides,
            max_repair_rounds=args.max_repair_rounds,
            enforce_copyfit_contract=args.enforce_copyfit_contract,
            native_text_verify=args.native_text_verify,
            native_auto_refit_rounds=args.native_auto_refit_rounds,
            enable_preflight=args.preflight,
            enable_doctor=args.doctor,
            skip_qa=args.skip_qa,
            profile=args.profile,
            quality_mode=args.quality_mode,
            svg_only_export=args.svg_only_export,
            incremental=args.incremental,
            changed_slides=args.changed_slide,
        )
    except Exception as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    if report is None:
        print("Build completed without QA.")
        return 0

    print(f"Build QA {'passed' if report.ok else 'failed'}: errors={report.errors}, warnings={report.warnings}")
    print(Path(report.project) / "qa" / "report.md")
    return 0 if report.ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
