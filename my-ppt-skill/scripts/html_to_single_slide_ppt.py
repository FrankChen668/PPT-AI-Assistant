#!/usr/bin/env python3
from __future__ import annotations

import argparse
import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def pick_browser() -> Path:
    candidates = []
    for name in ("msedge", "chrome", "chromium", "google-chrome"):
        found = shutil.which(name)
        if found:
            candidates.append(Path(found))
    candidates.extend(
        [
            Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
            Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
        ]
    )
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("No Edge/Chrome browser found for headless screenshot.")


def html_to_png(html_path: Path, png_path: Path) -> None:
    png_path.parent.mkdir(parents=True, exist_ok=True)
    browser = pick_browser()
    cmd = [
        str(browser),
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--no-sandbox",
        "--force-device-scale-factor=2",
        "--window-size=1600,3600",
        f"--screenshot={png_path.resolve()}",
        html_path.resolve().as_uri(),
    ]
    result = subprocess.run(cmd, check=False, capture_output=True, text=True, encoding="utf-8", errors="replace")
    if result.returncode != 0 or not png_path.exists():
        raise RuntimeError(f"Failed to screenshot HTML.\nstdout: {result.stdout}\nstderr: {result.stderr}")


def png_to_single_slide_ppt(png_path: Path, pptx_path: Path) -> None:
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(pptx_path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert an HTML page to single-slide PPTX by screenshot.")
    parser.add_argument("html", type=Path, help="Path to html file.")
    parser.add_argument("output", type=Path, help="Output pptx path.")
    parser.add_argument("--png", type=Path, help="Optional screenshot output path.")
    args = parser.parse_args()

    html_path = args.html.resolve()
    if not html_path.exists():
        raise FileNotFoundError(f"HTML not found: {html_path}")

    png_path = args.png.resolve() if args.png else args.output.with_suffix(".png").resolve()
    html_to_png(html_path, png_path)
    png_to_single_slide_ppt(png_path, args.output.resolve())
    print(f"Wrote screenshot: {png_path}")
    print(f"Wrote PPTX: {args.output.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
